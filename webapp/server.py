import os
import re
import json
import hmac
import hashlib
import uuid
import base64
import logging
import asyncio
from urllib.parse import parse_qsl
from io import BytesIO
from pathlib import Path
from aiohttp import web
from aiogram.types import FSInputFile
try:
    from PIL import Image as PilImage
    _PIL_OK = True
except Exception:
    _PIL_OK = False

import webapp

try:
    from services.template_service import TemplateService
    _TEMPLATE_SERVICE = TemplateService()
except Exception:
    _TEMPLATE_SERVICE = None

logger = logging.getLogger(__name__)

EDITOR_HTML = Path(__file__).parent / "editor.html"
NS_RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
NS_A   = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

SLIDE_W = 960
SLIDE_H = 540


def _emu_px(emu, total_emu, total_px):
    return round(emu / total_emu * total_px) if emu is not None else 0


def _compress_img_b64(img_bytes, content_type, max_w=480, quality=65):
    """Resize + compress image for preview. Returns (data_uri_src)."""
    try:
        if _PIL_OK:
            img = PilImage.open(BytesIO(img_bytes))
            if img.mode in ('RGBA', 'P', 'LA'):
                bg = PilImage.new('RGB', img.size, (255, 255, 255))
                bg.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = bg
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            if img.width > max_w:
                ratio = max_w / img.width
                img = img.resize((max_w, int(img.height * ratio)), PilImage.LANCZOS)
            buf = BytesIO()
            img.save(buf, format='JPEG', quality=quality, optimize=True)
            b64 = base64.b64encode(buf.getvalue()).decode()
            return f"data:image/jpeg;base64,{b64}"
    except Exception:
        pass
    b64 = base64.b64encode(img_bytes).decode()
    return f"data:{content_type};base64,{b64}"


def _letterbox_image(img_bytes, slot_w_emu, slot_h_emu):
    """Resize image to match slot aspect ratio with white letterboxing. Preserves aspect ratio."""
    if not _PIL_OK or not slot_w_emu or not slot_h_emu:
        return img_bytes
    try:
        img = PilImage.open(BytesIO(img_bytes))
        if img.mode in ('RGBA', 'P', 'LA'):
            bg_conv = PilImage.new('RGB', img.size, (255, 255, 255))
            bg_conv.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
            img = bg_conv
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        slot_ar = slot_w_emu / slot_h_emu
        img_ar  = img.width / img.height
        # Canvas in slot proportions at decent resolution
        canvas_w = 1920
        canvas_h = round(canvas_w / slot_ar)
        if img_ar > slot_ar:
            new_w = canvas_w
            new_h = max(1, round(canvas_w / img_ar))
        else:
            new_h = canvas_h
            new_w = max(1, round(canvas_h * img_ar))
        img_r = img.resize((new_w, new_h), PilImage.LANCZOS)
        canvas = PilImage.new('RGB', (canvas_w, canvas_h), (255, 255, 255))
        canvas.paste(img_r, ((canvas_w - new_w) // 2, (canvas_h - new_h) // 2))
        out = BytesIO()
        canvas.save(out, format='JPEG', quality=92, optimize=True)
        return out.getvalue()
    except Exception as ex:
        logger.warning(f"_letterbox_image failed: {ex}")
        return img_bytes


def _run_font_size(run):
    try:
        if run.font.size:
            return round(run.font.size.pt)
    except Exception:
        pass
    try:
        rPr = run._r.find(f"{{{NS_A}}}rPr")
        if rPr is None:
            rPr = run._r
        sz = rPr.get("sz")
        if sz:
            return round(int(sz) / 100)
    except Exception:
        pass
    return None


def _run_color(run):
    try:
        if run.font.color and run.font.color.type:
            rgb = run.font.color.rgb
            return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
    except Exception:
        pass
    return None


def _cell_fill(cell):
    """Extract a cell's solid fill color as #rrggbb, if any."""
    try:
        fill = cell.fill
        if fill and getattr(fill, "type", None) is not None:
            rgb = fill.fore_color.rgb
            if rgb is not None:
                return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
    except Exception:
        pass
    return None


def _para_align(para):
    try:
        a = str(para.alignment)
        if "CENTER" in a:
            return "center"
        if "RIGHT" in a:
            return "right"
        if "JUSTIFY" in a:
            return "justify"
    except Exception:
        pass
    return "left"


def _extract_pptx_data(file_path: str) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    sw = prs.slide_width
    sh = prs.slide_height
    slides_out = []

    for sld_id in prs.slides._sldIdLst:
        rId = sld_id.get(NS_RID)
        try:
            slide = prs.slides.part.related_slide(rId)
        except Exception:
            continue

        shapes_out = []
        for shape_idx, shape in enumerate(slide.shapes):
            left_px = _emu_px(shape.left if shape.left else 0, sw, SLIDE_W)
            top_px  = _emu_px(shape.top  if shape.top  else 0, sh, SLIDE_H)
            w_px    = _emu_px(shape.width  if shape.width  else sw, sw, SLIDE_W)
            h_px    = _emu_px(shape.height if shape.height else sh, sh, SLIDE_H)

            s = {
                "idx": shape_idx,
                "name": shape.name,
                "left": left_px,
                "top":  top_px,
                "w":    w_px,
                "h":    h_px,
                "type": "other",
            }

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    ct = shape.image.content_type or "image/jpeg"
                    s["type"] = "image"
                    s["src"]  = _compress_img_b64(img_bytes, ct)
                except Exception:
                    s["type"] = "image"
                    s["src"]  = ""

            elif getattr(shape, "has_table", False):
                tbl = shape.table
                rows_out = []
                for row_idx, row in enumerate(tbl.rows):
                    cells_out = []
                    for cell in row.cells:
                        txt = cell.text.replace('\x0b', '\n').replace('\r', '\n').strip()
                        bold = False
                        color = "#111111"
                        align = "left"
                        try:
                            for para in cell.text_frame.paragraphs:
                                a = _para_align(para)
                                if a:
                                    align = a
                                for run in para.runs:
                                    if run.font.bold:
                                        bold = True
                                    c = _run_color(run)
                                    if c:
                                        color = c
                        except Exception:
                            pass
                        fill = _cell_fill(cell)
                        cell_d = {"text": txt, "bold": bold, "color": color, "align": align}
                        if fill:
                            cell_d["fill"] = fill
                        cells_out.append(cell_d)
                    rows_out.append(cells_out)
                s["type"]  = "table"
                s["rows"]  = rows_out

            elif hasattr(shape, "text_frame"):
                tf = shape.text_frame
                paras = []
                for para in tf.paragraphs:
                    size = None
                    bold = False
                    color = None
                    align = _para_align(para)
                    for run in para.runs:
                        if size is None:
                            size = _run_font_size(run)
                        if not bold and run.font.bold:
                            bold = True
                        if color is None:
                            color = _run_color(run)

                    if size is None:
                        if top_px < SLIDE_H * 0.2:
                            size = 28
                        else:
                            size = 18

                    paras.append({
                        "text":  para.text.replace('\x0b', '\n').replace('\r', '\n'),
                        "size":  size,
                        "bold":  bold,
                        "color": color or "#111111",
                        "align": align,
                    })

                s["type"]      = "text"
                s["paras"]     = paras
                s["full_text"] = tf.text.replace('\x0b', '\n').replace('\r', '\n')
                s["wrap"]      = bool(tf.word_wrap)

            shapes_out.append(s)

        slides_out.append({"shapes": shapes_out})

    return {"slides": slides_out, "sw": SLIDE_W, "sh": SLIDE_H}


def _extract_docx_data(file_path: str) -> dict:
    """Extract paragraphs, tables, and inline images from a docx in document order.

    Returns:
        dict with:
          - items: ordered list of {type: 'para'|'table'|'image', ...} for rendering
          - paras: legacy paragraph list (backward compat)
    """
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(file_path)

    # Lists in body-order; we walk body children and advance counters deterministically
    # (avoids relying on lxml element identity/hashing which can be brittle).
    paras_seq = list(doc.paragraphs)
    tables_seq = list(doc.tables)
    p_cursor = 0
    t_cursor = 0

    items = []
    paras_out = []

    def _para_meta(para, idx):
        style_name = (para.style.name or "Normal") if para.style else "Normal"
        is_h1 = "Heading 1" in style_name
        is_h2 = "Heading 2" in style_name
        is_h3 = "Heading 3" in style_name
        size = 28 if is_h1 else (22 if is_h2 else (18 if is_h3 else 13))
        bold = is_h1 or is_h2 or is_h3
        color = "#1a1a1a"
        align = _para_align(para)
        for run in para.runs:
            if run.font.size:
                try:
                    size = round(run.font.size.pt)
                except Exception:
                    pass
            if run.font.bold:
                bold = True
            c = _run_color(run)
            if c:
                color = c
            break
        return {
            "idx": idx,
            "text": para.text,
            "style": style_name,
            "size": size,
            "bold": bold,
            "color": color,
            "align": align,
        }

    img_counter = [0]  # mutable counter shared with closure

    def _extract_inline_images(para_elem, para_idx):
        """Find all <a:blip> embeds in a paragraph and yield image item dicts.

        Increment img_counter for EVERY blip with r:embed (even if extraction
        fails) so the counter stays in lock-step with `_rebuild_docx_visual`'s
        traversal. Otherwise a corrupt blip would shift downstream indexes and
        the user could edit the wrong image.
        """
        for blip in para_elem.iter(qn('a:blip')):
            rId = blip.get(qn('r:embed'))
            if not rId:
                continue
            idx = img_counter[0]
            img_counter[0] += 1
            try:
                rel = doc.part.rels[rId]
                img_part = rel.target_part
                img_bytes = img_part.blob
                ct = getattr(img_part, "content_type", None) or "image/jpeg"
                src = _compress_img_b64(img_bytes, ct, max_w=620, quality=72)
                yield {"type": "image", "img_idx": idx, "src": src, "para_idx": para_idx}
            except Exception as ex:
                logger.warning(f"inline image extract failed idx={idx}: {ex}")

    def _extract_table(tbl, t_idx):
        rows_out = []
        for row in tbl.rows:
            cells_out = []
            for cell in row.cells:
                txt = cell.text.replace('\x0b', '\n').replace('\r', '\n').strip()
                bold = False
                color = "#111111"
                try:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.font.bold:
                                bold = True
                            c = _run_color(run)
                            if c:
                                color = c
                except Exception:
                    pass
                cells_out.append({"text": txt, "bold": bold, "color": color})
            rows_out.append(cells_out)
        return {"type": "table", "t_idx": t_idx, "rows": rows_out}

    # Iterate body children in document order so paragraphs, tables, and images interleave correctly.
    for child in doc.element.body.iterchildren():
        tag = child.tag
        if tag == qn('w:p'):
            if p_cursor >= len(paras_seq):
                continue
            idx = p_cursor
            para = paras_seq[idx]
            p_cursor += 1
            meta = _para_meta(para, idx)
            items.append({"type": "para", **meta})
            paras_out.append(meta)
            for img_item in _extract_inline_images(child, idx):
                items.append(img_item)
        elif tag == qn('w:tbl'):
            if t_cursor >= len(tables_seq):
                continue
            t_idx = t_cursor
            tbl = tables_seq[t_idx]
            t_cursor += 1
            items.append(_extract_table(tbl, t_idx))

    return {"items": items, "paras": paras_out}


def _px_to_emu(px, ref_px, slide_emu):
    """Convert pixel coordinate (in 960x540 reference space) to EMU."""
    return int(px / ref_px * slide_emu)


def _apply_shape_pos(shape, change, slide_w_emu, slide_h_emu):
    """Apply optional left/top/w/h from change dict to shape."""
    from pptx.util import Emu
    if "left" in change:
        shape.left  = _px_to_emu(change["left"], 960, slide_w_emu)
    if "top" in change:
        shape.top   = _px_to_emu(change["top"],  540, slide_h_emu)
    if "w" in change:
        shape.width  = _px_to_emu(change["w"],   960, slide_w_emu)
    if "h" in change:
        shape.height = _px_to_emu(change["h"],   540, slide_h_emu)


def _add_new_shape_to_slide(slide, ns, slide_w_emu, slide_h_emu):
    """Add a brand-new text or image shape to a slide."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from io import BytesIO

    left   = _px_to_emu(ns.get("left", 160), 960, slide_w_emu)
    top    = _px_to_emu(ns.get("top",  180), 540, slide_h_emu)
    width  = _px_to_emu(ns.get("w",    640), 960, slide_w_emu)
    height = _px_to_emu(ns.get("h",    120), 540, slide_h_emu)

    if ns.get("type") == "text":
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = ns.get("text", "")
        run.font.size = Pt(18)
        try:
            run.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
        except Exception:
            pass

    elif ns.get("type") == "image":
        src = ns.get("src", "")
        if src and "base64," in src:
            b64_data = src.split("base64,")[1]
            img_bytes = base64.b64decode(b64_data)
            try:
                slide.shapes.add_picture(BytesIO(img_bytes), left, top, width, height)
            except Exception as e:
                logger.error(f"Add new image shape error: {e}")


_SAFE_TPL_FILE_RE = re.compile(r'^[A-Za-z0-9._\- ]+\.(jpg|jpeg|png|webp)$', re.IGNORECASE)


def _resolve_template_bytes(template_id: str):
    """Return raw image bytes for a template_id from TemplateService, or None.

    Hardened against path traversal: the template file name must match a strict
    whitelist (alphanumerics + dot/dash/underscore/space, image extension only)
    and the resolved path must stay inside attached_assets/.
    """
    if not _TEMPLATE_SERVICE or not template_id:
        return None
    if not isinstance(template_id, str) or len(template_id) > 64:
        return None
    tpl = _TEMPLATE_SERVICE.templates.get(template_id)
    if not tpl:
        return None
    fname = tpl.get("file")
    if not isinstance(fname, str) or not _SAFE_TPL_FILE_RE.match(fname):
        logger.warning(f"Rejected unsafe template file name: {fname!r}")
        return None
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    assets_dir = os.path.realpath(os.path.join(base_dir, "attached_assets"))
    bg_path = os.path.realpath(os.path.join(assets_dir, fname))
    if not bg_path.startswith(assets_dir + os.sep):
        logger.warning(f"Rejected template path escaping assets dir: {bg_path}")
        return None
    if not os.path.exists(bg_path):
        return None
    try:
        with open(bg_path, "rb") as f:
            return f.read()
    except Exception:
        return None


def _apply_bg_to_slide(slide, img_bytes: bytes, slide_w_emu: int, slide_h_emu: int):
    """Replace the existing background image (full-slide picture at 0,0) or
    add a new background picture if none exists. The new picture is moved to
    the back of the shape tree so it sits behind all content."""
    if not img_bytes:
        return
    # Try to find an existing full-slide image shape and swap its blob.
    target = None
    try:
        for shape in slide.shapes:
            if not getattr(shape, "shape_type", None):
                continue
            try:
                if shape.shape_type != 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue
            except Exception:
                continue
            try:
                left = shape.left or 0
                top  = shape.top  or 0
                w    = shape.width or 0
                h    = shape.height or 0
            except Exception:
                continue
            if left <= int(slide_w_emu * 0.02) and top <= int(slide_h_emu * 0.02) \
                    and w >= int(slide_w_emu * 0.96) and h >= int(slide_h_emu * 0.96):
                target = shape
                break
    except Exception:
        target = None

    # Background pictures are placed at full slide width/height, so the image is
    # stretched by PowerPoint to fill — no letterboxing needed (would only add
    # ugly bars when aspect ratios differ).
    lb = img_bytes

    if target is not None:
        try:
            blip = target._element.find(f".//{{{NS_A}}}blip")
            if blip is not None:
                img_rId = blip.get(f"{{{NS_R}}}embed")
                if img_rId:
                    img_part = slide.part.related_part(img_rId)
                    img_part._blob = lb
                    img_part.content_type = "image/jpeg"
                    return
        except Exception as e:
            logger.error(f"Bg replace error: {e}")

    # Otherwise add new picture covering the whole slide and move it to the back.
    try:
        pic = slide.shapes.add_picture(BytesIO(lb), 0, 0,
                                       width=slide_w_emu, height=slide_h_emu)
        try:
            tree = slide.shapes._spTree
            tree.remove(pic._element)
            insert_at = 2 if len(tree) >= 2 else 0
            tree.insert(insert_at, pic._element)
        except Exception:
            pass
    except Exception as e:
        logger.error(f"Bg add error: {e}")


def _rebuild_pptx_visual(file_path: str, visual_changes: dict,
                         new_shapes_data: dict = None,
                         slide_order: list = None,
                         bg_changes: dict = None) -> str:
    from pptx import Presentation
    from config import TEMP_DIR

    prs = Presentation(file_path)
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    # Capture original slide sldId elements before any additions
    original_sldIds = list(prs.slides._sldIdLst)

    # ── 1. Apply visual_changes to original slides ──────────────────────────
    for slide_idx_str, shape_changes in visual_changes.items():
        slide_idx = int(slide_idx_str)
        if slide_idx >= len(original_sldIds):
            continue
        rId = original_sldIds[slide_idx].get(NS_RID)
        try:
            slide = prs.slides.part.related_slide(rId)
        except Exception:
            continue

        shapes_list = list(slide.shapes)

        for shape_idx_str, change in shape_changes.items():
            shape_idx = int(shape_idx_str)
            if shape_idx >= len(shapes_list):
                continue
            shape = shapes_list[shape_idx]

            # Apply position/size if provided
            _apply_shape_pos(shape, change, slide_w_emu, slide_h_emu)

            ctype = change.get("type", "pos")
            if ctype == "text":
                new_text = change.get("text", "")
                if hasattr(shape, "text_frame"):
                    tf = shape.text_frame
                    if tf.paragraphs:
                        first_para = tf.paragraphs[0]
                        if first_para.runs:
                            for run in first_para.runs:
                                run.text = ""
                            first_para.runs[0].text = new_text
                        else:
                            run = first_para.add_run()
                            run.text = new_text
                        for para in list(tf.paragraphs)[1:]:
                            p_el = para._p
                            p_el.getparent().remove(p_el)

            elif ctype == "table":
                cells_changes = change.get("cells", {})
                if getattr(shape, "has_table", False) and cells_changes:
                    tbl = shape.table
                    for cell_key, new_text in cells_changes.items():
                        try:
                            ri, ci = int(cell_key.split("_")[0]), int(cell_key.split("_")[1])
                            cell = tbl.cell(ri, ci)
                            tf = cell.text_frame
                            if tf.paragraphs:
                                first_para = tf.paragraphs[0]
                                if first_para.runs:
                                    for run in first_para.runs:
                                        run.text = ""
                                    first_para.runs[0].text = new_text
                                else:
                                    first_para.add_run().text = new_text
                                for para in list(tf.paragraphs)[1:]:
                                    para._p.getparent().remove(para._p)
                        except Exception as e:
                            logger.error(f"Table cell update error {cell_key}: {e}")

            elif ctype == "image":
                src = change.get("src", "")
                if src and "base64," in src:
                    b64_data = src.split("base64,")[1]
                    img_bytes = base64.b64decode(b64_data)
                    # Letterbox to preserve aspect ratio inside the slot
                    try:
                        lb_bytes = _letterbox_image(img_bytes, shape.width, shape.height)
                    except Exception:
                        lb_bytes = img_bytes
                    try:
                        blip = shape._element.find(f".//{{{NS_A}}}blip")
                        if blip is not None:
                            img_rId = blip.get(f"{{{NS_R}}}embed")
                            if img_rId:
                                img_part = slide.part.related_part(img_rId)
                                img_part._blob = lb_bytes
                                img_part.content_type = "image/jpeg"
                    except Exception as e:
                        logger.error(f"Image replace error slide={slide_idx} shape={shape_idx}: {e}")

    # ── 2. Create new blank slides and build sldId map ──────────────────────
    new_slide_sldId_map = {}  # entry['id'] -> sldId element
    if slide_order:
        for entry in slide_order:
            if entry.get("type") == "new":
                eid = entry["id"]
                # Use blank layout (index 6, fallback to last available)
                layout_idx = min(6, len(prs.slide_layouts) - 1)
                prs.slides.add_slide(prs.slide_layouts[layout_idx])
                new_sldId = list(prs.slides._sldIdLst)[-1]
                new_slide_sldId_map[eid] = new_sldId

    # ── 3. Add new shapes to slides ─────────────────────────────────────────
    if new_shapes_data:
        for disp_idx_str, ns_list in new_shapes_data.items():
            disp_idx = int(disp_idx_str)
            entry = slide_order[disp_idx] if slide_order and disp_idx < len(slide_order) else None
            if entry is None:
                # fallback: treat disp_idx as orig idx
                if disp_idx < len(original_sldIds):
                    rId = original_sldIds[disp_idx].get(NS_RID)
                    try:
                        slide = prs.slides.part.related_slide(rId)
                    except Exception:
                        continue
                else:
                    continue
            elif entry.get("type") == "orig":
                orig_idx = entry["idx"]
                if orig_idx >= len(original_sldIds):
                    continue
                rId = original_sldIds[orig_idx].get(NS_RID)
                try:
                    slide = prs.slides.part.related_slide(rId)
                except Exception:
                    continue
            else:
                sldId_el = new_slide_sldId_map.get(entry.get("id"))
                if sldId_el is None:
                    continue
                rId = sldId_el.get(NS_RID)
                try:
                    slide = prs.slides.part.related_slide(rId)
                except Exception:
                    continue

            for ns in (ns_list or []):
                try:
                    _add_new_shape_to_slide(slide, ns, slide_w_emu, slide_h_emu)
                except Exception as e:
                    logger.error(f"Add new shape error disp={disp_idx}: {e}")

    # ── 3b. Apply background changes ────────────────────────────────────────
    if bg_changes:
        for slide_idx_str, payload in bg_changes.items():
            try:
                slide_idx = int(slide_idx_str)
            except Exception:
                continue
            if slide_idx >= len(original_sldIds):
                continue
            rId = original_sldIds[slide_idx].get(NS_RID)
            try:
                slide = prs.slides.part.related_slide(rId)
            except Exception:
                continue
            img_bytes = None
            if isinstance(payload, dict):
                tid = payload.get("template_id")
                src = payload.get("src", "")
                if tid:
                    img_bytes = _resolve_template_bytes(tid)
                elif src and "base64," in src:
                    try:
                        img_bytes = base64.b64decode(src.split("base64,")[1])
                    except Exception:
                        img_bytes = None
            elif isinstance(payload, str) and "base64," in payload:
                try:
                    img_bytes = base64.b64decode(payload.split("base64,")[1])
                except Exception:
                    img_bytes = None
            if img_bytes:
                try:
                    _apply_bg_to_slide(slide, img_bytes, slide_w_emu, slide_h_emu)
                except Exception as e:
                    logger.error(f"Apply bg slide={slide_idx}: {e}")

    # ── 4. Reorder slides per slide_order ───────────────────────────────────
    if slide_order and len(slide_order) > 1:
        sldIdLst = prs.slides._sldIdLst
        ordered = []
        for entry in slide_order:
            if entry.get("type") == "orig":
                idx = entry["idx"]
                if idx < len(original_sldIds):
                    ordered.append(original_sldIds[idx])
            else:
                sldId_el = new_slide_sldId_map.get(entry.get("id"))
                if sldId_el is not None:
                    ordered.append(sldId_el)
        if ordered:
            for el in list(sldIdLst):
                sldIdLst.remove(el)
            for el in ordered:
                sldIdLst.append(el)

    fname = f"edited_{uuid.uuid4().hex[:8]}.pptx"
    out = os.path.join(TEMP_DIR, fname)
    os.makedirs(TEMP_DIR, exist_ok=True)
    prs.save(out)
    return out


def _clone_image_part(doc_part, img_part):
    """Clone a docx ImagePart so a new blip can target an independent copy."""
    from docx.parts.image import ImagePart
    from docx.opc.constants import CONTENT_TYPE
    from docx.opc.packuri import PackURI
    # Find an unused partname like /word/media/imageN.<ext>
    base = "/word/media/image"
    ext = ".jpg"
    try:
        ct = (img_part.content_type or "").lower()
        if "png" in ct: ext = ".png"
        elif "gif" in ct: ext = ".gif"
        elif "jpeg" in ct or "jpg" in ct: ext = ".jpg"
    except Exception:
        pass
    existing = {p.partname for p in doc_part.package.iter_parts()}
    n = 1
    while True:
        candidate = PackURI(f"{base}{n}{ext}")
        if candidate not in existing:
            break
        n += 1
    new_part = ImagePart(
        candidate,
        img_part.content_type or "image/jpeg",
        img_part.blob,
        doc_part.package,
    )
    return new_part


def _rebuild_docx_visual(
    file_path: str,
    para_changes: dict,
    table_changes: dict = None,
    image_changes: dict = None,
) -> str:
    """Apply paragraph, table-cell, and inline-image changes to a docx.

    - table_changes keys: "t_idx_ri_ci" → new text.
    - image_changes keys: "img_idx" (string of int) → "data:image/...;base64,..."
    """
    from docx import Document
    from docx.oxml.ns import qn
    from config import TEMP_DIR

    doc = Document(file_path)

    # ── Paragraph text updates ─────────────────────────────────────────────
    for idx_str, new_text in (para_changes or {}).items():
        try:
            idx = int(idx_str)
        except (TypeError, ValueError):
            continue
        if not isinstance(new_text, str):
            continue
        if 0 <= idx < len(doc.paragraphs):
            para = doc.paragraphs[idx]
            if para.runs:
                para.runs[0].text = new_text
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.add_run(new_text)

    # ── Table cell updates ─────────────────────────────────────────────────
    if table_changes:
        for key, new_text in table_changes.items():
            try:
                parts = str(key).split('_')
                if len(parts) != 3:
                    continue
                t_idx, ri, ci = int(parts[0]), int(parts[1]), int(parts[2])
                if not isinstance(new_text, str):
                    continue
                # Reject negative indices — Python would silently wrap-around and modify
                # unintended cells.
                if t_idx < 0 or ri < 0 or ci < 0:
                    continue
                if t_idx >= len(doc.tables):
                    continue
                tbl = doc.tables[t_idx]
                if ri >= len(tbl.rows):
                    continue
                row = tbl.rows[ri]
                if ci >= len(row.cells):
                    continue
                cell = row.cells[ci]
                if not cell.paragraphs:
                    cell.add_paragraph(new_text)
                    continue
                # Write into first paragraph, preserve its run formatting.
                p0 = cell.paragraphs[0]
                if p0.runs:
                    p0.runs[0].text = new_text
                    for run in p0.runs[1:]:
                        run.text = ""
                else:
                    p0.add_run(new_text)
                # Clear extra paragraphs inside the cell (avoid stray old text).
                for extra in cell.paragraphs[1:]:
                    for run in extra.runs:
                        run.text = ""
            except Exception as ex:
                logger.warning(f"table cell update failed key={key}: {ex}")

    # ── Inline image replacements ─────────────────────────────────────────
    if image_changes:
        # Build {int_idx: bytes} from the data: URLs.
        img_payloads = {}
        for k, src in image_changes.items():
            try:
                idx = int(k)
            except (TypeError, ValueError):
                continue
            if idx < 0 or not isinstance(src, str) or "base64," not in src:
                continue
            try:
                b64 = src.split("base64,", 1)[1]
                img_payloads[idx] = base64.b64decode(b64)
            except Exception as ex:
                logger.warning(f"image decode failed idx={idx}: {ex}")

        if img_payloads:
            # Walk the body in the SAME order as extraction so img_idx aligns.
            # Increment counter for EVERY blip with r:embed regardless of success
            # to mirror _extract_inline_images exactly.
            counter = 0
            # Track which target parts we've already mutated so a doc with two
            # blips pointing at the same media part doesn't overwrite the first
            # edit when the second is processed.
            mutated_parts = {}
            for child in doc.element.body.iterchildren():
                if child.tag != qn('w:p'):
                    continue
                for blip in child.iter(qn('a:blip')):
                    rId = blip.get(qn('r:embed'))
                    if not rId:
                        continue
                    cur = counter
                    counter += 1
                    if cur not in img_payloads:
                        continue
                    try:
                        rel = doc.part.rels[rId]
                        img_part = rel.target_part
                        # If two image items share one underlying part, clone the
                        # part for this blip so edits stay independent.
                        part_id = id(img_part)
                        if part_id in mutated_parts:
                            try:
                                new_part = _clone_image_part(doc.part, img_part)
                                # Re-point this blip's relationship to the new part.
                                doc.part.rels[rId].target_part = new_part
                                img_part = new_part
                            except Exception as cex:
                                logger.warning(f"image part clone failed idx={cur}: {cex}")
                        mutated_parts[id(img_part)] = True
                        new_bytes = img_payloads[cur]
                        # Normalize to JPEG for consistent embedding.
                        try:
                            if _PIL_OK:
                                pim = PilImage.open(BytesIO(new_bytes))
                                out_buf = BytesIO()
                                if pim.mode in ('RGBA', 'P', 'LA'):
                                    bg_im = PilImage.new('RGB', pim.size, (255, 255, 255))
                                    bg_im.paste(pim, mask=pim.split()[-1] if pim.mode in ('RGBA', 'LA') else None)
                                    pim = bg_im
                                elif pim.mode != 'RGB':
                                    pim = pim.convert('RGB')
                                pim.save(out_buf, format='JPEG', quality=88, optimize=True)
                                new_bytes = out_buf.getvalue()
                        except Exception as nex:
                            logger.warning(f"image normalize failed idx={cur}: {nex}")
                        img_part._blob = new_bytes
                        img_part.content_type = "image/jpeg"
                    except Exception as ex:
                        logger.warning(f"image replace failed idx={cur}: {ex}")

    fname = f"edited_{uuid.uuid4().hex[:8]}.docx"
    out = os.path.join(TEMP_DIR, fname)
    os.makedirs(TEMP_DIR, exist_ok=True)
    doc.save(out)
    return out


def _create_edited_docx(content: str, topic: str) -> str:
    from docx import Document
    from config import TEMP_DIR
    doc = Document()
    lines = content.split("\n")
    first = True
    for line in lines:
        if first and line.strip():
            doc.add_heading(line.strip(), level=1)
            first = False
        else:
            doc.add_paragraph(line)
    fname = f"edited_{uuid.uuid4().hex[:8]}.docx"
    out = os.path.join(TEMP_DIR, fname)
    os.makedirs(TEMP_DIR, exist_ok=True)
    doc.save(out)
    return out


def _create_edited_pptx(content: str, original_path: str) -> str:
    from pptx import Presentation
    from config import TEMP_DIR
    prs = Presentation(original_path)
    sld_ids = list(prs.slides._sldIdLst)
    sections = re.split(r"=== Slayd \d+ ===\s*\n?", content)
    sections = [s.strip() for s in sections if s.strip()]
    for i, sld_id in enumerate(sld_ids):
        if i >= len(sections):
            break
        rId = sld_id.get(NS_RID)
        try:
            slide = prs.slides.part.related_slide(rId)
        except Exception:
            continue
        new_text = sections[i]
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            if not tf.text.strip():
                continue
            if tf.paragraphs:
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = new_text
                else:
                    tf.paragraphs[0].add_run().text = new_text
            break
    fname = f"edited_{uuid.uuid4().hex[:8]}.pptx"
    out = os.path.join(TEMP_DIR, fname)
    os.makedirs(TEMP_DIR, exist_ok=True)
    prs.save(out)
    return out


_INIT_DATA_MAX_AGE_SEC = 24 * 60 * 60  # reject initData older than 24h (replay defence)


def _verify_telegram_init_data(init_data: str) -> dict | None:
    """Verify Telegram Mini App initData HMAC and return parsed user dict on success.

    Spec: https://core.telegram.org/bots/webapps#validating-data-received-via-the-mini-app
    Also enforces an `auth_date` freshness window to limit replay.
    """
    if not init_data or not isinstance(init_data, str):
        return None
    try:
        bot_token = os.environ.get("BOT_TOKEN", "")
        if not bot_token:
            logger.error("BOT_TOKEN missing — cannot verify Telegram initData")
            return None
        parsed = dict(parse_qsl(init_data, keep_blank_values=True))
        recv_hash = parsed.pop("hash", None)
        if not recv_hash:
            return None
        data_check = "\n".join(f"{k}={parsed[k]}" for k in sorted(parsed.keys()))
        secret = hmac.new(b"WebAppData", bot_token.encode(), hashlib.sha256).digest()
        calc = hmac.new(secret, data_check.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(calc, recv_hash):
            return None
        # Replay defence: auth_date must be present and recent.
        try:
            import time as _t
            auth_date = int(parsed.get("auth_date", "0"))
            if auth_date <= 0 or (_t.time() - auth_date) > _INIT_DATA_MAX_AGE_SEC:
                logger.warning(f"initData too old or missing auth_date: {auth_date}")
                return None
        except (ValueError, TypeError):
            return None
        user_raw = parsed.get("user")
        if not user_raw:
            return None
        try:
            parsed["_user"] = json.loads(user_raw)
        except Exception:
            return None
        return parsed
    except Exception as ex:
        logger.warning(f"initData verify error: {ex}")
        return None


def _authorize_token_request(request: web.Request, token: str) -> tuple[dict | None, web.Response | None]:
    """Resolve token info and verify the requester is its owner via Telegram initData.

    initData is REQUIRED — there is no legacy fallback. Possession of the token
    alone is not enough; the request must come from the Telegram user the token
    was issued to.
    """
    info = webapp.DOC_TOKENS.get(token)
    if not info:
        return None, web.json_response({"error": "expired"}, status=404)
    init_data = request.headers.get("X-Telegram-Init-Data", "")
    if not init_data:
        logger.warning(f"Token request missing initData token={token[:8]}...")
        return None, web.json_response({"error": "unauthorized"}, status=401)
    verified = _verify_telegram_init_data(init_data)
    if not verified:
        return None, web.json_response({"error": "unauthorized"}, status=401)
    u = verified.get("_user") or {}
    uid = u.get("id")
    owner_id = info.get("user_id")
    # Fail closed: both ids must be present and equal.
    if uid is None or owner_id is None or int(uid) != int(owner_id):
        logger.warning(
            f"Token owner mismatch: token_owner={owner_id} request_user={uid}"
        )
        return None, web.json_response({"error": "forbidden"}, status=403)
    return info, None


# ── JSON schema validation for /api/save ──────────────────────────────
def _validate_save_payload(body: dict) -> str | None:
    """Return None if payload is shape-valid, else an error message string."""
    if not isinstance(body, dict):
        return "Body must be an object"
    ft = body.get("file_type")
    if ft not in ("pptx", "docx"):
        return "file_type must be 'pptx' or 'docx'"
    if "visual_changes" in body:
        vc = body["visual_changes"]
        if not isinstance(vc, dict):
            return "visual_changes must be an object"
        for sk, shape_changes in vc.items():
            if not str(sk).lstrip("-").isdigit():
                return f"visual_changes key '{sk}' must be int-like"
            if not isinstance(shape_changes, dict):
                return f"visual_changes[{sk}] must be an object"
            for shk, ch in shape_changes.items():
                if not isinstance(ch, dict):
                    return f"visual_changes[{sk}][{shk}] must be an object"
                ctype = ch.get("type", "pos")
                if ctype not in ("text", "table", "image", "pos"):
                    return f"unknown change type '{ctype}'"
                if ctype == "text" and not isinstance(ch.get("text", ""), str):
                    return "text change must have string 'text'"
                if ctype == "table":
                    cells = ch.get("cells", {})
                    if not isinstance(cells, dict):
                        return "table change 'cells' must be an object"
                    for ck, cv in cells.items():
                        if not isinstance(cv, str):
                            return f"table cell '{ck}' must be string"
                if ctype == "image":
                    src = ch.get("src", "")
                    if not isinstance(src, str) or len(src) > 8 * 1024 * 1024:
                        return "image src must be string under 8MB"
        ns = body.get("new_shapes", {})
        if not isinstance(ns, dict):
            return "new_shapes must be an object"
        so = body.get("slide_order")
        if so is not None and not isinstance(so, list):
            return "slide_order must be an array"
        bg = body.get("bg_changes")
        if bg is not None and not isinstance(bg, dict):
            return "bg_changes must be an object"
    if ft == "docx":
        for fld in ("para_changes", "table_changes", "image_changes"):
            v = body.get(fld)
            if v is not None and not isinstance(v, dict):
                return f"{fld} must be an object"
    if "content" in body and not isinstance(body["content"], str):
        return "content must be a string"
    return None


async def handle_editor(request: web.Request) -> web.Response:
    try:
        content = EDITOR_HTML.read_text(encoding="utf-8")
        return web.Response(text=content, content_type="text/html")
    except Exception as e:
        logger.error(f"Editor read error: {e}")
        return web.Response(text="Not found", status=500)


async def handle_doc_api(request: web.Request) -> web.Response:
    token = request.match_info.get("token", "")
    info, err = _authorize_token_request(request, token)
    if err is not None:
        return err
    file_path = info["file_path"]
    if not os.path.exists(file_path):
        return web.json_response({"error": "file missing"}, status=404)
    return web.json_response({
        "topic": info.get("topic", ""),
        "file_type": "pptx" if file_path.endswith(".pptx") else "docx",
        "user_lang": info.get("user_lang", "uz"),
    })


async def handle_slides_api(request: web.Request) -> web.Response:
    token = request.match_info.get("token", "")
    info, err = _authorize_token_request(request, token)
    if err is not None:
        return err
    file_path = info["file_path"]
    if not os.path.exists(file_path):
        return web.json_response({"error": "file missing"}, status=404)
    try:
        loop = asyncio.get_event_loop()
        if file_path.endswith(".pptx"):
            data = await loop.run_in_executor(None, _extract_pptx_data, file_path)
            data["file_type"] = "pptx"
        else:
            data = await loop.run_in_executor(None, _extract_docx_data, file_path)
            data["file_type"] = "docx"
        data["topic"] = info.get("topic", "")
        data["user_lang"] = info.get("user_lang", "uz")
        return web.json_response(data)
    except Exception as e:
        logger.error(f"Slides API error token={token}: {e}")
        # Don't leak internal details to the client.
        return web.json_response({"error": "Failed to load document"}, status=500)


async def handle_save_api(request: web.Request) -> web.Response:
    token = request.match_info.get("token", "")
    info, err = _authorize_token_request(request, token)
    if err is not None:
        return err

    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    schema_err = _validate_save_payload(body)
    if schema_err:
        logger.warning(f"Rejected /api/save payload token={token[:8]}: {schema_err}")
        return web.json_response({"ok": False, "error": "Invalid payload"}, status=400)

    topic = info.get("topic", "")
    chat_id = info.get("chat_id")
    file_path = info.get("file_path", "")
    file_type = body.get("file_type", "docx")

    try:
        loop = asyncio.get_event_loop()

        if "visual_changes" in body:
            visual_changes = body["visual_changes"]
            if file_type == "pptx":
                new_shapes_data = body.get("new_shapes", {})
                slide_order = body.get("slide_order", None)
                bg_changes = body.get("bg_changes", None)
                new_path = await loop.run_in_executor(
                    None, _rebuild_pptx_visual, file_path,
                    visual_changes, new_shapes_data, slide_order, bg_changes
                )
            else:
                para_changes = body.get("para_changes", {})
                table_changes = body.get("table_changes", {})
                image_changes = body.get("image_changes", {})
                new_path = await loop.run_in_executor(
                    None, _rebuild_docx_visual, file_path,
                    para_changes, table_changes, image_changes,
                )
        elif "content" in body:
            content = body["content"]
            if file_type == "pptx":
                new_path = await loop.run_in_executor(None, _create_edited_pptx, content, file_path)
            else:
                new_path = await loop.run_in_executor(None, _create_edited_docx, content, topic)
        else:
            return web.json_response({"ok": False, "error": "No content or changes"}, status=400)

    except Exception as e:
        import traceback
        logger.error(f"File build error token={token}: {e}\n{traceback.format_exc()}")
        # Don't echo the raw exception to the client — it can leak internals.
        return web.json_response({"ok": False, "error": "File build failed"}, status=500)

    try:
        bot = webapp.BOT
        if bot and chat_id:
            import html as _html
            safe_topic = _html.escape(topic)[:900]
            caption_text = f"✏️ {safe_topic}"
            await bot.send_document(
                chat_id=chat_id,
                document=FSInputFile(new_path),
                caption=caption_text,
                parse_mode=None,
            )
        # Delete the old generated file before replacing with new edited file
        old_path = webapp.DOC_TOKENS[token].get("file_path", "")
        if old_path and old_path != new_path and os.path.exists(old_path):
            try:
                os.remove(old_path)
            except Exception:
                pass
        webapp.DOC_TOKENS[token]["file_path"] = new_path
        webapp.save_tokens_to_disk()
        logger.info(f"Edited file sent chat_id={chat_id} token={token[:8]}...")
    except Exception as e:
        logger.error(f"Send error token={token}: {e}")
        return web.json_response({"ok": False, "error": "Send failed"}, status=500)

    return web.json_response({"ok": True})


async def handle_root(request: web.Request) -> web.Response:
    raise web.HTTPFound("/edit")


def create_web_app() -> web.Application:
    app = web.Application(client_max_size=60 * 1024 * 1024)
    app.router.add_get("/", handle_root)
    app.router.add_get("/edit", handle_editor)
    app.router.add_get("/api/doc/{token}", handle_doc_api)
    app.router.add_get("/api/slides/{token}", handle_slides_api)
    app.router.add_post("/api/save/{token}", handle_save_api)
    app.router.add_get("/api/templates", handle_templates_api)
    app.router.add_get("/api/template-image/{tid}", handle_template_image_api)
    return app


async def handle_templates_api(request: web.Request) -> web.Response:
    if not _TEMPLATE_SERVICE:
        return web.json_response({"templates": []})
    out = []
    for tid, t in _TEMPLATE_SERVICE.templates.items():
        out.append({
            "id": tid,
            "name": t.get("name", {}),
            "url": f"/api/template-image/{tid}",
        })
    return web.json_response({"templates": out})


async def handle_template_image_api(request: web.Request) -> web.Response:
    tid = request.match_info.get("tid", "")
    if not _TEMPLATE_SERVICE:
        return web.Response(status=404)
    tpl = _TEMPLATE_SERVICE.templates.get(tid)
    if not tpl:
        return web.Response(status=404)
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    bg_path = os.path.join(base_dir, "attached_assets", tpl["file"])
    if not os.path.exists(bg_path):
        return web.Response(status=404)
    ct = "image/jpeg"
    low = bg_path.lower()
    if low.endswith(".png"):  ct = "image/png"
    elif low.endswith(".webp"): ct = "image/webp"
    return web.FileResponse(bg_path, headers={
        "Content-Type": ct,
        "Cache-Control": "public, max-age=86400",
    })


async def start_web_server(port: int = 5000):
    app = create_web_app()
    runner = web.AppRunner(app)
    await runner.setup()
    site = web.TCPSite(runner, "0.0.0.0", port)
    await site.start()
    logger.info(f"Web server started on port {port}")
    try:
        await asyncio.Event().wait()  # run forever until cancelled
    except asyncio.CancelledError:
        pass
    finally:
        await runner.cleanup()
        logger.info("Web server stopped cleanly")
