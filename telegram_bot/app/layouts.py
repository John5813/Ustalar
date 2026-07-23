"""
Layout banki: har funksiya bitta slaydni python-pptx orqali "qo'lda" chizadi
(standart placeholder layout ISHLATILMAYDI — barcha shakl/matn aniq koordinata bilan qo'shiladi).

Yangi layoutlar:
  bar_chart     — vertikal bar diagrammasi, Y o'qi bilan
  xy_chart      — chiziqli (line) diagramma, X va Y o'qlari bilan
  math_formula  — matematik formulalar va izohlar slaydasi
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_background(slide, color_hex: str):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_rect(slide, x, y, w, h, fill_hex=None, line=False, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill_hex:
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    else:
        shp.fill.background()
    if not line:
        shp.line.fill.background()
    return shp


def add_circle(slide, x, y, d, fill_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shp.line.fill.background()
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, italic=False,
             color_hex="000000", font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = hex_to_rgb(color_hex)
    return tb


def motif_badge(slide, x, y, d, label, palette, fonts, motif):
    """Butun taqdimotda takrorlanadigan motiv: raqam-badge, ikonka-doira yoki ramka."""
    if motif == "rounded_frame":
        shp = add_rect(slide, x, y, d, d, fill_hex=palette["accent"], radius=True)
    else:
        shp = add_circle(slide, x, y, d, palette["primary"])
    add_text(slide, x, y, d, d, str(label), size=int(d * 22), bold=True,
             color_hex="FFFFFF", font=fonts["body"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------- title_dark
def render_title_dark(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, palette["primary"])
    add_circle(slide, 10.8, -1.2, 3.6, palette["secondary"])
    add_circle(slide, -1.4, 5.6, 3.2, palette["secondary"])
    add_text(slide, 0.9, 2.5, 10.5, 1.3, s.title or "", size=40, bold=True,
              color_hex="FFFFFF", font=fonts["heading"])
    add_rect(slide, 0.9, 3.75, 0.9, 0.06, fill_hex=palette["accent"])
    if s.hook_line:
        add_text(slide, 0.9, 3.95, 10.5, 0.8, s.hook_line, size=17, italic=True,
                  color_hex=palette["secondary"], font=fonts["body"])
    if s.subtitle:
        add_text(slide, 0.9, 4.75, 10.5, 0.5, s.subtitle, size=13,
                  color_hex="CCCCCC", font=fonts["body"])


# ----------------------------------------------------------- conclusion_dark
def render_conclusion_dark(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, palette["primary"])
    add_text(slide, 0.9, 0.7, 8, 0.8, s.title or "Xulosa", size=32, bold=True,
              color_hex="FFFFFF", font=fonts["heading"])
    add_rect(slide, 0.9, 1.45, 0.9, 0.06, fill_hex=palette["accent"])
    takeaways = s.key_takeaways or []
    y = 1.75
    for i, t in enumerate(takeaways[:4], start=1):
        add_circle(slide, 0.9, y, 0.35, palette["accent"])
        add_text(slide, 0.9, y, 0.35, 0.35, str(i), size=11, bold=True,
                  color_hex="FFFFFF", font=fonts["body"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, 1.45, y, 10.0, 0.4, t, size=14, color_hex="EEEEEE", font=fonts["body"])
        y += 0.75
    if s.closing_thought:
        add_rect(slide, 0.9, 6.3, 11.5, 0.75, fill_hex=palette["accent"], radius=True)
        add_text(slide, 1.1, 6.3, 11.1, 0.75, s.closing_thought, size=14, italic=True,
                  color_hex="FFFFFF", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)


# ----------------------------------------------------------- agenda_numbered
def render_agenda_numbered(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_rect(slide, 0.0, 0.0, 13.333, 1.3, fill_hex=palette["primary"])
    add_text(slide, 0.8, 0.3, 11.5, 0.8, s.title or "", size=28, bold=True,
              color_hex="FFFFFF", font=fonts["heading"])
    items = s.items or []
    y = 1.65
    for i, it in enumerate(items[:6], start=1):
        add_rect(slide, 0.8, y, 11.7, 0.7, fill_hex="F4F6FA", radius=True)
        motif_badge(slide, 0.9, y + 0.05, 0.6, i, palette, fonts, motif)
        add_text(slide, 1.7, y + 0.05, 4.0, 0.35, it.heading, size=13, bold=True,
                  color_hex=palette["primary"], font=fonts["heading"])
        add_text(slide, 1.7, y + 0.38, 9.5, 0.3, it.body, size=11,
                  color_hex="555555", font=fonts["body"])
        y += 0.85


# -------------------------------------------------------------- two_card_compare
def render_two_card_compare(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_text(slide, 0.8, 0.4, 11.5, 0.7, s.title or "", size=28, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])
    items = s.items or []
    colors = [palette["primary"], palette["accent"]]
    x = 0.8
    for i, it in enumerate(items[:2]):
        add_rect(slide, x, 1.4, 5.7, 5.5, fill_hex=colors[i % 2], radius=True)
        add_text(slide, x + 0.3, 1.7, 5.1, 0.6, it.heading, size=18, bold=True,
                  color_hex="FFFFFF", font=fonts["heading"])
        add_rect(slide, x + 0.3, 2.35, 5.1, 0.05, fill_hex="FFFFFF")
        add_text(slide, x + 0.3, 2.55, 5.1, 3.5, it.body, size=13,
                  color_hex="FFFFFF", font=fonts["body"])
        if it.example:
            add_text(slide, x + 0.3, 6.0, 5.1, 0.5, f"Misol: {it.example}", size=11,
                      italic=True, color_hex="EEEEEE", font=fonts["body"])
        x += 6.1


# ---------------------------------------------------------- three/four card
def _card_grid(slide, s, palette, fonts, motif, n):
    set_background(slide, "FFFFFF")
    add_text(slide, 0.8, 0.5, 10, 0.6, s.title or "", size=28, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])
    items = (s.items or [])[:n]
    gap = 0.3
    total_w = 11.7
    cw = (total_w - gap * (n - 1)) / n
    x = 0.8
    for i, it in enumerate(items, start=1):
        add_rect(slide, x, 1.6, cw, 4.4, fill_hex="F4F6FA", radius=True)
        motif_badge(slide, x + 0.3, 1.95, 0.6, i, palette, fonts, motif)
        add_text(slide, x + 0.3, 2.75, cw - 0.6, 0.9, it.heading, size=15, bold=True,
                  color_hex=palette["primary"], font=fonts["heading"])
        add_text(slide, x + 0.3, 3.6, cw - 0.6, 2.1, it.body, size=11.5,
                  color_hex="333333", font=fonts["body"])
        x += cw + gap


def render_three_card_grid(slide, s, palette, fonts, motif, image_path=None):
    _card_grid(slide, s, palette, fonts, motif, 3)


def render_four_card_grid(slide, s, palette, fonts, motif, image_path=None):
    _card_grid(slide, s, palette, fonts, motif, 4)


# -------------------------------------------------------------- comparison_table
def render_comparison_table(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_text(slide, 0.8, 0.5, 10.5, 0.6, s.title or "", size=26, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])
    rows = s.items or []
    row_h = 0.65
    y = 1.5
    add_rect(slide, 0.8, y, 11.7, row_h, fill_hex=palette["primary"])
    add_text(slide, 1.0, y, 5.6, row_h, "Belgisi", size=13, bold=True,
              color_hex="FFFFFF", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 6.6, y, 5.7, row_h, "Tafsilot", size=13, bold=True,
              color_hex="FFFFFF", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)
    y += row_h
    for i, it in enumerate(rows):
        if i % 2 == 0:
            add_rect(slide, 0.8, y, 11.7, row_h, fill_hex="F4F6FA")
        add_text(slide, 1.0, y, 5.6, row_h, it.heading, size=12.5, bold=True,
                  color_hex="222222", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, 6.6, y, 5.7, row_h, it.body, size=12.5,
                  color_hex="444444", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)
        y += row_h


# -------------------------------------------------------------- icon_row_list
def render_icon_row_list(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_rect(slide, 0.0, 0.0, 13.333, 1.3, fill_hex=palette["primary"])
    add_text(slide, 0.8, 0.3, 11.5, 0.8, s.title or "", size=28, bold=True,
              color_hex="FFFFFF", font=fonts["heading"])
    if s.body:
        add_text(slide, 0.8, 1.5, 11.5, 0.5, s.body, size=13, italic=True,
                  color_hex="555555", font=fonts["body"])
    items = s.items or []
    y = 2.2
    for i, it in enumerate(items[:5]):
        motif_badge(slide, 0.8, y, 0.55, i + 1, palette, fonts, motif)
        add_text(slide, 1.55, y, 10.8, 0.3, it.heading, size=14, bold=True,
                  color_hex=palette["primary"], font=fonts["heading"])
        add_text(slide, 1.55, y + 0.32, 10.8, 0.5, it.body, size=12,
                  color_hex="444444", font=fonts["body"])
        y += 0.95


# -------------------------------------------------------------- stat_callout
def render_stat_callout(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, palette["primary"])
    add_text(slide, 1.0, 0.6, 11.3, 0.7, s.title or "", size=22, bold=True,
              color_hex="CCCCCC", font=fonts["heading"])
    items = s.items or []
    if items:
        it = items[0]
        add_text(slide, 1.0, 1.6, 11.3, 2.5, it.heading, size=88, bold=True,
                  color_hex=palette["accent"], font=fonts["heading"], align=PP_ALIGN.CENTER)
        add_rect(slide, 2.5, 4.2, 8.3, 0.06, fill_hex=palette["secondary"])
        add_text(slide, 1.0, 4.4, 11.3, 1.5, it.body, size=18,
                  color_hex="EEEEEE", font=fonts["body"], align=PP_ALIGN.CENTER)


# -------------------------------------------------------------- stat_row_triple
def render_stat_row_triple(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_rect(slide, 0.0, 0.0, 13.333, 1.3, fill_hex=palette["primary"])
    add_text(slide, 0.8, 0.3, 11.5, 0.8, s.title or "", size=28, bold=True,
              color_hex="FFFFFF", font=fonts["heading"])
    items = (s.items or [])[:3]
    x = 0.8
    for it in items:
        add_rect(slide, x, 2.0, 3.7, 3.5, fill_hex="F4F6FA", radius=True)
        add_text(slide, x + 0.1, 2.3, 3.5, 1.4, it.heading, size=36, bold=True,
                  color_hex=palette["accent"], font=fonts["heading"], align=PP_ALIGN.CENTER)
        add_rect(slide, x + 0.5, 3.7, 2.7, 0.05, fill_hex=palette["primary"])
        add_text(slide, x + 0.1, 3.85, 3.5, 1.4, it.body, size=12,
                  color_hex="555555", font=fonts["body"], align=PP_ALIGN.CENTER)
        x += 4.0


# -------------------------------------------------------------- timeline_process
def render_timeline_process(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_text(slide, 0.8, 0.4, 11.5, 0.7, s.title or "", size=28, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])
    items = s.items or []
    n = min(len(items), 5)
    if n == 0:
        return
    # Gorizontal chiziq
    line_y = 3.5
    line_x_start = 1.2
    line_x_end = 12.1
    add_rect(slide, line_x_start, line_y - 0.02, line_x_end - line_x_start, 0.06,
              fill_hex=palette["secondary"])
    step = (line_x_end - line_x_start) / n
    for i, it in enumerate(items[:n]):
        cx = line_x_start + step * i + step / 2
        # Doira
        add_circle(slide, cx - 0.32, line_y - 0.32, 0.64, palette["primary"])
        add_text(slide, cx - 0.32, line_y - 0.32, 0.64, 0.64, str(i + 1), size=14, bold=True,
                  color_hex="FFFFFF", font=fonts["body"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Sarlavha (ustida)
        add_text(slide, cx - 1.0, line_y - 1.6, 2.0, 1.0, it.heading, size=12, bold=True,
                  color_hex=palette["primary"], font=fonts["heading"], align=PP_ALIGN.CENTER)
        # Tavsif (ostida)
        add_text(slide, cx - 1.1, line_y + 0.5, 2.2, 1.5, it.body, size=10,
                  color_hex="555555", font=fonts["body"], align=PP_ALIGN.CENTER)


# -------------------------------------------------------------- image_half_bleed
def render_image_half_bleed(slide, s, palette, fonts, motif, image_path=None):
    from pptx.util import Inches
    set_background(slide, "FFFFFF")
    # Chap tomon: rasm yoki rang blok
    if image_path:
        try:
            slide.shapes.add_picture(image_path, Inches(0), Inches(0),
                                     Inches(6.4), Inches(7.5))
        except Exception:
            add_rect(slide, 0.0, 0.0, 6.4, 7.5, fill_hex=palette["primary"])
    else:
        add_rect(slide, 0.0, 0.0, 6.4, 7.5, fill_hex=palette["primary"])
    # O'ng tomon: matn
    add_text(slide, 7.0, 1.5, 5.7, 1.0, s.title or "", size=26, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])
    add_rect(slide, 7.0, 2.55, 0.7, 0.06, fill_hex=palette["accent"])
    add_text(slide, 7.0, 2.8, 5.7, 3.8, s.body or "", size=14,
              color_hex="333333", font=fonts["body"])
    if s.title:
        add_text(slide, 1.2, 5.1, 10.9, 0.5, s.title, size=13, color_hex="555555",
                  font=fonts["body"], align=PP_ALIGN.CENTER)


# -------------------------------------------------------------- quote_callout
def render_quote_callout(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, palette["secondary"])
    # Katta tirnoq belgisi
    add_text(slide, 0.7, 0.5, 2.0, 2.0, "\u201c", size=120, bold=True,
              color_hex=palette["accent"], font=fonts["heading"])
    add_text(slide, 1.5, 1.8, 10.0, 3.0, s.body or "", size=22, italic=True,
              color_hex=palette["primary"], font=fonts["heading"], align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 3.0, 5.3, 7.3, 0.06, fill_hex=palette["accent"])
    add_text(slide, 1.5, 5.5, 10.3, 0.6, f"— {s.title or ''}", size=14, bold=True,
              color_hex=palette["primary"], font=fonts["body"], align=PP_ALIGN.CENTER)


# -------------------------------------------------------------- definition_spotlight
def render_definition_spotlight(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_text(slide, 0.8, 1.6, 5.2, 2.0, s.title or "", size=38, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])
    add_rect(slide, 6.4, 1.6, 0.05, 4.0, fill_hex=palette["accent"])
    add_text(slide, 6.8, 1.8, 5.7, 3.8, s.body or "", size=15, color_hex="333333", font=fonts["body"])


# -------------------------------------------------------------- mini_case_study
def render_mini_case_study(slide, s, palette, fonts, motif, image_path=None):
    set_background(slide, "FFFFFF")
    add_text(slide, 0.8, 0.5, 10.5, 0.6, s.title or "", size=26, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])
    items = s.items or []
    labels = ["Muammo", "Yechim"]
    x = 0.8
    for i, it in enumerate(items[:2]):
        add_rect(slide, x, 1.6, 5.55, 4.4, fill_hex="F4F6FA", radius=True)
        add_text(slide, x + 0.3, 1.9, 5.0, 0.4, labels[i] if i < 2 else it.heading, size=13, bold=True,
                  color_hex=palette["accent"], font=fonts["body"])
        add_text(slide, x + 0.3, 2.4, 5.0, 3.3, it.body, size=13, color_hex="333333", font=fonts["body"])
        x += 5.85


# ============================================================
# YANGI LAYOUTLAR: diagramma va matematik formulalar
# ============================================================

# -------------------------------------------------------------- bar_chart
def render_bar_chart(slide, s, palette, fonts, motif, image_path=None):
    """Vertikal bar diagrammasi — Y o'qi va X belgili barlar."""
    set_background(slide, "FFFFFF")

    # Sarlavha
    add_text(slide, 0.7, 0.35, 12.0, 0.65, s.title or "", size=26, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])

    chart = s.chart_data
    if not chart or not chart.points:
        add_text(slide, 1.0, 2.0, 11.0, 1.0, "(diagramma ma'lumotlari yo'q)", size=14,
                  color_hex="999999", font=fonts["body"])
        return

    points = chart.points
    n = len(points)
    max_val = max(p.value for p in points) if points else 1
    if max_val == 0:
        max_val = 1

    # Diagramma chegaralari (inch)
    chart_left   = 1.4    # Y o'qi joyi
    chart_right  = 12.6
    chart_top    = 1.35
    chart_bottom = 5.8
    chart_w = chart_right - chart_left
    chart_h = chart_bottom - chart_top

    # Y o'qi chizig'i
    add_rect(slide, chart_left - 0.03, chart_top, 0.03, chart_h, fill_hex="AAAAAA")
    # X o'qi chizig'i
    add_rect(slide, chart_left, chart_bottom, chart_w, 0.03, fill_hex="AAAAAA")

    # Y o'qi sarlavhasi (chapga aylantirilgan matn — textbox burchagi o'zgartirib bo'lmaydi
    # python-pptx da, shuning uchun vertikal qilib yozamiz)
    y_label = f"{chart.y_label}" + (f" ({chart.unit})" if chart.unit else "")
    add_text(slide, 0.05, 2.5, 1.2, 0.5, y_label, size=10, color_hex="555555",
              font=fonts["body"], align=PP_ALIGN.CENTER)

    # Y o'qi graduirovkasi (5 ta chiziq)
    tick_count = 5
    for i in range(tick_count + 1):
        val = max_val * i / tick_count
        y_pos = chart_bottom - chart_h * i / tick_count
        add_rect(slide, chart_left - 0.12, y_pos - 0.01, 0.12, 0.02, fill_hex="CCCCCC")
        label_str = f"{val:.1f}" if val != int(val) else str(int(val))
        if chart.unit:
            label_str += f" {chart.unit}"
        add_text(slide, chart_left - 1.1, y_pos - 0.15, 0.95, 0.3, label_str,
                  size=9, color_hex="555555", font=fonts["body"], align=PP_ALIGN.RIGHT)

    # Barlar
    bar_gap_ratio = 0.25
    total_slots = n
    slot_w = chart_w / total_slots
    bar_w = slot_w * (1 - bar_gap_ratio)
    gap_w = slot_w * bar_gap_ratio / 2

    bar_colors = [palette["primary"], palette["accent"], palette["secondary"]]

    for i, pt in enumerate(points):
        bar_height = chart_h * (pt.value / max_val)
        bx = chart_left + i * slot_w + gap_w
        by = chart_bottom - bar_height
        color = bar_colors[i % len(bar_colors)]
        add_rect(slide, bx, by, bar_w, bar_height, fill_hex=color.lstrip("#"))

        # Qiymat yuqorida
        val_str = f"{pt.value:.1f}" if pt.value != int(pt.value) else str(int(pt.value))
        add_text(slide, bx, by - 0.35, bar_w, 0.3, val_str, size=10, bold=True,
                  color_hex=palette["primary"], font=fonts["body"], align=PP_ALIGN.CENTER)

        # X belgisi
        add_text(slide, bx - 0.1, chart_bottom + 0.08, bar_w + 0.2, 0.5, pt.label,
                  size=9, color_hex="444444", font=fonts["body"], align=PP_ALIGN.CENTER)

    # X o'qi sarlavhasi
    add_text(slide, chart_left, chart_bottom + 0.65, chart_w, 0.4, chart.x_label,
              size=11, bold=True, color_hex="333333", font=fonts["body"], align=PP_ALIGN.CENTER)

    # Izoh (body)
    if s.body:
        add_rect(slide, 0.7, 6.3, 11.9, 0.85, fill_hex="F4F6FA", radius=True)
        add_text(slide, 1.0, 6.35, 11.3, 0.75, s.body, size=11, italic=True,
                  color_hex="555555", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)


# -------------------------------------------------------------- xy_chart
def render_xy_chart(slide, s, palette, fonts, motif, image_path=None):
    """Chiziqli (line) diagramma — X va Y o'qlari bilan, nuqtalar ulangan."""
    set_background(slide, "FFFFFF")

    # Sarlavha
    add_text(slide, 0.7, 0.35, 12.0, 0.65, s.title or "", size=26, bold=True,
              color_hex=palette["primary"], font=fonts["heading"])

    chart = s.chart_data
    if not chart or not chart.points or len(chart.points) < 2:
        add_text(slide, 1.0, 2.0, 11.0, 1.0, "(diagramma ma'lumotlari yo'q)", size=14,
                  color_hex="999999", font=fonts["body"])
        return

    points = chart.points
    n = len(points)
    values = [p.value for p in points]
    max_val = max(values) if values else 1
    min_val = min(values) if values else 0
    val_range = max_val - min_val if max_val != min_val else 1

    # Diagramma chegaralari
    chart_left   = 1.5
    chart_right  = 12.5
    chart_top    = 1.35
    chart_bottom = 5.8
    chart_w = chart_right - chart_left
    chart_h = chart_bottom - chart_top

    # O'qlar
    add_rect(slide, chart_left - 0.03, chart_top, 0.03, chart_h, fill_hex="888888")
    add_rect(slide, chart_left, chart_bottom, chart_w, 0.03, fill_hex="888888")

    # Y o'qi sarlavhasi
    y_label = f"{chart.y_label}" + (f" ({chart.unit})" if chart.unit else "")
    add_text(slide, 0.05, 2.5, 1.3, 0.5, y_label, size=10, color_hex="555555",
              font=fonts["body"], align=PP_ALIGN.CENTER)

    # Y graduirovka
    tick_count = 5
    for i in range(tick_count + 1):
        val = min_val + val_range * i / tick_count
        y_pos = chart_bottom - chart_h * i / tick_count
        # Gorizontal to'r chizig'i
        add_rect(slide, chart_left, y_pos - 0.005, chart_w, 0.01, fill_hex="E8E8E8")
        add_rect(slide, chart_left - 0.1, y_pos - 0.01, 0.1, 0.02, fill_hex="AAAAAA")
        label_str = f"{val:.1f}" if val != int(val) else str(int(val))
        if chart.unit:
            label_str += f" {chart.unit}"
        add_text(slide, chart_left - 1.2, y_pos - 0.15, 1.05, 0.3, label_str,
                  size=9, color_hex="555555", font=fonts["body"], align=PP_ALIGN.RIGHT)

    # Nuqtalar pozitsiyalarini hisoblash
    def get_xy(i, pt):
        px = chart_left + chart_w * i / (n - 1) if n > 1 else chart_left + chart_w / 2
        py = chart_bottom - chart_h * (pt.value - min_val) / val_range
        return px, py

    # Nuqtalar orasidagi chiziqlar (tor to'rtburchak bilan ifodalanadi)
    for i in range(n - 1):
        x1, y1 = get_xy(i, points[i])
        x2, y2 = get_xy(i + 1, points[i + 1])
        dx = x2 - x1
        dy = y2 - y1
        import math
        length = math.sqrt(dx * dx + dy * dy)
        if length < 0.001:
            continue
        # python-pptx da aniq chiziq yo'q — ingichka to'rtburchak bilan taqlidlaymiz
        # faqat gorizontal va vertikal uchun aniq, qiya uchun taxminiy
        mid_x = (x1 + x2) / 2 - 0.02
        mid_y = min(y1, y2) - 0.005
        seg_h = abs(dy) if abs(dy) > 0.01 else 0.01
        seg_w = abs(dx) if abs(dx) > 0.01 else 0.01
        # Oddiyroq: ikkita segment — gorizontal + vertikal stepped line
        add_rect(slide, x1, min(y1, y2) - 0.008, dx if dx > 0 else -dx,
                  0.016, fill_hex=palette["primary"].lstrip("#"))
        add_rect(slide, x2 - 0.008, min(y1, y2) - 0.008, 0.016,
                  dy if dy > 0 else -dy, fill_hex=palette["primary"].lstrip("#"))

    # Nuqtalar va X belgilari
    for i, pt in enumerate(points):
        px, py = get_xy(i, pt)
        # Doira nuqta
        dot_r = 0.12
        add_circle(slide, px - dot_r, py - dot_r, dot_r * 2, palette["accent"].lstrip("#"))
        # Qiymat yuqorida
        val_str = f"{pt.value:.1f}" if pt.value != int(pt.value) else str(int(pt.value))
        add_text(slide, px - 0.4, py - 0.45, 0.8, 0.3, val_str, size=9, bold=True,
                  color_hex=palette["primary"], font=fonts["body"], align=PP_ALIGN.CENTER)
        # X belgisi
        add_text(slide, px - 0.5, chart_bottom + 0.08, 1.0, 0.45, pt.label,
                  size=9, color_hex="444444", font=fonts["body"], align=PP_ALIGN.CENTER)

    # X o'qi sarlavhasi
    add_text(slide, chart_left, chart_bottom + 0.62, chart_w, 0.4, chart.x_label,
              size=11, bold=True, color_hex="333333", font=fonts["body"], align=PP_ALIGN.CENTER)

    # Izoh
    if s.body:
        add_rect(slide, 0.7, 6.3, 11.9, 0.85, fill_hex="F4F6FA", radius=True)
        add_text(slide, 1.0, 6.35, 11.3, 0.75, s.body, size=11, italic=True,
                  color_hex="555555", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)


# -------------------------------------------------------------- math_formula
def render_math_formula(slide, s, palette, fonts, motif, image_path=None):
    """Matematik formulalar va izohlar — taqdimotga ilmiy chuqurlik beradi."""
    set_background(slide, "FFFFFF")

    # Yuqori chiziq va sarlavha
    add_rect(slide, 0.0, 0.0, 13.333, 1.15, fill_hex=palette["primary"])
    add_text(slide, 0.8, 0.25, 11.7, 0.7, s.title or "", size=26, bold=True,
              color_hex="FFFFFF", font=fonts["heading"])

    blocks = s.math_blocks or []
    if not blocks:
        add_text(slide, 1.0, 2.0, 11.3, 1.0, "(formulalar yo'q)", size=14,
                  color_hex="999999", font=fonts["body"])
        return

    n = len(blocks)
    # Har blok uchun balandlik
    avail_h = 5.8  # 1.15 dan 7.0 gacha
    block_h = min(avail_h / n, 2.0)
    y = 1.3

    for i, mb in enumerate(blocks):
        # Formula foni — kontrast blok
        add_rect(slide, 0.8, y, 11.7, block_h - 0.15, fill_hex="F0F4FF", radius=True)

        # Formula matni — katta, monospace kabi bold
        formula_display = mb.formula
        add_text(slide, 1.1, y + 0.12, 5.5, block_h * 0.55,
                  formula_display, size=26, bold=True, font="Courier New",
                  color_hex=palette["primary"], align=PP_ALIGN.LEFT)

        # Vertikal ajratgich
        add_rect(slide, 6.9, y + 0.15, 0.04, block_h - 0.35, fill_hex=palette["accent"])

        # Izoh matni
        add_text(slide, 7.1, y + 0.12, 5.1, block_h - 0.3,
                  mb.description, size=13, color_hex="333333",
                  font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)

        # Raqam badge
        motif_badge(slide, 0.82, y + (block_h - 0.6) / 2, 0.0, i + 1, palette, fonts, motif)

        y += block_h

    # Pastki izoh (body)
    if s.body and y < 6.8:
        add_rect(slide, 0.8, max(y + 0.1, 6.5), 11.7, 0.7, fill_hex="F4F6FA", radius=True)
        add_text(slide, 1.1, max(y + 0.15, 6.55), 11.1, 0.6, s.body, size=11,
                  italic=True, color_hex="555555", font=fonts["body"], anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# LAYOUT REGISTRI
# ============================================================

LAYOUT_REGISTRY = {
    "title_dark":          render_title_dark,
    "conclusion_dark":     render_conclusion_dark,
    "agenda_numbered":     render_agenda_numbered,
    "two_card_compare":    render_two_card_compare,
    "three_card_grid":     render_three_card_grid,
    "four_card_grid":      render_four_card_grid,
    "comparison_table":    render_comparison_table,
    "icon_row_list":       render_icon_row_list,
    "stat_callout":        render_stat_callout,
    "stat_row_triple":     render_stat_row_triple,
    "timeline_process":    render_timeline_process,
    "image_half_bleed":    render_image_half_bleed,
    "quote_callout":       render_quote_callout,
    "definition_spotlight":render_definition_spotlight,
    "mini_case_study":     render_mini_case_study,
    # Yangi diagramma/formula layoutlari
    "bar_chart":           render_bar_chart,
    "xy_chart":            render_xy_chart,
    "math_formula":        render_math_formula,
}

# Agar layout_type ro'yxatda topilmasa yoki rasm talab qilinib topilmasa ishlatiladigan zaxira
FALLBACK_LAYOUT = "icon_row_list"
