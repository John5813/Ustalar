import logging
import os
import uuid

from pptx import Presentation
from pptx.util import Inches

from . import config
from .layouts import LAYOUT_REGISTRY, FALLBACK_LAYOUT, SLIDE_W, SLIDE_H
from .models import Brief
from .image_client import generate_image

log = logging.getLogger("renderer")


def build_presentation(brief: Brief) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # bo'sh layout — o'zimiz to'liq chizamiz

    palette = brief.palette.model_dump()
    fonts = brief.font_pair.model_dump()
    motif = brief.motif

    for s in brief.slides:
        slide = prs.slides.add_slide(blank_layout)

        image_path = None
        if s.image_prompt:
            image_path = generate_image(s.image_prompt)  # xato bo'lsa None qaytadi, try/catch ichida

        render_fn = LAYOUT_REGISTRY.get(s.layout_type)
        if render_fn is None:
            log.warning("Noma'lum layout_type='%s' (slayd %s), fallback ishlatiladi", s.layout_type, s.index)
            render_fn = LAYOUT_REGISTRY[FALLBACK_LAYOUT]

        # Agar image_half_bleed uchun rasm topilmasa, matn-asosli layoutga tushamiz
        if s.layout_type == "image_half_bleed" and image_path is None:
            log.info("Slayd %s uchun rasm topilmadi, icon_row_list'ga o'tildi", s.index)
            render_fn = LAYOUT_REGISTRY[FALLBACK_LAYOUT]

        render_fn(slide, s, palette, fonts, motif, image_path=image_path)

    os.makedirs(config.WORK_DIR, exist_ok=True)
    out_path = os.path.join(config.WORK_DIR, f"presentation_{uuid.uuid4().hex[:10]}.pptx")
    prs.save(out_path)
    return out_path
