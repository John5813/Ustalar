import os
import logging
import re
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from typing import Dict, Optional, List
import asyncio
from config import DOCUMENTS_DIR, TEMP_DIR
from services.together_service import get_together_service
from services.ai_service import clean_text
from services.icon_service import find_icon_path_for_column

logger = logging.getLogger(__name__)


def _split_into_paragraphs(text: str, target_count: int = 2, min_sentences: int = 6) -> list:
    """Split text into ~target_count paragraphs at safe sentence boundaries.

    Splits only at `. `/`! `/`? ` followed by an uppercase letter (Latin/Cyrillic/Uzbek O'/G'),
    so abbreviations like "Resp." or "T.B." don't break a paragraph mid-sentence.
    Splits only when there are MORE THAN 5 sentences (min_sentences=6 by default);
    otherwise returns the whole text as one paragraph.
    """
    text = (text or "").strip()
    if not text:
        return []
    # Split at sentence boundary followed by capital letter (covers Latin, Cyrillic, Uzbek apostrophes)
    sentences = re.split(r"(?<=[.!?])\s+(?=[A-ZА-ЯЁЎҚҒҲ])", text)
    if len(sentences) < min_sentences or target_count <= 1:
        return [text]
    chunk = max(1, len(sentences) // target_count)
    paragraphs = []
    for k in range(target_count):
        start = k * chunk
        end = (k + 1) * chunk if k < target_count - 1 else len(sentences)
        para = " ".join(sentences[start:end]).strip()
        if para:
            paragraphs.append(para)
    return paragraphs or [text]


def _extras_for_cycle(extras: list, section_num: int) -> list:
    """Return the subset of extras for this section based on a 3-step cycle.

    Cycle (1-based, repeating):
      pos 1 (sections 1, 4, 7 …) → formulas only
      pos 2 (sections 2, 5, 8 …) → images + tables
      pos 3 (sections 3, 6, 9 …) → tables only
    Extras not in the cycle (e.g. 'glossary', 'statistics') are ignored here
    because glossary is added at document end and statistics follows tables.
    """
    pos = ((section_num - 1) % 3) + 1
    if pos == 1:
        return [e for e in extras if e == "formulas"]
    elif pos == 2:
        return [e for e in extras if e in ("images", "tables", "statistics")]
    else:
        return [e for e in extras if e in ("tables", "statistics")]


class DocumentService:
    def __init__(self):
        self.documents_dir = DOCUMENTS_DIR
        self.temp_dir = TEMP_DIR
        self._last_used_icons: set = set()
        try:
            self.together = get_together_service()
        except Exception as e:
            logger.warning(f"Together AI not available: {e}")
            self.together = None
        
        # Verify bayoo-docx footnote support is available
        try:
            test_doc = Document()
            test_para = test_doc.add_paragraph("test")
            self.footnotes_available = hasattr(test_para, 'add_footnote')
            if self.footnotes_available:
                logger.info("Word-native footnotes (snoska) support is available via bayoo-docx")
            else:
                logger.warning("Word-native footnotes not available - will use inline citations as fallback")
        except Exception as e:
            logger.warning(f"Could not verify footnote support: {e}")
            self.footnotes_available = False

    def _calculate_auto_font_size(self, text: str, width_inches: float, height_inches: float, 
                                   max_font_pt: int = 24, min_font_pt: int = 14) -> int:
        """Calculate optimal font size to fit text within boundaries.
        
        Accounts for:
        - Character count per line based on font size
        - Explicit line breaks (\\n)
        - Word wrapping estimation
        - Long words that may not wrap properly
        """
        if not text:
            return max_font_pt
        
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        paragraphs = text.split('\n')
        
        current_font = max_font_pt
        while current_font >= min_font_pt:
            chars_per_inch = 8.0 * (24 / current_font)
            chars_per_line = int(width_inches * chars_per_inch)
            
            line_height_inches = current_font / 72 * 1.5
            max_lines = int(height_inches / line_height_inches)
            
            total_lines_needed = 0
            for para in paragraphs:
                if not para.strip():
                    total_lines_needed += 1
                    continue
                
                words = para.split()
                line_chars = 0
                para_lines = 1
                
                for word in words:
                    word_len = len(word)
                    if word_len > chars_per_line:
                        para_lines += (word_len // chars_per_line) + 1
                        line_chars = word_len % chars_per_line
                    elif line_chars + word_len + 1 <= chars_per_line:
                        line_chars += word_len + 1
                    else:
                        para_lines += 1
                        line_chars = word_len
                
                total_lines_needed += para_lines
            
            if total_lines_needed <= max_lines:
                return current_font
            
            current_font -= 1
        
        return min_font_pt

    async def _prefetch_slide_images(self, slides_data: list, topic: str, language: str) -> dict:
        """Pre-generate slide images with bounded parallelism (max 4 concurrent).

        Unbounded `asyncio.gather` over 10–20 slides used to spike memory and
        trip Together AI's concurrency limits. A semaphore keeps at most 4
        in-flight requests at once while preserving the parallelism win.
        """
        if not self.together:
            return {}

        sem = asyncio.Semaphore(4)

        async def _bounded(coro):
            async with sem:
                return await coro

        tasks = {}
        for i, slide_data in enumerate(slides_data):
            layout = slide_data.get('layout', 'text')
            if layout == 'cover':
                tasks[f'cover_{i}'] = (i, 'cover', _bounded(self.together.generate_cover_image(topic, language)))
            elif layout in ('right_image', 'left_image'):
                tasks[f'slide_{i}'] = (i, 'slide', _bounded(self.together.generate_slide_image(
                    topic, slide_data.get('title', ''), language)))
            elif layout == 'horizontal_image':
                tasks[f'horiz_{i}'] = (i, 'horiz', _bounded(self.together.generate_panoramic_image(
                    topic, slide_data.get('title', ''), language)))

        keys = list(tasks.keys())
        coros = [tasks[k][2] for k in keys]
        results = await asyncio.gather(*coros, return_exceptions=True)

        image_map = {}
        for key, result in zip(keys, results):
            if isinstance(result, Exception):
                logger.error(f"Image prefetch failed for {key}: {result}")
                continue
            slide_idx = tasks[key][0]
            image_map[slide_idx if key != 'infographic' else 'infographic'] = result

        logger.info(f"Prefetched {len(image_map)} images (max 4 concurrent)")
        return image_map

    async def create_presentation_with_smart_images(self, topic: str, content: Dict, author_name: str, language: str = "uz", template_service=None, template_id: str = None) -> str:
        """Create PowerPoint presentation with new layout system and Together AI images

        NEW STRUCTURE:
        1. Muqova - O'ng: Mavzu + Ism, Chap: 50% rasm
        2. Reja - 4 asosiy punkt
        3. Kirish - ~50 so'z
        4+. Asosiy slaidlar
        N+1. Xulosa
        N+2. Adabiyotlar
        N+3. Rahmat
        """
        try:
            prs = Presentation()
            prs.slide_width = PptxInches(13.333)
            prs.slide_height = PptxInches(7.5)

            slides_data = content.get('slides', [])

            # Pre-generate ALL images in parallel before building slides
            logger.info("Pre-fetching all slide images in parallel...")
            image_map = await self._prefetch_slide_images(slides_data, topic, language)

            used_icons: set[str] = set()  # tracks icon filenames used across the whole presentation

            try:
                for i, slide_data in enumerate(slides_data):
                    await asyncio.sleep(0)  # yield to event loop between slides
                    pre_img = image_map.get(i)
                    await self._create_slide_by_layout(prs, slide_data, i, author_name, topic, language, template_service, template_id, pre_fetched_image=pre_img, used_icons=used_icons)
            finally:
                # Always clean up prefetched images, even if slide build crashes,
                # so we don't leak temp files when generation fails mid-way.
                for img_val in image_map.values():
                    img_path = img_val[0] if isinstance(img_val, tuple) else img_val
                    if img_path and isinstance(img_path, str) and os.path.exists(img_path):
                        try:
                            os.remove(img_path)
                        except Exception:
                            pass

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(prs.save, file_path)
            logger.info(f"Presentation saved: {file_path}")

            # Log icon usage summary
            self._last_used_icons = used_icons
            if used_icons:
                logger.info(f"Icon summary: {len(used_icons)} unique icon(s) used — {', '.join(sorted(used_icons))}")
            else:
                logger.info("Icon summary: no icons were added to this presentation")

            return file_path

        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise

    async def _create_fullpage_infographic_slide(self, prs, topic: str, language: str, template_service=None, template_id: str = None, pre_fetched_url: str = None):
        """Create a full-page infographic slide using nano-banana-2 between slides 4 and 5."""
        try:
            import aiohttp, time

            image_url = pre_fetched_url

            if not image_url:
                from services.fal_service import generate_image_nano
                lang_name = {'uz': 'Uzbek', 'ru': 'Russian', 'en': 'English'}.get(language, 'Uzbek')
                prompt = (
                    f"Professional educational infographic poster clearly explaining '{topic}'. "
                    f"Include text labels, key terms, and annotations in {lang_name} language. "
                    "Data charts, statistics, diagrams, icons, flowcharts, key concepts all related to this specific topic. "
                    "Colorful modern academic design, high quality, wide landscape format."
                )
                logger.info(f"Generating full-page infographic slide for: {topic}")
                image_url = await generate_image_nano(prompt, aspect_ratio="16_9")

            if not image_url:
                logger.warning("No infographic image URL available, skipping slide")
                return

            logger.info(f"Infographic slide image URL: {image_url[:80] if image_url else 'None'}")

            filename = f"pres_infographic_{int(time.time())}.png"
            os.makedirs(self.temp_dir, exist_ok=True)
            filepath = os.path.join(self.temp_dir, filename)

            async with aiohttp.ClientSession() as session:
                async with session.get(image_url, timeout=aiohttp.ClientTimeout(total=60)) as resp:
                    if resp.status != 200:
                        logger.warning(f"Could not download infographic slide image: HTTP {resp.status}")
                        return
                    with open(filepath, "wb") as f:
                        f.write(await resp.read())

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if template_service and template_id:
                try:
                    template_service.apply_template_to_slide(slide, template_id)
                except Exception:
                    pass

            # Fill entire slide with the image
            slide.shapes.add_picture(
                filepath,
                PptxInches(0), PptxInches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            logger.info(f"Full-page infographic slide added to presentation")

        except Exception as e:
            logger.warning(f"Could not create infographic slide: {e}", exc_info=True)

    async def _create_slide_by_layout(self, prs, slide_data: Dict, slide_idx: int, author_name: str, topic: str, language: str, template_service=None, template_id: str = None, pre_fetched_image=None, used_icons: set | None = None):
        """Create slide based on layout type"""
        layout = slide_data.get('layout', 'text_only')
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        if template_service and template_id:
            template_service.apply_template_to_slide(slide, template_id)
        
        if used_icons is None:
            used_icons = set()

        if layout == 'cover':
            await self._create_cover_slide(slide, slide_data, author_name, topic, language, pre_fetched_image=pre_fetched_image)
        elif layout == 'plan':
            self._create_plan_slide(slide, slide_data)
        elif layout == 'intro':
            self._create_intro_slide(slide, slide_data)
        elif layout == 'two_column':
            self._create_two_column_slide(slide, slide_data, language, used_icons=used_icons)
        elif layout == 'right_image':
            await self._create_right_image_slide(slide, slide_data, topic, language, pre_fetched_image=pre_fetched_image)
        elif layout == 'left_image':
            await self._create_left_image_slide(slide, slide_data, topic, language, pre_fetched_image=pre_fetched_image)
        elif layout == 'three_column':
            self._create_three_column_slide(slide, slide_data, language, used_icons=used_icons)
        elif layout == 'horizontal_image':
            await self._create_horizontal_image_slide(slide, slide_data, topic, language, pre_fetched_image=pre_fetched_image)
        elif layout == 'text_with_numbers':
            self._create_text_with_numbers_slide(slide, slide_data)
        elif layout == 'conclusion':
            self._create_conclusion_slide(slide, slide_data)
        elif layout == 'references':
            self._create_references_slide(slide, slide_data)
        elif layout == 'thanks':
            self._create_thanks_slide(slide, language)
        elif layout == 'table':
            await self._create_table_slide(slide, slide_data, topic, language)
        else:
            self._create_default_slide(slide, slide_data)

    async def _create_cover_slide(self, slide, slide_data: Dict, author_name: str, topic: str, language: str, pre_fetched_image=None):
        """1-varoq: Chap 50% rasm, O'ng mavzu + ism"""
        if self.together:
            try:
                image_path = pre_fetched_image if pre_fetched_image else await self.together.generate_cover_image(topic, language)
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(
                        image_path,
                        PptxInches(0), PptxInches(0),
                        PptxInches(6.666), PptxInches(7.5)
                    )
                    if not pre_fetched_image:
                        try:
                            os.remove(image_path)
                        except Exception:
                            pass
            except Exception as e:
                logger.error(f"Error generating cover image: {e}")
        
        title_box = slide.shapes.add_textbox(
            PptxInches(7), PptxInches(2.5),
            PptxInches(5.8), PptxInches(3)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p1 = tf.paragraphs[0]
        p1.text = topic
        cover_font_size = self._calculate_auto_font_size(
            topic,
            width_inches=5.8,
            height_inches=2.2,
            max_font_pt=42,
            min_font_pt=18
        )
        p1.font.size = PptxPt(cover_font_size)
        p1.font.bold = True
        p1.font.name = 'Times New Roman'
        p1.font.color.rgb = RGBColor(0, 0, 0)
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = author_name if author_name else "________________"
        p2.font.size = PptxPt(26)
        p2.font.bold = True
        p2.font.name = 'Times New Roman'
        p2.alignment = PP_ALIGN.CENTER

    def _create_plan_slide(self, slide, slide_data: Dict):
        """2-varoq: Reja - 4 asosiy punkt"""
        self._add_slide_title(slide, slide_data.get('title', 'Reja'))
        
        plan_items = slide_data.get('plan_items', [])
        width_in = 11
        height_in = 5
        content_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(2),
            PptxInches(width_in), PptxInches(height_in)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        all_text = ' '.join(plan_items[:4])
        optimal_font = self._calculate_auto_font_size(all_text, width_in, height_in, 26, 18)
        
        for i, item in enumerate(plan_items[:4]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(item, str):
                p.text = item if item.startswith(str(i+1)) else f"{i+1}. {item}"
            else:
                p.text = f"{i+1}. {str(item)}"
            p.font.size = PptxPt(optimal_font)
            p.font.bold = True
            p.font.name = 'Times New Roman'
            p.alignment = PP_ALIGN.LEFT
            p.space_after = PptxPt(24)

    def _create_intro_slide(self, slide, slide_data: Dict):
        """3-varoq: Kirish - ~50 so'z"""
        self._add_slide_title(slide, slide_data.get('title', 'Kirish'))
        self._add_justified_content(slide, slide_data.get('content', ''), 
                                    PptxInches(1), PptxInches(2), 
                                    PptxInches(11), PptxInches(5))

    def _create_two_column_slide(self, slide, slide_data: Dict, language: str = 'uz', used_icons: set | None = None):
        """Shablon 1: 2 ustunli - har ustun 30 so'z + ikonka"""
        if used_icons is None:
            used_icons = set()
        self._add_slide_title(slide, slide_data.get('title', ''))

        columns = slide_data.get('columns', [])
        content = slide_data.get('content', '')

        if isinstance(content, dict):
            content = content.get('text', content.get('content', str(content)))
        if isinstance(content, list):
            content = ' '.join(str(item) for item in content)
        if not isinstance(content, str):
            content = str(content)

        if not columns and content:
            logger.warning("two_column slide missing 'columns' array - using content as fallback")
            columns = [
                {'text': content},
                {'text': ''}
            ]

        max_font = 23 if language in ['ru', 'en'] else 24
        width_in = 5.8
        icon_size = 1.0        # icon width & height in inches
        icon_y = 1.40          # icon top position
        text_y = 2.55          # text box top (shifted down to make room for icon)
        height_in = 3.9        # text box height (reduced)
        column_positions = [0.5, 6.8]

        for i, col in enumerate(columns[:2]):
            x_start = column_positions[i]
            x_pos = PptxInches(x_start)

            if isinstance(col, dict):
                col_keyword = col.get('keyword', '')
                col_text = col.get('column_content', col.get('text', col.get('content', '')))
            else:
                col_keyword = ''
                col_text = str(col)

            col_text = clean_text(col_text)

            # --- Icon (deduplicated across the whole presentation) ---
            if getattr(self, 'use_icons', True):
                try:
                    icon_path = find_icon_path_for_column(col_keyword, col_text, used=used_icons)
                    if icon_path and os.path.isfile(icon_path):
                        icon_x = PptxInches(x_start + (width_in - icon_size) / 2)
                        slide.shapes.add_picture(
                            icon_path,
                            icon_x, PptxInches(icon_y),
                            PptxInches(icon_size), PptxInches(icon_size)
                        )
                        used_icons.add(os.path.basename(icon_path))
                except Exception as e:
                    logger.warning(f"Could not add icon to two_column slide: {e}")

            # --- Text box ---
            box = slide.shapes.add_textbox(x_pos, PptxInches(text_y), PptxInches(width_in), PptxInches(height_in))
            tf = box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            optimal_font = self._calculate_auto_font_size(col_text, width_in, height_in, max_font, 12)
            p.text = col_text
            p.font.size = PptxPt(optimal_font)
            p.font.name = 'Times New Roman'
            p.alignment = PP_ALIGN.LEFT

    async def _create_right_image_slide(self, slide, slide_data: Dict, topic: str, language: str, pre_fetched_image=None):
        """Shablon 2: O'ng 50% rasm, chap matn"""
        self._add_slide_title(slide, slide_data.get('title', ''))
        
        self._add_justified_content(slide, slide_data.get('content', ''),
                                    PptxInches(0.5), PptxInches(2),
                                    PptxInches(5.8), PptxInches(4.5), align_left=True)
        
        if self.together:
            try:
                image_path = pre_fetched_image if pre_fetched_image else await self.together.generate_slide_image(
                    topic, slide_data.get('title', ''), language
                )
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(
                        image_path,
                        PptxInches(6.8), PptxInches(1.5),
                        PptxInches(6.2), PptxInches(5.5)
                    )
                    if not pre_fetched_image:
                        try:
                            os.remove(image_path)
                        except Exception:
                            pass
            except Exception as e:
                logger.error(f"Error generating right image: {e}")

    async def _create_left_image_slide(self, slide, slide_data: Dict, topic: str, language: str, pre_fetched_image=None):
        """Shablon 3: Chap 50% rasm, o'ng matn"""
        self._add_slide_title(slide, slide_data.get('title', ''))
        
        if self.together:
            try:
                image_path = pre_fetched_image if pre_fetched_image else await self.together.generate_slide_image(
                    topic, slide_data.get('title', ''), language
                )
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(
                        image_path,
                        PptxInches(0.3), PptxInches(1.5),
                        PptxInches(6.2), PptxInches(5.5)
                    )
                    if not pre_fetched_image:
                        try:
                            os.remove(image_path)
                        except Exception:
                            pass
            except Exception as e:
                logger.error(f"Error generating left image: {e}")
        
        self._add_justified_content(slide, slide_data.get('content', ''),
                                    PptxInches(6.8), PptxInches(2),
                                    PptxInches(5.8), PptxInches(4.5), align_left=True)

    def _create_three_column_slide(self, slide, slide_data: Dict, language: str = 'uz', used_icons: set | None = None):
        """Shablon 4: 3 ustunli - har ustunda kalit so'z + tarif + ikonka"""
        if used_icons is None:
            used_icons = set()
        self._add_slide_title(slide, slide_data.get('title', ''))

        columns = slide_data.get('columns', [])
        content = slide_data.get('content', '')

        if isinstance(content, dict):
            content = content.get('text', content.get('content', str(content)))
        if isinstance(content, list):
            content = ' '.join(str(item) for item in content)
        if not isinstance(content, str):
            content = str(content)

        if not columns and content:
            logger.warning("three_column slide missing 'columns' array - using content as fallback")
            columns = [
                {'column_content': content},
                {'column_content': ''},
                {'column_content': ''}
            ]

        max_font = 22 if language in ['ru', 'en'] else 23
        width_in = 4.0
        icon_size = 0.7        # smaller icon for 3-column layout
        icon_y = 1.45
        text_y = 2.40
        height_in = 4.0        # reduced height to fit icon above

        column_positions = [0.3, 4.5, 8.7]

        for i, col in enumerate(columns[:3]):
            x_start = column_positions[i]
            x_pos = PptxInches(x_start)

            if isinstance(col, dict):
                keyword = col.get('keyword', '')
                col_text = col.get('column_content', col.get('text', col.get('content', '')))
            else:
                keyword = ''
                col_text = str(col)

            keyword = clean_text(keyword)
            col_text = clean_text(col_text)

            # --- Icon (deduplicated across the whole presentation) ---
            if getattr(self, 'use_icons', True):
                try:
                    icon_path = find_icon_path_for_column(keyword, col_text, used=used_icons)
                    if icon_path and os.path.isfile(icon_path):
                        icon_x = PptxInches(x_start + (width_in - icon_size) / 2)
                        slide.shapes.add_picture(
                            icon_path,
                            icon_x, PptxInches(icon_y),
                            PptxInches(icon_size), PptxInches(icon_size)
                        )
                        used_icons.add(os.path.basename(icon_path))
                except Exception as e:
                    logger.warning(f"Could not add icon to three_column slide: {e}")

            # --- Text box ---
            box = slide.shapes.add_textbox(x_pos, PptxInches(text_y), PptxInches(width_in), PptxInches(height_in))
            tf = box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            optimal_font = self._calculate_auto_font_size(col_text, width_in, height_in - 1, max_font, 12)

            if keyword:
                p_keyword = tf.paragraphs[0]
                p_keyword.text = keyword
                keyword_font = self._calculate_auto_font_size(
                    keyword, width_inches=width_in, height_inches=0.8,
                    max_font_pt=26, min_font_pt=16
                )
                p_keyword.font.size = PptxPt(keyword_font)
                p_keyword.font.bold = True
                p_keyword.font.name = 'Times New Roman'
                p_keyword.alignment = PP_ALIGN.LEFT

                p_desc = tf.add_paragraph()
                p_desc.text = col_text
                p_desc.font.size = PptxPt(optimal_font)
                p_desc.font.name = 'Times New Roman'
                p_desc.alignment = PP_ALIGN.LEFT
            else:
                p = tf.paragraphs[0]
                p.text = col_text
                p.font.size = PptxPt(optimal_font)
                p.font.name = 'Times New Roman'
                p.alignment = PP_ALIGN.LEFT

    async def _create_horizontal_image_slide(self, slide, slide_data: Dict, topic: str, language: str, pre_fetched_image=None):
        """Shablon 5: Pastda 21:9 gorizontal rasm, ustida matn"""
        self._add_slide_title(slide, slide_data.get('title', ''))
        
        self._add_justified_content(slide, slide_data.get('content', ''),
                                    PptxInches(0.5), PptxInches(1.5),
                                    PptxInches(12), PptxInches(2))
        
        if self.together:
            try:
                image_path = pre_fetched_image if pre_fetched_image else await self.together.generate_panoramic_image(
                    topic, slide_data.get('title', ''), language
                )
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(
                        image_path,
                        PptxInches(0.3), PptxInches(4),
                        PptxInches(12.7), PptxInches(3.3)
                    )
                    if not pre_fetched_image:
                        try:
                            os.remove(image_path)
                        except Exception:
                            pass
            except Exception as e:
                logger.error(f"Error generating horizontal image: {e}")

    def _create_text_with_numbers_slide(self, slide, slide_data: Dict):
        """Shablon 6: Oddiy matn, raqamlar bilan - 50 so'z"""
        self._add_slide_title(slide, slide_data.get('title', ''))
        self._add_justified_content(slide, slide_data.get('content', ''),
                                    PptxInches(1), PptxInches(2),
                                    PptxInches(11), PptxInches(5))

    def _create_conclusion_slide(self, slide, slide_data: Dict):
        """Xulosa slayd - ~50 so'z"""
        self._add_slide_title(slide, slide_data.get('title', 'Xulosa'))
        self._add_justified_content(slide, slide_data.get('content', ''),
                                    PptxInches(1), PptxInches(2),
                                    PptxInches(11), PptxInches(5))

    def _create_references_slide(self, slide, slide_data: Dict):
        """Adabiyotlar ro'yxati slayd"""
        self._add_slide_title(slide, slide_data.get('title', 'Foydalangan adabiyotlar'))
        
        references = slide_data.get('references', [])
        if isinstance(references, str):
            references = [references]
        
        content_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2),
            PptxInches(12), PptxInches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        ref_strings = []
        for ref in references[:4]:
            if isinstance(ref, dict):
                ref_strings.append(ref.get('text', ref.get('title', ref.get('source', str(ref)))))
            else:
                ref_strings.append(str(ref))

        all_refs_text = '\n'.join(f"{i+1}. {r}" for i, r in enumerate(ref_strings))
        ref_font_size = self._calculate_auto_font_size(all_refs_text, width_inches=12, height_inches=5, max_font_pt=24, min_font_pt=14)

        for i, ref_text in enumerate(ref_strings):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i+1}. {ref_text}" if not ref_text.startswith(str(i+1)) else ref_text
            p.font.size = PptxPt(ref_font_size)
            p.font.name = 'Times New Roman'
            p.alignment = PP_ALIGN.LEFT
            p.space_after = PptxPt(12)

    def _create_thanks_slide(self, slide, language: str):
        """Rahmat slayd - E'tiboringiz uchun rahmat!"""
        thanks_texts = {
            'uz': "E'tiboringiz uchun rahmat!",
            'ru': "Спасибо за внимание!",
            'en': "Thank you for your attention!"
        }
        
        text = thanks_texts.get(language, thanks_texts['uz'])
        
        thanks_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(3),
            PptxInches(11), PptxInches(2)
        )
        tf = thanks_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptxPt(48)
        p.font.bold = True
        p.font.name = 'Times New Roman'
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    async def _create_table_slide(self, slide, slide_data: Dict, topic: str, language: str):
        """Create table slide with white background table (no caption)"""
        from pptx.util import Inches as PptxInches, Pt as PptxPt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Add slide title
        title = slide_data.get('title', '')
        if not title:
            if language == 'uz':
                title = f"{topic} tahlili"
            elif language == 'ru':
                title = f"Анализ: {topic}"
            else:
                title = f"Analysis: {topic}"
        
        self._add_slide_title(slide, title)
        
        # Get table data - generate from AI if not provided
        table_data = slide_data.get('table_data', {})
        if not table_data or not table_data.get('rows'):
            try:
                from services.ai_service import get_ai_service
                ai_service = get_ai_service()
                table_data = await ai_service.generate_table_data(topic, 1, language)
            except Exception as e:
                logger.warning(f"Could not generate AI table data: {e}")
                table_data = {}
        
        if isinstance(table_data, dict):
            headers = table_data.get('headers', ['Ko\'rsatkich', '2023', '2024', 'O\'zgarish'])
            rows = table_data.get('rows', [])
        else:
            headers = ['Ko\'rsatkich', '2023', '2024', 'O\'zgarish']
            rows = []
        
        if not rows:
            rows = [
                ['Ma\'lumot 1', '100', '120', '+20%'],
                ['Ma\'lumot 2', '50', '65', '+30%'],
                ['Ma\'lumot 3', '200', '180', '-10%'],
                ['Ma\'lumot 4', '75', '90', '+20%'],
                ['Ma\'lumot 5', '150', '175', '+17%']
            ]
        
        # Create table (centered on slide)
        num_rows = min(len(rows) + 1, 6)  # Max 5 data rows + header
        num_cols = len(headers)
        
        table_width = PptxInches(10)
        table_height = PptxInches(4)
        left = PptxInches(1.666)
        top = PptxInches(2)
        
        # Auto-fit font size: reduce by 1pt if text would overflow cell
        col_width_inch = 10.0 / num_cols
        row_height_inch = 4.0 / num_rows
        all_texts = [str(h) for h in headers[:num_cols]]
        for row in rows[:num_rows - 1]:
            all_texts.extend(str(c) for c in row[:num_cols])
        max_len = max((len(t) for t in all_texts), default=1)

        font_size = 14
        while font_size > 8:
            char_w = font_size * 0.007
            line_h = font_size * 0.016
            chars_per_line = max(1, col_width_inch / char_w)
            lines_fit = max(1, row_height_inch / line_h)
            if max_len <= chars_per_line * lines_fit:
                break
            font_size -= 1

        table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
        
        # Style the table - white background, black borders
        from pptx.oxml import parse_xml
        
        for row_idx in range(num_rows):
            for col_idx in range(num_cols):
                cell = table.cell(row_idx, col_idx)
                # White background
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                # Add black borders to cell
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                    border_xml = f'''
                    <a:{border_name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="12700" cap="flat" cmpd="sng">
                        <a:solidFill><a:srgbClr val="000000"/></a:solidFill>
                    </a:{border_name}>
                    '''
                    border_elem = parse_xml(border_xml)
                    existing = tcPr.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{border_name}')
                    if existing is not None:
                        tcPr.remove(existing)
                    tcPr.append(border_elem)
                
                # Set text
                if row_idx == 0:
                    cell.text = str(headers[col_idx]) if col_idx < len(headers) else ''
                else:
                    data_row = rows[row_idx - 1] if row_idx - 1 < len(rows) else []
                    cell.text = str(data_row[col_idx]) if col_idx < len(data_row) else ''
                
                # Format text with auto-fitted font size
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = PptxPt(font_size)
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        if row_idx == 0:
                            run.font.bold = True

    def _create_default_slide(self, slide, slide_data: Dict):
        """Default text slide"""
        self._add_slide_title(slide, slide_data.get('title', ''))
        self._add_justified_content(slide, slide_data.get('content', ''),
                                    PptxInches(1), PptxInches(2),
                                    PptxInches(11), PptxInches(5))

    def _add_slide_title(self, slide, title: str):
        """Add title to slide - qalin qora, auto-size 42→18pt Times New Roman"""
        title_box = slide.shapes.add_textbox(
            PptxInches(0.3), PptxInches(0.3),
            PptxInches(12.7), PptxInches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.text = title
        title_font_size = self._calculate_auto_font_size(
            title,
            width_inches=12.7,
            height_inches=1.0,
            max_font_pt=42,
            min_font_pt=18
        )
        p.font.size = PptxPt(title_font_size)
        p.font.bold = True
        p.font.name = 'Times New Roman'
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    def _add_justified_content(self, slide, content, left: float, top: float, width: float, height: float, align_left: bool = False, max_font: int = 24, min_font: int = 14):
        """Add justified content text with auto-fit font sizing.
        
        Uses Times New Roman font and LEFT alignment to prevent text gaps.
        If text is too long, font automatically reduces (min 14pt).
        """
        if isinstance(content, list):
            content = ' '.join(str(item) for item in content)
        elif not isinstance(content, str):
            content = str(content) if content else ''
        
        width_inches = width / 914400 if width > 100 else width
        height_inches = height / 914400 if height > 100 else height
        
        optimal_font = self._calculate_auto_font_size(content, width_inches, height_inches, max_font, min_font)
            
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = PptxPt(optimal_font)
        p.font.name = 'Times New Roman'
        p.alignment = PP_ALIGN.LEFT

    async def create_presentation(self, topic: str, content: Dict, images: Dict, author_name: str) -> str:
        """Legacy method - redirects to new method"""
        return await self.create_presentation_with_smart_images(topic, content, author_name)

    async def create_presentation_with_layouts(self, topic: str, content: Dict, author_name: str) -> str:
        """Legacy method - redirects to new method"""
        return await self.create_presentation_with_smart_images(topic, content, author_name)

    async def create_presentation_from_template(self, topic: str, content: Dict, author_name: str, template_path: str) -> str:
        """Create presentation from template - falls back to standard creation"""
        return await self.create_presentation_with_smart_images(topic, content, author_name)

    async def _add_section_extras(
        self,
        doc,
        section_title: str,
        topic: str,
        lang: str,
        extras: list,
        formula_data: dict | None = None,
        section_idx: int = 0,
    ) -> None:
        """Add selected extras (image, formulas, stats, table) after a section.
        formula_data: pre-fetched formula dict; fetched here if None and 'formulas' in extras.
        section_idx: used to alternate image type (even=infographic, odd=scene).
        """
        from services.ai_service import get_ai_service
        from services.together_service import get_together_service
        import io as _io
        from docx.shared import Inches as _Inches

        if not extras:
            return

        ai = get_ai_service()

        # ── Pre-fetch formulas ───────────────────────────────────────────
        if "formulas" in extras and formula_data is None:
            formula_data = await ai.generate_section_formulas(section_title, topic, lang)

        # ══ ORDER: image1(infographic) → image2(scene) → formulas+masala → tables → statistics ══

        async def _add_bridge(block_type: str) -> None:
            text = await ai.generate_bridge_sentence(block_type, section_title, topic, lang)
            p = doc.add_paragraph()
            p.paragraph_format.first_line_indent = Inches(0.3)
            p.paragraph_format.space_before = Pt(4)
            p.paragraph_format.space_after = Pt(2)
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            r = p.add_run(text)
            r.font.size = Pt(13)
            r.font.name = "Times New Roman"

        async def _embed_image(img_path: str, caption: str) -> None:
            with open(img_path, "rb") as _f:
                img_bytes = _f.read()
            img_para = doc.add_paragraph()
            img_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            img_para.paragraph_format.space_before = Pt(6)
            img_para.paragraph_format.space_after = Pt(2)
            img_para.paragraph_format.line_spacing = 1.0
            img_run = img_para.add_run()
            img_run.add_picture(_io.BytesIO(img_bytes), width=_Inches(5.5))
            cap_para = doc.add_paragraph()
            cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            cap_para.paragraph_format.space_before = Pt(0)
            cap_para.paragraph_format.space_after = Pt(6)
            cap_para.paragraph_format.line_spacing = 1.0
            cap_run = cap_para.add_run(caption)
            cap_run.font.size = Pt(11)
            cap_run.font.italic = True
            cap_run.font.name = "Times New Roman"

        # ── 1. Image (alternates: even sections → infographic, odd → scene) ─
        if "images" in extras:
            img_type = "infographic" if section_idx % 2 == 0 else "scene"
            if lang == "ru":
                cap = f"Рис. {section_title}"
            elif lang == "en":
                cap = f"Fig. {section_title}"
            else:
                cap = f"{'Infografika' if img_type == 'infographic' else 'Rasm'}. {section_title}"
            bridge_key = "before_image1" if img_type == "infographic" else "before_image2"
            try:
                together = get_together_service()
                img_path, _ = await together.generate_flux_pro_image(
                    topic, section_title, lang, image_type=img_type
                )
                if img_path and os.path.exists(img_path):
                    await _add_bridge(bridge_key)
                    await _embed_image(img_path, cap)
                    try:
                        os.remove(img_path)
                    except Exception:
                        pass
            except Exception as img_err:
                logger.warning(f"Could not embed {img_type} image: {img_err}")

        # ── 2. Formulas ───────────────────────────────────────────────────
        def _render_latex(latex_str: str):
            """Render a LaTeX math string to PNG bytes using matplotlib."""
            fig = None
            plt = None
            try:
                import matplotlib
                matplotlib.use("Agg")
                import matplotlib.pyplot as plt
                import io as _mpl_io
                fig = plt.figure(figsize=(6, 0.7))
                fig.patch.set_facecolor("white")
                ax = fig.add_axes([0, 0, 1, 1])
                ax.set_axis_off()
                ax.text(
                    0.5, 0.5,
                    f"${latex_str}$",
                    ha="center", va="center",
                    fontsize=18,
                    transform=ax.transAxes,
                    color="black",
                )
                buf = _mpl_io.BytesIO()
                fig.savefig(buf, format="png", dpi=130, bbox_inches="tight",
                            facecolor="white", edgecolor="none")
                buf.seek(0)
                data = buf.read()
                buf.close()
                return data
            except Exception as _e:
                logger.warning(f"LaTeX render failed ({latex_str}): {_e}")
                return None
            finally:
                try:
                    if plt is not None:
                        if fig is not None:
                            plt.close(fig)
                        plt.close("all")
                except Exception:
                    pass

        if "formulas" in extras and formula_data:
            formulas = formula_data.get("formulas", [])
            example = formula_data.get("example", {})
            if formulas:
                await _add_bridge("before_formulas")
                if lang == "ru":
                    f_label = "Формулы:"
                elif lang == "en":
                    f_label = "Formulas:"
                else:
                    f_label = "Formulalar:"
                lbl_para = doc.add_paragraph()
                lbl_run = lbl_para.add_run(f_label)
                lbl_run.font.bold = True
                lbl_run.font.size = Pt(13)
                lbl_run.font.name = "Times New Roman"
                for f in formulas:
                    if f.get("name"):
                        f_name_para = doc.add_paragraph()
                        f_name_para.paragraph_format.first_line_indent = Inches(0.3)
                        f_name_run = f_name_para.add_run(f["name"] + ":")
                        f_name_run.font.bold = True
                        f_name_run.font.italic = False
                        f_name_run.font.size = Pt(13)
                        f_name_run.font.name = "Times New Roman"
                    # Render formula: prefer latex image, fallback to plain text
                    latex_src = f.get("latex") or f.get("formula", "")
                    if latex_src:
                        img_bytes = _render_latex(latex_src)
                        if img_bytes:
                            f_disp_para = doc.add_paragraph()
                            f_disp_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            f_disp_para.paragraph_format.space_before = Pt(2)
                            f_disp_para.paragraph_format.space_after = Pt(1)
                            f_disp_para.paragraph_format.line_spacing = 1.0
                            f_disp_run = f_disp_para.add_run()
                            f_disp_run.add_picture(_io.BytesIO(img_bytes), width=_Inches(3.5))
                        else:
                            f_disp_para = doc.add_paragraph()
                            f_disp_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            f_disp_para.paragraph_format.space_before = Pt(2)
                            f_disp_para.paragraph_format.space_after = Pt(2)
                            f_disp_para.paragraph_format.line_spacing = 1.0
                            f_disp_run = f_disp_para.add_run(f.get("formula", latex_src))
                            f_disp_run.font.size = Pt(16)
                            f_disp_run.font.bold = True
                            f_disp_run.font.name = "Courier New"
                    if f.get("explanation"):
                        exp_para = doc.add_paragraph()
                        exp_para.paragraph_format.first_line_indent = Inches(0.5)
                        exp_para.paragraph_format.space_before = Pt(0)
                        exp_para.paragraph_format.space_after = Pt(2)
                        exp_para.paragraph_format.line_spacing = 1.0
                        exp_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        exp_run = exp_para.add_run(f["explanation"])
                        exp_run.font.size = Pt(12)
                        exp_run.font.italic = True
                        exp_run.font.name = "Times New Roman"
                steps = example.get("steps") or (
                    # fallback: wrap old-style "solution" string as a single text step
                    [{"text": example["solution"]}] if example.get("solution") else []
                )
                if example and (example.get("task") or steps):
                    if lang == "ru":
                        ex_label = "Пример:"
                    elif lang == "en":
                        ex_label = "Example:"
                    else:
                        ex_label = "Misol:"
                    ex_lbl = doc.add_paragraph()
                    ex_lbl_run = ex_lbl.add_run(ex_label)
                    ex_lbl_run.font.bold = True
                    ex_lbl_run.font.size = Pt(13)
                    ex_lbl_run.font.name = "Times New Roman"
                    if example.get("task"):
                        task_para = doc.add_paragraph()
                        task_para.paragraph_format.first_line_indent = Inches(0.3)
                        task_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        task_run = task_para.add_run(example["task"])
                        task_run.font.size = Pt(13)
                        task_run.font.name = "Times New Roman"
                    # Render each step: latex → matplotlib image (left), text → plain paragraph
                    for step in steps:
                        if step.get("latex"):
                            step_img = _render_latex(step["latex"])
                            if step_img:
                                step_para = doc.add_paragraph()
                                step_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
                                step_para.paragraph_format.left_indent = Inches(0.3)
                                step_para.paragraph_format.space_before = Pt(1)
                                step_para.paragraph_format.space_after = Pt(0)
                                step_para.paragraph_format.line_spacing = 1.0
                                step_para.add_run().add_picture(
                                    _io.BytesIO(step_img), width=_Inches(3.2)
                                )
                            else:
                                step_para = doc.add_paragraph()
                                step_para.paragraph_format.left_indent = Inches(0.3)
                                step_para.paragraph_format.space_before = Pt(1)
                                step_para.paragraph_format.space_after = Pt(0)
                                step_para.paragraph_format.line_spacing = 1.0
                                step_run = step_para.add_run(step["latex"])
                                step_run.font.size = Pt(13)
                                step_run.font.name = "Courier New"
                        elif step.get("text"):
                            step_text = step["text"]
                            # If "text" field actually contains LaTeX, render it as image
                            _looks_latex = "\\" in step_text and any(
                                c in step_text for c in ("{", "_", "^", "\\frac", "\\sqrt")
                            )
                            if _looks_latex:
                                step_img = _render_latex(step_text)
                                if step_img:
                                    step_para = doc.add_paragraph()
                                    step_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
                                    step_para.paragraph_format.left_indent = Inches(0.3)
                                    step_para.paragraph_format.space_before = Pt(1)
                                    step_para.paragraph_format.space_after = Pt(0)
                                    step_para.paragraph_format.line_spacing = 1.0
                                    step_para.add_run().add_picture(
                                        _io.BytesIO(step_img), width=_Inches(3.2)
                                    )
                                else:
                                    step_para = doc.add_paragraph()
                                    step_para.paragraph_format.left_indent = Inches(0.3)
                                    step_para.paragraph_format.space_before = Pt(1)
                                    step_para.paragraph_format.space_after = Pt(0)
                                    step_para.paragraph_format.line_spacing = 1.0
                                    step_run = step_para.add_run(step_text)
                                    step_run.font.size = Pt(12)
                                    step_run.font.name = "Courier New"
                            else:
                                step_para = doc.add_paragraph()
                                step_para.paragraph_format.left_indent = Inches(0.3)
                                step_para.paragraph_format.space_before = Pt(1)
                                step_para.paragraph_format.space_after = Pt(0)
                                step_para.paragraph_format.line_spacing = 1.0
                                step_run = step_para.add_run(step_text)
                                step_run.font.size = Pt(12)
                                step_run.font.italic = True
                                step_run.font.name = "Times New Roman"

        # ── 3. Comparison table ───────────────────────────────────────────
        if "tables" in extras:
            tbl_data = await ai.generate_comparison_table(section_title, topic, lang)
            headers = tbl_data.get("headers", [])
            rows = tbl_data.get("rows", [])[:3]
            if headers and rows:
                await _add_bridge("before_table")
                from docx.shared import RGBColor as _RGB
                from docx.enum.table import WD_TABLE_ALIGNMENT
                tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                tbl.style = "Table Grid"

                # Set table to full usable page width (100%) and center it
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                tblPr = tbl._tbl.tblPr
                tblW = OxmlElement("w:tblW")
                tblW.set(qn("w:w"), "5000")   # 100% of usable page width
                tblW.set(qn("w:type"), "pct")
                tblPr.append(tblW)

                hdr_cells = tbl.rows[0].cells
                for i, h in enumerate(headers):
                    hdr_cells[i].text = str(h)
                    hdr_cells[i].paragraphs[0].runs[0].font.bold = True
                    hdr_cells[i].paragraphs[0].runs[0].font.size = Pt(10)
                    tc = hdr_cells[i]._tc
                    tcPr = tc.get_or_add_tcPr()
                    shd = OxmlElement("w:shd")
                    shd.set(qn("w:val"), "clear")
                    shd.set(qn("w:color"), "auto")
                    shd.set(qn("w:fill"), "D6E4F0")
                    tcPr.append(shd)
                for r_idx, row_data in enumerate(rows):
                    row_cells = tbl.rows[r_idx + 1].cells
                    for c_idx, cell_val in enumerate(row_data[:len(headers)]):
                        cell_text = str(cell_val)
                        # Hard cap at 120 chars (about 2 short sentences)
                        if len(cell_text) > 120:
                            cell_text = cell_text[:117] + "..."
                        row_cells[c_idx].text = cell_text
                        row_cells[c_idx].paragraphs[0].runs[0].font.size = Pt(9)
                doc.add_paragraph()

        # ── 4. Statistics ─────────────────────────────────────────────────
        if "statistics" in extras:
            stats_text = await ai.generate_section_statistics(section_title, topic, lang)
            # Strip any AI preamble lines that don't start with a bullet
            if stats_text:
                lines = stats_text.strip().splitlines()
                bullet_lines = [l for l in lines if l.strip().startswith("•")]
                if bullet_lines:
                    stats_text = "\n".join(bullet_lines)
                elif lines:
                    # No bullets found — keep as is but drop intro-style lines
                    stats_text = "\n".join(
                        l for l in lines
                        if not any(l.lower().startswith(w) for w in ("here", "below", "the following", "sure"))
                    ).strip()
            if stats_text and stats_text.strip():
                await _add_bridge("before_statistics")
                if lang == "ru":
                    stats_label = "Statistika i faktlar:"
                elif lang == "en":
                    stats_label = "Statistics & facts:"
                else:
                    stats_label = "Statistika va faktlar:"
                lbl_para = doc.add_paragraph()
                lbl_run = lbl_para.add_run(stats_label)
                lbl_run.font.bold = True
                lbl_run.font.size = Pt(13)
                lbl_run.font.name = "Times New Roman"
                stats_para = doc.add_paragraph()
                stats_para.paragraph_format.first_line_indent = Inches(0.3)
                stats_para.paragraph_format.line_spacing = 1.5
                stats_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                s_run = stats_para.add_run(stats_text.strip())
                s_run.font.size = Pt(13)
                s_run.font.italic = True
                s_run.font.name = "Times New Roman"

    async def _add_glossary_section(self, doc, topic: str, lang: str) -> None:
        """Add a glossary page at the end of the document."""
        from services.ai_service import get_ai_service
        ai = get_ai_service()
        terms = await ai.generate_glossary(topic, lang)
        if not terms:
            return
        doc.add_page_break()
        if lang == "ru":
            gloss_label = "ГЛОССАРИЙ"
        elif lang == "en":
            gloss_label = "GLOSSARY"
        else:
            gloss_label = "LUG'AT"
        g_title = doc.add_paragraph()
        g_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        g_run = g_title.add_run(gloss_label)
        g_run.font.bold = True
        g_run.font.size = Pt(14)
        g_run.font.name = "Times New Roman"
        for item in terms:
            term = item.get("term", "")
            definition = item.get("definition", "")
            if term and definition:
                t_para = doc.add_paragraph()
                t_para.paragraph_format.first_line_indent = Inches(0.0)
                t_para.paragraph_format.line_spacing = 1.5
                t_bold = t_para.add_run(term + " — ")
                t_bold.font.bold = True
                t_bold.font.size = Pt(13)
                t_bold.font.name = "Times New Roman"
                t_def = t_para.add_run(definition)
                t_def.font.size = Pt(13)
                t_def.font.name = "Times New Roman"

    async def create_independent_work(self, topic: str, content: Dict, extras: list = None) -> str:
        """Create independent work document with professional footnotes (snoska)"""
        try:
            doc = Document()
            style = doc.styles['Normal']
            font = style.font
            font.name = 'Times New Roman'
            font.size = Pt(14)
            
            paragraph_format = style.paragraph_format
            paragraph_format.line_spacing = 1.5
            paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            
            for section in doc.sections:
                section.top_margin = Inches(0.79)
                section.bottom_margin = Inches(0.79)
                section.left_margin = Inches(1.18)
                section.right_margin = Inches(0.39)

            user_lang = content.get('language', 'uz')
            author_name = content.get('author_name', '')
            await self._create_independent_work_title_page(doc, topic, user_lang, author_name)

            doc.add_page_break()

            toc_texts = self._get_toc_texts(user_lang)
            
            toc_para = doc.add_paragraph()
            toc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            toc_run = toc_para.add_run(toc_texts.get('reja', 'REJA').upper())
            toc_run.font.size = Pt(14)
            toc_run.font.bold = True
            all_sections = content.get('sections', [])

            toc_item = doc.add_paragraph()
            toc_item.add_run(toc_texts['kirish'])

            numbered_count = 0
            for idx, section in enumerate(all_sections):
                if idx == 0 or idx == len(all_sections) - 1:
                    continue
                numbered_count += 1
                toc_item = doc.add_paragraph()
                toc_item.add_run(f"{numbered_count}. {section['title']}")

            toc_item = doc.add_paragraph()
            toc_item.add_run(toc_texts['xulosa'])

            if content.get('references'):
                toc_item = doc.add_paragraph()
                toc_item.add_run(toc_texts['adabiyotlar'])

            doc.add_page_break()

            for section in doc.sections:
                self._add_page_number(section)

            # Get references for footnotes (use all references for cycling)
            references = content.get('references', [])
            footnote_counter = 1

            numbered_section_count = 0
            for idx, section in enumerate(all_sections):
                title = section['title']
                
                if idx == len(all_sections) - 1:
                    doc.add_page_break()
                
                section_title = doc.add_paragraph()
                section_title.alignment = WD_ALIGN_PARAGRAPH.CENTER

                if idx == 0:
                    section_title_run = section_title.add_run(toc_texts['kirish'].upper())
                elif idx == len(all_sections) - 1:
                    section_title_run = section_title.add_run(toc_texts['xulosa'].upper())
                else:
                    numbered_section_count += 1
                    section_title_run = section_title.add_run(f"{numbered_section_count}. {title}")

                section_title_run.font.bold = True
                section_title_run.font.size = Pt(14)

                # Add content with footnote references
                section_content = section['content']
                content_para = doc.add_paragraph()
                content_para.paragraph_format.first_line_indent = Inches(0.5)
                content_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                content_para.paragraph_format.line_spacing = 1.5
                
                content_run = content_para.add_run(section_content)
                content_run.font.size = Pt(14)
                content_run.font.name = 'Times New Roman'
                
                is_main = idx > 0 and idx < len(all_sections) - 1

                # Add footnote to main body sections (not intro/conclusion)
                if is_main and references:
                    # Filter out __CATEGORY__ headers — they are bibliography section
                    # headings, not actual citations. Never use them as snoska text.
                    citable_refs = [r for r in references if not r.startswith("__CATEGORY__")]
                    if citable_refs:
                        ref_idx = (footnote_counter - 1) % len(citable_refs)
                        self._add_footnote(content_para, citable_refs[ref_idx], footnote_counter)
                        footnote_counter += 1
                
                # Add extras only to main body sections, not kirish/xulosa
                # Use cycling pattern: pos1→formulas, pos2→image+table, pos3→table
                if extras and is_main:
                    cycle_extras = _extras_for_cycle(extras, numbered_section_count)
                    if cycle_extras:
                        await self._add_section_extras(doc, title, topic, user_lang, cycle_extras, section_idx=idx)
                elif numbered_section_count == 2 and self.together:
                    # Fallback: old Together image only when no extras selected
                    try:
                        image_path, image_prompt = await self.together.generate_flux_pro_image(
                            topic, title, user_lang
                        )
                        if image_path and os.path.exists(image_path):
                            doc.add_page_break()
                            img_para = doc.add_paragraph()
                            img_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            img_run = img_para.add_run()
                            img_run.add_picture(image_path, width=Inches(5.0))
                            caption_para = doc.add_paragraph()
                            caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            caption_para.paragraph_format.space_after = Pt(6)
                            if user_lang == 'uz':
                                caption_text = f"1-rasm. {title}"
                            elif user_lang == 'ru':
                                caption_text = f"Рисунок 1. {title}"
                            else:
                                caption_text = f"Figure 1. {title}"
                            cap_run = caption_para.add_run(caption_text)
                            cap_run.font.size = Pt(12)
                            cap_run.font.italic = True
                            cap_run.font.name = 'Times New Roman'
                            image_description = await self.together.generate_image_description(
                                topic, title, user_lang, image_path
                            )
                            if image_description:
                                from docx.enum.text import WD_LINE_SPACING
                                desc_para = doc.add_paragraph()
                                desc_para.paragraph_format.first_line_indent = Inches(0.5)
                                desc_para.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
                                desc_para.paragraph_format.space_after = Pt(12)
                                desc_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                                desc_run = desc_para.add_run(clean_text(image_description))
                                desc_run.font.size = Pt(14)
                                desc_run.font.name = 'Times New Roman'
                            doc.add_page_break()
                            try:
                                os.remove(image_path)
                            except Exception:
                                pass
                    except Exception as img_error:
                        logger.warning(f"Could not add image for independent work: {img_error}")

                # Add informational table after section 3 (on separate page, only when no extras)
                if not extras and numbered_section_count == 3:
                    table_data = content.get('table_data_3', content.get('table_data_2', []))
                    if table_data:
                        self._add_info_table(doc, topic, table_data, user_lang)

            if content.get('references'):
                doc.add_page_break()
                ref_title = doc.add_paragraph()
                ref_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                ref_title_run = ref_title.add_run(toc_texts['adabiyotlar'].upper())
                ref_title_run.font.bold = True
                ref_title_run.font.size = Pt(14)

                for idx, ref in enumerate(references, 1):
                    ref_para = doc.add_paragraph()
                    ref_para.paragraph_format.first_line_indent = Inches(0.5)
                    ref_para.paragraph_format.line_spacing = 1.5
                    ref_para.add_run(f"{idx}. {ref}")

            if extras and "glossary" in extras:
                await self._add_glossary_section(doc, topic, user_lang)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"independent_work_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(doc.save, file_path)
            logger.info(f"Independent work saved: {file_path}")
            return file_path

        except Exception as e:
            logger.error(f"Error creating independent work: {e}")
            raise

    async def create_referat(self, topic: str, content: Dict, extras: list = None) -> str:
        """Create referat document"""
        try:
            doc = Document()
            style = doc.styles['Normal']
            font = style.font
            font.name = 'Times New Roman'
            font.size = Pt(14)

            paragraph_format = style.paragraph_format
            paragraph_format.line_spacing = 1.5

            for idx, section in enumerate(doc.sections):
                section.top_margin = Inches(0.79)
                section.bottom_margin = Inches(0.79)
                section.left_margin = Inches(1.18)
                section.right_margin = Inches(0.39)
                section.footer.is_linked_to_previous = False
                if idx == 0:
                    section.different_first_page_header_footer = True

            user_lang = content.get('language', 'uz')
            author_name = content.get('author_name', '')
            await self._create_referat_title_page(doc, topic, user_lang, author_name)

            doc.add_page_break()

            toc_texts = self._get_toc_texts(user_lang)
            
            toc_para = doc.add_paragraph()
            toc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            toc_run = toc_para.add_run(toc_texts.get('reja', 'REJA').upper())
            toc_run.font.size = Pt(14)
            toc_run.font.bold = True
            all_sections = content.get('sections', [])

            toc_item = doc.add_paragraph()
            toc_item.add_run(toc_texts['kirish'])

            numbered_count = 0
            for idx, section in enumerate(all_sections):
                if idx == 0 or idx == len(all_sections) - 1:
                    continue
                numbered_count += 1
                toc_item = doc.add_paragraph()
                toc_item.add_run(f"{numbered_count}. {section['title']}")

            toc_item = doc.add_paragraph()
            toc_item.add_run(toc_texts['xulosa'])

            if content.get('references'):
                toc_item = doc.add_paragraph()
                toc_item.add_run(toc_texts['adabiyotlar'])

            doc.add_page_break()

            for section in doc.sections:
                self._add_page_number(section)

            references = content.get('references', [])
            footnote_counter = 1
            numbered_section_count = 0
            for idx, section in enumerate(all_sections):
                title = section['title']
                is_main_section = idx > 0 and idx < len(all_sections) - 1

                if idx == len(all_sections) - 1:
                    doc.add_page_break()
                
                section_title = doc.add_paragraph()
                section_title.alignment = WD_ALIGN_PARAGRAPH.CENTER

                if idx == 0:
                    section_title_run = section_title.add_run(toc_texts['kirish'].upper())
                elif idx == len(all_sections) - 1:
                    section_title_run = section_title.add_run(toc_texts['xulosa'].upper())
                else:
                    numbered_section_count += 1
                    section_title_run = section_title.add_run(f"{numbered_section_count}. {title}")

                section_title_run.font.bold = True
                section_title_run.font.size = Pt(14)

                content_para = doc.add_paragraph(section['content'])
                content_para.paragraph_format.first_line_indent = Inches(0.5)
                content_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                content_para.paragraph_format.line_spacing = 1.5

                # Add footnote to main body sections only (not intro/conclusion)
                if is_main_section and references:
                    # Filter out __CATEGORY__ headers — they are bibliography section
                    # headings, not actual citations. Never use them as snoska text.
                    citable_refs = [r for r in references if not r.startswith("__CATEGORY__")]
                    if citable_refs:
                        ref_idx = (footnote_counter - 1) % len(citable_refs)
                        self._add_footnote(content_para, citable_refs[ref_idx], footnote_counter)
                        footnote_counter += 1

                # Add extras only to main body sections, not kirish/xulosa
                # Use cycling pattern: pos1→formulas, pos2→image+table, pos3→table
                if extras and is_main_section:
                    cycle_extras = _extras_for_cycle(extras, numbered_section_count)
                    if cycle_extras:
                        await self._add_section_extras(doc, title, topic, user_lang, cycle_extras, section_idx=idx)

            if content.get('references'):
                doc.add_page_break()
                ref_title = doc.add_paragraph()
                ref_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                ref_title_run = ref_title.add_run(toc_texts['adabiyotlar'].upper())
                ref_title_run.font.bold = True
                ref_title_run.font.size = Pt(14)

                references = content['references'][:5]
                references_reversed = list(reversed(references))
                for idx, ref in enumerate(references_reversed, 1):
                    ref_para = doc.add_paragraph()
                    ref_para.paragraph_format.first_line_indent = Inches(0.5)
                    ref_para.paragraph_format.line_spacing = 1.5
                    ref_para.add_run(f"{idx}. {ref}")

            if extras and "glossary" in extras:
                await self._add_glossary_section(doc, topic, user_lang)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"referat_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(doc.save, file_path)
            logger.info(f"Referat saved: {file_path}")
            return file_path

        except Exception as e:
            logger.error(f"Error creating referat: {e}")
            raise

    async def _create_referat_title_page(self, doc, topic: str, language: str = 'uz', author_name: str = ''):
        """Create referat title page"""
        try:
            texts = self._get_referat_template_texts(language)

            para1 = doc.add_paragraph()
            para1.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run1 = para1.add_run("_" * 50)
            run1.font.size = Pt(14)
            run1.font.name = 'Times New Roman'

            para2 = doc.add_paragraph()
            para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run2 = para2.add_run("_" * 20 + f" {texts['from_subject']}")
            run2.font.size = Pt(14)
            run2.font.name = 'Times New Roman'

            for _ in range(4):
                doc.add_paragraph()

            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(f"{texts['referat']}:")
            title_run.font.size = Pt(36)
            title_run.font.bold = True
            title_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_run = topic_para.add_run(f"{texts['topic']}: {topic}")
            topic_run.font.size = Pt(14)
            topic_run.font.name = 'Times New Roman'

            for _ in range(2):
                doc.add_paragraph()

            signatures_para = doc.add_paragraph()
            signatures_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            bajardi_run = signatures_para.add_run(f"{texts['prepared_by']}: ")
            bajardi_run.font.size = Pt(14)
            bajardi_run.font.name = 'Times New Roman'

            if author_name:
                author_run = signatures_para.add_run(f"{author_name}")
                author_run.font.size = Pt(14)
                author_run.font.name = 'Times New Roman'
                author_run.font.bold = True
            else:
                kurs_run = signatures_para.add_run(f"_____ {texts['course']}")
                kurs_run.font.size = Pt(14)
                kurs_run.font.name = 'Times New Roman'

            signatures_para.add_run("               ")

            qabul_run = signatures_para.add_run(f"{texts['accepted_by']}: ")
            qabul_run.font.size = Pt(14)
            qabul_run.font.name = 'Times New Roman'

            qabul_line_run = signatures_para.add_run("_" * 15)
            qabul_line_run.font.size = Pt(14)
            qabul_line_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            city_para = doc.add_paragraph()
            city_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            city_run = city_para.add_run(texts['city'])
            city_run.font.size = Pt(14)
            city_run.font.name = 'Times New Roman'

        except Exception as e:
            logger.error(f"Error creating referat title page: {e}")

    def _get_referat_template_texts(self, language: str) -> Dict[str, str]:
        """Get language-specific texts for referat template"""
        if language == 'ru':
            return {
                'from_subject': 'по предмету',
                'referat': 'РЕФЕРАТ',
                'topic': 'Тема',
                'prepared_by': 'Выполнил',
                'course': 'курс',
                'accepted_by': 'Принял',
                'city': 'Ташкент'
            }
        elif language == 'en':
            return {
                'from_subject': 'on the subject',
                'referat': 'REPORT',
                'topic': 'Topic',
                'prepared_by': 'Prepared by',
                'course': 'course',
                'accepted_by': 'Accepted by',
                'city': 'Tashkent'
            }
        else:
            return {
                'from_subject': 'fanidan',
                'referat': 'REFERAT',
                'topic': 'Mavzu',
                'prepared_by': 'Bajardi',
                'course': 'kurs',
                'accepted_by': 'Qabul qildi',
                'city': 'Toshkent'
            }

    async def _create_independent_work_title_page(self, doc, topic: str, language: str = 'uz', author_name: str = ''):
        """Create independent work title page"""
        try:
            texts = self._get_independent_work_template_texts(language)

            faculty_para = doc.add_paragraph()
            faculty_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            faculty_run = faculty_para.add_run("_" * 30 + f" {texts['faculty']}")
            faculty_run.font.size = Pt(14)
            faculty_run.font.name = 'Times New Roman'

            subject_para = doc.add_paragraph()
            subject_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            subject_run = subject_para.add_run("_" * 30 + f" {texts['from_subject']}")
            subject_run.font.size = Pt(14)
            subject_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(texts['independent_work'])
            title_run.font.size = Pt(32)
            title_run.font.bold = True
            title_run.font.name = 'Times New Roman'

            for _ in range(2):
                doc.add_paragraph()

            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            topic_run = topic_para.add_run(f"{texts['topic']}: {topic}")
            topic_run.font.size = Pt(14)
            topic_run.font.name = 'Times New Roman'

            for _ in range(4):
                doc.add_paragraph()

            signatures_para = doc.add_paragraph()
            signatures_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

            bajardi_run = signatures_para.add_run(f"{texts['prepared_by']}: ")
            bajardi_run.font.size = Pt(14)
            bajardi_run.font.name = 'Times New Roman'

            if author_name:
                author_run = signatures_para.add_run(f"{author_name}")
                author_run.font.size = Pt(14)
                author_run.font.name = 'Times New Roman'
                author_run.font.bold = True
            else:
                bajardi_line_run = signatures_para.add_run("_" * 18)
                bajardi_line_run.font.size = Pt(14)
                bajardi_line_run.font.name = 'Times New Roman'

            signatures_para.add_run("         ")

            qabul_run = signatures_para.add_run(f"{texts['accepted_by']}: ")
            qabul_run.font.size = Pt(14)
            qabul_run.font.name = 'Times New Roman'

            qabul_line_run = signatures_para.add_run("_" * 15)
            qabul_line_run.font.size = Pt(14)
            qabul_line_run.font.name = 'Times New Roman'

        except Exception as e:
            logger.error(f"Error creating independent work title page: {e}")

    def _get_independent_work_template_texts(self, language: str) -> Dict[str, str]:
        """Get language-specific texts for independent work template"""
        if language == 'ru':
            return {
                'faculty': 'факультета',
                'from_subject': 'по предмету',
                'independent_work': 'Самостоятельная работа',
                'topic': 'Тема',
                'prepared_by': 'Выполнил',
                'accepted_by': 'Принял'
            }
        elif language == 'en':
            return {
                'faculty': 'faculty',
                'from_subject': 'on the subject',
                'independent_work': 'Independent work',
                'topic': 'Topic',
                'prepared_by': 'Prepared by',
                'accepted_by': 'Accepted by'
            }
        else:
            return {
                'faculty': 'fakulteti',
                'from_subject': 'fanidan',
                'independent_work': 'Mustaqil ish',
                'topic': 'Mavzu',
                'prepared_by': 'Bajardi',
                'accepted_by': 'Qabul qildi'
            }

    def _get_toc_texts(self, language: str) -> dict:
        """Get language-specific texts for table of contents"""
        if language == 'ru':
            return {
                'reja': 'План',
                'kirish': 'Введение',
                'xulosa': 'Заключение',
                'adabiyotlar': 'Использованная литература'
            }
        elif language == 'en':
            return {
                'reja': 'Contents',
                'kirish': 'Introduction',
                'xulosa': 'Conclusion',
                'adabiyotlar': 'References'
            }
        else:
            return {
                'reja': 'Reja',
                'kirish': 'Kirish',
                'xulosa': 'Xulosa',
                'adabiyotlar': 'Foydalangan adabiyotlar'
            }

    def _add_page_number(self, section):
        """Add page number to footer"""
        try:
            footer = section.footer
            paragraph = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

            run = paragraph.add_run()

            fldChar1 = OxmlElement('w:fldChar')
            fldChar1.set(qn('w:fldCharType'), 'begin')
            run._r.append(fldChar1)

            instrText = OxmlElement('w:instrText')
            instrText.set(qn('xml:space'), 'preserve')
            instrText.text = "PAGE"
            run._r.append(instrText)

            fldChar2 = OxmlElement('w:fldChar')
            fldChar2.set(qn('w:fldCharType'), 'separate')
            run._r.append(fldChar2)

            fldChar3 = OxmlElement('w:fldChar')
            fldChar3.set(qn('w:fldCharType'), 'end')
            run._r.append(fldChar3)

            run.font.size = Pt(14)
            run.font.name = 'Times New Roman'
        except Exception as e:
            logger.error(f"Error adding page number: {e}")

    def _add_footnote(self, paragraph, footnote_text: str, footnote_num: int = 1):
        """Add a real Word page-bottom footnote (snoska).

        Always routes to `_add_word_footnote_xml` because:
        - bayoo-docx's monkey-patch is broken under python-docx >=1.0
          (`hasattr(p, 'add_footnote')` returns False → silent fallback to a
          bare Unicode superscript with no actual page-bottom note).
        - The XML method works on any python-docx version and produces
          reference marks + body text that Word/LibreOffice render natively.

        `footnote_num` is treated as a *desired display number*. The actual
        internal w:footnote @w:id is auto-allocated to keep IDs unique within
        the document (avoids id collisions if the same number is reused).
        """
        try:
            doc = self._doc_from_paragraph(paragraph)
            if doc is None:
                raise RuntimeError("paragraph has no parent document")
            self._add_word_footnote_xml(doc, paragraph, footnote_text or "", footnote_num)
        except Exception as e:
            logger.error(f"Error adding footnote: {e}")
            run = paragraph.add_run(self._get_superscript_number(footnote_num))
            run.font.size = Pt(8)
            run.font.name = 'Times New Roman'

    @staticmethod
    def _doc_from_paragraph(paragraph):
        """Resolve the parent Document from a Paragraph in a python-docx safe way."""
        try:
            part = paragraph.part
            doc = getattr(part, "document", None)
            if doc is not None:
                return doc
            from docx.parts.document import DocumentPart
            if isinstance(part, DocumentPart):
                return part.document
        except Exception:
            pass
        return None

    def _add_word_footnote_xml(self, doc, paragraph, footnote_text: str, footnote_id: int = 1):
        """Create a proper Word page-bottom footnote via direct XML manipulation.

        Works with standard python-docx (no bayoo-docx required).

        Behaviour:
        - Lazily creates `/word/footnotes.xml` with separator + continuation
          footnotes (ids -1 and 0) on first call.
        - Allocates a unique w:footnote/@w:id by scanning existing ids — the
          caller-supplied `footnote_id` is only a hint; we never collide.
        - Cleans control chars from `footnote_text` (raw chars break OOXML).
        - Auto-appends a trailing period if text is non-empty and unpunctuated.
        - Registers the part on the package via `relate_to` so save() emits it.

        Returns the actual footnote id used (useful for callers that want to
        keep their own visible counter in sync).
        """
        try:
            from lxml import etree
            from docx.oxml.ns import qn
            from docx.opc.part import XmlPart
            from docx.opc.packuri import PackURI

            W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
            XML_SPACE = 'http://www.w3.org/XML/1998/namespace'
            FOOTNOTES_RT = (
                'http://schemas.openxmlformats.org/officeDocument/2006/'
                'relationships/footnotes'
            )
            FOOTNOTES_CT = (
                'application/vnd.openxmlformats-officedocument.'
                'wordprocessingml.footnotes+xml'
            )

            def w(tag):
                return f'{{{W}}}{tag}'

            # ── Locate or create the footnotes part ──────────────────────
            footnotes_elm = None
            for rel in doc.part.rels.values():
                if rel.reltype == FOOTNOTES_RT:
                    footnotes_elm = rel.target_part._element
                    break

            if footnotes_elm is None:
                fn_root_xml = (
                    f'<w:footnotes xmlns:w="{W}">'
                    f'<w:footnote w:type="separator" w:id="-1">'
                    f'<w:p><w:pPr><w:spacing w:after="0" w:line="240" w:lineRule="auto"/></w:pPr>'
                    f'<w:r><w:separator/></w:r></w:p></w:footnote>'
                    f'<w:footnote w:type="continuationSeparator" w:id="0">'
                    f'<w:p><w:pPr><w:spacing w:after="0" w:line="240" w:lineRule="auto"/></w:pPr>'
                    f'<w:r><w:continuationSeparator/></w:r></w:p></w:footnote>'
                    f'</w:footnotes>'
                )
                footnotes_elm = etree.fromstring(fn_root_xml.encode('utf-8'))
                fn_part = XmlPart(
                    PackURI('/word/footnotes.xml'),
                    FOOTNOTES_CT,
                    footnotes_elm,
                    doc.part.package,
                )
                doc.part.relate_to(fn_part, FOOTNOTES_RT)

            # ── Allocate a unique footnote id (avoid -1, 0, and existing) ─
            existing_ids = set()
            for fn_el in footnotes_elm.findall(w('footnote')):
                try:
                    existing_ids.add(int(fn_el.get(w('id'))))
                except (TypeError, ValueError):
                    pass
            new_id = max([1, footnote_id] + [i + 1 for i in existing_ids if i >= 1])
            while new_id in existing_ids:
                new_id += 1

            # ── Sanitize footnote text ───────────────────────────────────
            safe_text = self._sanitize_xml_text(footnote_text)
            if safe_text and safe_text[-1] not in '.!?…':
                safe_text += '.'

            # ── Build w:footnote element ─────────────────────────────────
            fn = etree.SubElement(footnotes_elm, w('footnote'))
            fn.set(w('id'), str(new_id))

            fp = etree.SubElement(fn, w('p'))
            fpp = etree.SubElement(fp, w('pPr'))
            fps = etree.SubElement(fpp, w('pStyle'))
            fps.set(w('val'), 'FootnoteText')
            # Tight single-line spacing for footnote body
            sp = etree.SubElement(fpp, w('spacing'))
            sp.set(w('after'), '0')
            sp.set(w('line'), '240')
            sp.set(w('lineRule'), 'auto')
            jc = etree.SubElement(fpp, w('jc'))
            jc.set(w('val'), 'both')

            # Reference mark run (the small superscript number)
            fr = etree.SubElement(fp, w('r'))
            frp = etree.SubElement(fr, w('rPr'))
            frs = etree.SubElement(frp, w('rStyle'))
            frs.set(w('val'), 'FootnoteReference')
            etree.SubElement(fr, w('footnoteRef'))

            # Footnote body text run
            fr2 = etree.SubElement(fp, w('r'))
            fr2p = etree.SubElement(fr2, w('rPr'))
            fr2sz = etree.SubElement(fr2p, w('sz'))
            fr2sz.set(w('val'), '20')          # 10pt
            fr2szcs = etree.SubElement(fr2p, w('szCs'))
            fr2szcs.set(w('val'), '20')
            fr2fn = etree.SubElement(fr2p, w('rFonts'))
            fr2fn.set(w('ascii'), 'Times New Roman')
            fr2fn.set(w('hAnsi'), 'Times New Roman')
            fr2fn.set(w('cs'), 'Times New Roman')
            ft = etree.SubElement(fr2, w('t'))
            ft.set(f'{{{XML_SPACE}}}space', 'preserve')
            ft.text = f' {safe_text}' if safe_text else ' '

            # ── Insert w:footnoteReference into the paragraph ────────────
            ref_r = etree.SubElement(paragraph._element, w('r'))
            ref_rpr = etree.SubElement(ref_r, w('rPr'))
            ref_rs = etree.SubElement(ref_rpr, w('rStyle'))
            ref_rs.set(w('val'), 'FootnoteReference')
            # Belt-and-suspenders: explicit superscript + 10pt so even apps
            # that don't honor the FootnoteReference style render correctly.
            ref_va = etree.SubElement(ref_rpr, w('vertAlign'))
            ref_va.set(w('val'), 'superscript')
            ref_sz = etree.SubElement(ref_rpr, w('sz'))
            ref_sz.set(w('val'), '20')
            ref_szcs = etree.SubElement(ref_rpr, w('szCs'))
            ref_szcs.set(w('val'), '20')
            ref_fn_el = etree.SubElement(ref_r, w('footnoteReference'))
            ref_fn_el.set(w('id'), str(new_id))

            return new_id

        except Exception as e:
            logger.error(f"Error adding XML footnote: {e}")
            sup_run = paragraph.add_run(self._get_superscript_number(footnote_id))
            sup_run.font.size = Pt(10)
            sup_run.font.name = 'Times New Roman'
            return footnote_id

    @staticmethod
    def _sanitize_xml_text(s: str) -> str:
        """Strip XML-illegal control chars and collapse whitespace.
        OOXML rejects U+0000-U+0008, U+000B, U+000C, U+000E-U+001F.
        """
        if not s:
            return ""
        out = []
        for ch in str(s):
            o = ord(ch)
            if o == 0x09 or o == 0x0A or o == 0x0D:
                out.append(' ')
            elif o < 0x20:
                continue
            else:
                out.append(ch)
        cleaned = ''.join(out)
        # collapse runs of whitespace
        return ' '.join(cleaned.split()).strip()

    def _get_superscript_number(self, num: int) -> str:
        """Return small superscript number (just ¹ format, not large)"""
        # Use only small superscript ¹ symbol, not large Unicode variants
        superscript_map = {
            '0': '⁰', '1': '¹', '2': '²', '3': '³', '4': '⁴',
            '5': '⁵', '6': '⁶', '7': '⁷', '8': '⁸', '9': '⁹'
        }
        return ''.join(superscript_map.get(c, c) for c in str(num))
    
    def _add_inline_footnote(self, paragraph, footnote_num: int, reference_text: str):
        """Fallback: Add inline footnote marker (superscript number) to paragraph
        
        This creates a simpler inline citation that works if native footnotes fail.
        Format: text¹ where the number is Unicode superscript on top-right
        """
        try:
            # Add Unicode superscript footnote number (¹²³⁴⁵⁶⁷⁸⁹)
            superscript_num = self._get_superscript_number(footnote_num)
            run = paragraph.add_run(superscript_num)
            run.font.size = Pt(10)
            run.font.name = 'Times New Roman'
            
            return footnote_num
            
        except Exception as e:
            logger.error(f"Error adding inline footnote: {e}")
            return footnote_num

    def _add_info_table(self, doc, topic: str, table_data, language: str = 'uz', chapter_num: int = 1):
        """Add a 5-row bordered informational table related to the topic
        
        Creates a professional table with borders showing comparative or analytical data.
        
        Args:
            doc: Document object
            topic: Topic to generate table caption for
            table_data: Dict with 'headers' and 'rows' keys, or list of rows
            language: Language for table caption
            chapter_num: Chapter number for table numbering
        """
        try:
            # Handle both dict format (new) and list format (old)
            if isinstance(table_data, dict):
                headers = table_data.get('headers', ['Tahlil', '2020', '2023', "O'zgarish"])
                rows = table_data.get('rows', [])
            else:
                # Old format - list of rows
                headers = ['Xususiyat', "Ko'rsatkich 1", "Ko'rsatkich 2", 'Natija']
                rows = table_data
            
            if not rows:
                logger.warning("No table data provided")
                return

            # Enforce 4 columns, 6 data rows
            headers = headers[:4]
            rows = [r[:4] for r in rows[:6]]

            # Add page break before table (table on separate page)
            doc.add_page_break()
            
            # Create table with rows + 1 header
            table = doc.add_table(rows=len(rows) + 1, cols=len(headers))
            table.style = 'Table Grid'
            
            # Header row
            header_row = table.rows[0]
            for i, header_text in enumerate(headers):
                if i < len(header_row.cells):
                    cell = header_row.cells[i]
                    cell.text = str(header_text)
                    for paragraph in cell.paragraphs:
                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.size = Pt(11)
                            run.font.name = 'Times New Roman'
            
            # Data rows
            for row_idx, row_data in enumerate(rows):
                row = table.rows[row_idx + 1]
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < len(row.cells):
                        cell = row.cells[col_idx]
                        cell.text = str(cell_text) if cell_text else ""
                        for paragraph in cell.paragraphs:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(11)
                                run.font.name = 'Times New Roman'
            
            # Get description from table_data if available (AI-generated)
            if isinstance(table_data, dict):
                base_desc = table_data.get('description', f"{topic} tahlili")
                if language == 'uz':
                    description = f"{chapter_num}-jadval. {base_desc}"
                elif language == 'ru':
                    description = f"Таблица {chapter_num}. {base_desc}"
                else:
                    description = f"Table {chapter_num}. {base_desc}"
            else:
                # Fallback descriptions
                captions = {
                    'uz': f"{chapter_num}-jadval. {topic} tahlili",
                    'ru': f"Таблица {chapter_num}. Анализ: {topic}", 
                    'en': f"Table {chapter_num}. Analysis: {topic}"
                }
                description = captions.get(language, captions['uz'])
            
            # Add table caption/description below
            caption_para = doc.add_paragraph()
            caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            caption_para.paragraph_format.space_before = Pt(6)
            caption_run = caption_para.add_run(description)
            caption_run.font.italic = True
            caption_run.font.size = Pt(11)
            caption_run.font.name = 'Times New Roman'
            
            # Add spacing after table
            doc.add_paragraph()
            
        except Exception as e:
            logger.error(f"Error adding info table: {e}")

    async def create_presentation_with_template_background(
        self, 
        topic: str, 
        content: Dict, 
        author_name: str, 
        template_id: str, 
        template_service, 
        language: str, 
        references: List = None, 
        plan_items: List = None
    ) -> str:
        """Create presentation with template background and custom content"""
        try:
            slides_data = content.get('slides', [])
            
            if references:
                for slide in slides_data:
                    if slide.get('layout') == 'references':
                        slide['references'] = references
                        break
                else:
                    slides_data.append({
                        'title': 'Adabiyotlar' if language == 'uz' else ('Литература' if language == 'ru' else 'References'),
                        'content': '',
                        'layout': 'references',
                        'references': references
                    })
            
            if plan_items:
                for slide in slides_data:
                    if slide.get('layout') == 'plan':
                        slide['plan_items'] = plan_items
                        break
            
            new_content = {'slides': slides_data}
            return await self.create_presentation_with_smart_images(topic, new_content, author_name, language, template_service, template_id)
            
        except Exception as e:
            logger.error(f"Error creating presentation with template background: {e}")
            raise

    async def create_new_presentation_system(
        self, 
        topic: str, 
        content: Dict, 
        author_name: str, 
        language: str
    ) -> str:
        """Create presentation using new layout system"""
        return await self.create_presentation_with_smart_images(topic, content, author_name, language)

    def _to_roman(self, n: int) -> str:
        """Convert integer to Roman numeral"""
        roman_map = [(10, 'X'), (9, 'IX'), (5, 'V'), (4, 'IV'), (1, 'I')]
        result = ""
        for val, rom in roman_map:
            while n >= val:
                result += rom
                n -= val
        return result

    async def create_thesis(self, topic: str, content: Dict, author_name: str, university: str, language: str = 'uz', faculty: str = '', group: str = '') -> str:
        """Create thesis document (4 pages) with specific formatting"""
        try:
            doc = Document()
            
            # Set margins
            for section in doc.sections:
                section.top_margin = Inches(0.79)
                section.bottom_margin = Inches(0.79)
                section.left_margin = Inches(1.18)
                section.right_margin = Inches(0.59)
            
            # 1. Title Page (Mavzu, Author, University)
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_uz = content.get('topic_uz', topic)
            topic_ru = content.get('topic_ru', '')
            topic_en = content.get('topic_en', '')
            title_lines = [topic_uz.upper()]
            if topic_ru:
                title_lines.append(topic_ru.upper())
            if topic_en:
                title_lines.append(topic_en.upper())
            title_run = title_para.add_run('\n'.join(title_lines))
            title_run.font.size = Pt(18)
            title_run.font.bold = True
            title_run.font.name = 'Times New Roman'
            
            doc.add_paragraph() # Spacing
            
            info_para = doc.add_paragraph()
            info_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            student_label = {'uz': 'guruhi talabasi', 'ru': 'группы студент', 'en': 'group student'}
            label = student_label.get(language, student_label['uz'])
            info_lines = [university]
            if faculty:
                info_lines.append(faculty)
            if group:
                info_lines.append(f"{group} {label} {author_name}")
            else:
                info_lines.append(author_name)
            info_run = info_para.add_run('\n'.join(info_lines))
            info_run.font.size = Pt(14)
            info_run.font.name = 'Times New Roman'
            info_run.font.italic = True
            
            doc.add_paragraph() # Spacing
            
            # 2. Trilingual Annotation, Keywords, Introduction
            labels = {
                'uz': {'lit_review': 'Adabiyotlar tahlili:', 'anal': 'Asosiy qism:', 'ref': 'Adabiyotlar roʻyxati:'},
                'ru': {'lit_review': 'Обзор литературы:', 'anal': 'Основная часть:', 'ref': 'Список литературы:'},
                'en': {'lit_review': 'Literature Review:', 'anal': 'Main Part:', 'ref': 'References:'}
            }
            l = labels.get(language, labels['uz'])
            
            FONT_SIZE = Pt(14)
            FONT_NAME = 'Times New Roman'
            FIRST_LINE_INDENT = Inches(0.49)
            
            refs = content.get('references', [])
            
            def add_heading_para(doc, text):
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                run = p.add_run(text)
                run.font.bold = True
                run.font.size = FONT_SIZE
                run.font.name = FONT_NAME
                return p
            
            def add_body_para(doc, text='', justify=True):
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY if justify else WD_ALIGN_PARAGRAPH.LEFT
                p.paragraph_format.first_line_indent = FIRST_LINE_INDENT
                if text:
                    run = p.add_run(text)
                    run.font.size = FONT_SIZE
                    run.font.name = FONT_NAME
                return p
            
            kw_labels = {'uz': 'Kalit so\u02bbzlar:', 'ru': 'Ключевые слова:', 'en': 'Keywords:'}
            add_heading_para(doc, 'Annotatsiya:')
            for lang_code in ['uz', 'ru', 'en']:
                annotation_text = content.get(f'annotation_{lang_code}', content.get('annotation', ''))
                if annotation_text:
                    p = add_body_para(doc)
                    run = p.add_run(annotation_text)
                    run.font.size = FONT_SIZE
                    run.font.name = FONT_NAME
                
                keywords_list = content.get(f'keywords_{lang_code}', content.get('keywords', []))
                if isinstance(keywords_list, list):
                    keywords_text = ', '.join(keywords_list)
                else:
                    keywords_text = str(keywords_list)
                p = doc.add_paragraph()
                run = p.add_run(f"{kw_labels[lang_code]} ")
                run.font.bold = True
                run.font.size = FONT_SIZE
                run.font.name = FONT_NAME
                run = p.add_run(keywords_text)
                run.font.size = FONT_SIZE
                run.font.name = FONT_NAME
                p.paragraph_format.first_line_indent = FIRST_LINE_INDENT
            
            intro_heading = {'uz': 'Kirish:', 'ru': 'Введение:', 'en': 'Introduction:'}
            add_heading_para(doc, intro_heading.get(language, 'Kirish:'))
            intro_text = content.get('introduction', content.get(f'intro_{language}', ''))
            if intro_text:
                p = add_body_para(doc)
                self._add_text_with_footnotes(p, intro_text, refs)
            
            add_heading_para(doc, l['lit_review'])
            p = add_body_para(doc)
            self._add_text_with_footnotes(p, content.get('literature_review', ''), refs)
            
            add_heading_para(doc, l['anal'])
            
            main_intro = content.get('main_intro', '')
            if main_intro:
                p = add_body_para(doc)
                self._add_text_with_footnotes(p, main_intro, refs)
            
            analysis = content.get('analysis', '')
            if isinstance(analysis, list):
                combined = ' '.join(f"{p.get('title', '')}. {p.get('content', '')}" for p in analysis if isinstance(p, dict))
                analysis = combined
            if analysis:
                paragraphs = [par.strip() for par in analysis.split('\n') if par.strip()]
                if len(paragraphs) <= 1:
                    paragraphs = [analysis]
                for para_text in paragraphs:
                    p = add_body_para(doc)
                    self._add_text_with_footnotes(p, para_text, refs)
            
            table_data = content.get('table', {})
            logger.info(f"Thesis table_data keys: {table_data.keys() if isinstance(table_data, dict) else type(table_data)}, has rows: {bool(table_data.get('rows') if isinstance(table_data, dict) else False)}")
            if table_data and isinstance(table_data, dict) and table_data.get('rows'):
                doc.add_paragraph()
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
                num_cols = len(headers) if headers else 1
                table = doc.add_table(rows=len(rows) + 1, cols=num_cols)
                table.style = 'Table Grid'
                
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = str(header)
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.font.bold = True
                            run.font.size = Pt(11)
                            run.font.name = FONT_NAME
                
                for i, row in enumerate(rows):
                    for j, val in enumerate(row):
                        if j < num_cols:
                            cell = table.cell(i + 1, j)
                            cell.text = str(val)
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.font.size = Pt(11)
                                    run.font.name = FONT_NAME

                explanation = content.get('table_explanation', '')
                if explanation:
                    p = add_body_para(doc)
                    p.paragraph_format.space_before = Pt(12)
                    self._add_text_with_footnotes(p, explanation, refs)

            table2_data = content.get('table2', {})
            if table2_data and table2_data.get('rows'):
                doc.add_paragraph()
                headers2 = table2_data.get('headers', [])
                rows2 = table2_data.get('rows', [])
                num_cols2 = len(headers2) if headers2 else 1
                table2 = doc.add_table(rows=len(rows2) + 1, cols=num_cols2)
                table2.style = 'Table Grid'
                
                for i, header in enumerate(headers2):
                    cell = table2.cell(0, i)
                    cell.text = str(header)
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.font.bold = True
                            run.font.size = Pt(11)
                            run.font.name = FONT_NAME
                
                for i, row in enumerate(rows2):
                    for j, val in enumerate(row):
                        if j < num_cols2:
                            cell = table2.cell(i + 1, j)
                            cell.text = str(val)
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.font.size = Pt(11)
                                    run.font.name = FONT_NAME

                explanation2 = content.get('table2_explanation', '')
                if explanation2:
                    p = add_body_para(doc)
                    p.paragraph_format.space_before = Pt(12)
                    self._add_text_with_footnotes(p, explanation2, refs)

            conclusion = content.get('conclusion', '')
            if conclusion:
                label_xulosa = {'uz': 'Xulosa:', 'ru': 'Заключение:', 'en': 'Conclusion:'}.get(language, 'Xulosa:')
                add_heading_para(doc, label_xulosa)
                p = add_body_para(doc)
                self._add_text_with_footnotes(p, conclusion, refs)

            doc.add_page_break()
            add_heading_para(doc, l['ref'])
            
            for i, ref in enumerate(content.get('references', []), 1):
                p = doc.add_paragraph()
                run = p.add_run(f"{i}. {ref}")
                run.font.size = FONT_SIZE
                run.font.name = FONT_NAME
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"thesis_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(doc.save, file_path)
            return file_path
        except Exception as e:
            logger.error(f"Error creating thesis: {e}")
            raise

    def _add_text_with_footnotes(self, paragraph, text: str, references: List[str] = None):
        """Helper to add text with snoska [n] support to a paragraph.

        Detects markers [1], [2] etc in text and creates real Word page-bottom
        footnotes via `_add_footnote`. Skips markers whose number exceeds the
        available references (avoids "Academic reference N" placeholders).
        Cycles through references if AI emitted markers like [99] by mistake.
        """
        if not text:
            return

        parts = re.split(r'(\[\d+\])', text)
        clean_refs = [r for r in (references or []) if r and not str(r).startswith('__CATEGORY__')]
        for part in parts:
            if not part:
                continue
            if re.match(r'\[\d+\]', part):
                marker_num = int(part[1:-1])
                if not clean_refs:
                    # No references available — drop the marker quietly
                    # rather than emitting a fake "Academic reference N".
                    continue
                ref_idx = (marker_num - 1) % len(clean_refs)
                ref_text = clean_refs[ref_idx]
                self._add_footnote(paragraph, ref_text, marker_num)
            else:
                run = paragraph.add_run(part)
                run.font.size = Pt(14)
                run.font.name = 'Times New Roman'

    async def _generate_diploma_infographic(self, topic: str, subsection_title: str, language: str) -> tuple:
        """Generate infographic image for diploma work subsection 2.

        Returns (image_path, image_description) or (None, None) on failure.
        Nano Banana / fal.ai disabled.
        """
        return None, None
        try:
            import aiohttp
            import time
            from services.fal_service import generate_image_nano

            if self.together:
                prompt = await self.together.generate_infographic_prompt(topic, subsection_title, language)
            else:
                prompt = (
                    f"Professional scientific infographic about '{subsection_title}' "
                    f"within the topic '{topic}', diagrams, flowcharts, icons, statistics, "
                    "no text, colorful modern academic design, high quality, landscape format."
                )

            logger.info(f"Generating diploma infographic via Nano Banana 2 for: {subsection_title}")
            image_url = await generate_image_nano(prompt, aspect_ratio="16_9")
            logger.info(f"Nano Banana 2 image URL: {image_url[:80] if image_url else 'None'}")

            filename = f"diploma_img_{int(time.time())}_{abs(hash(subsection_title)) % 10000}.png"
            os.makedirs(self.temp_dir, exist_ok=True)
            filepath = os.path.join(self.temp_dir, filename)

            async with aiohttp.ClientSession() as session:
                async with session.get(image_url, timeout=aiohttp.ClientTimeout(total=60)) as resp:
                    if resp.status == 200:
                        with open(filepath, "wb") as f:
                            f.write(await resp.read())
                        logger.info(f"Diploma infographic saved: {filepath}")
                    else:
                        logger.warning(f"Failed to download infographic, HTTP {resp.status}")
                        return None, None

            description = ""
            try:
                if self.together:
                    description = await self.together.generate_image_description(
                        topic, subsection_title, language, filepath
                    )
                else:
                    from services.ai_service import AIService as _AI, clean_text
                    _ai = _AI()
                    if language == "uz":
                        dp = f'"{topic}" mavzusining "{subsection_title}" bo\'limiga tegishli ilmiy infografik rasmni 60-80 so\'z bilan tavsif yozing. Faqat oddiy matn.'
                    elif language == "ru":
                        dp = f'Опишите научную инфографику по разделу "{subsection_title}" темы "{topic}" в 60-80 словах. Только обычный текст.'
                    else:
                        dp = f'Describe the scientific infographic for "{subsection_title}" of "{topic}" in 60-80 words. Plain text only.'
                    dr = await _ai._make_request(messages=[{"role": "user", "content": dp}], max_tokens=300, temperature=0.6)
                    description = clean_text(dr.strip())
            except Exception as de:
                logger.warning(f"Could not generate infographic description: {de}")

            return filepath, description

        except Exception as e:
            logger.warning(f"Could not generate diploma infographic: {e}", exc_info=True)
            return None, None

    async def create_diploma_work(self, topic: str, content: Dict, author_name: str, language: str = 'uz', extras: list = None) -> str:
        """Create diploma work document.

        Same structure as course work but:
        - Title page reads DIPLOM ISHI / ДИПЛОМНАЯ РАБОТА / DIPLOMA WORK
        - Each chapter's subsection 2 STARTS with a Nano Banana 2 infographic image
          followed by its caption and description, then the subsection text.
        """
        try:
            doc = Document()

            for section in doc.sections:
                section.top_margin = Inches(0.79)
                section.bottom_margin = Inches(0.79)
                section.left_margin = Inches(1.18)
                section.right_margin = Inches(0.59)

            await self._create_diploma_work_title_page(doc, topic, language, author_name)
            doc.add_page_break()

            self._create_diploma_work_toc(doc, content, language)
            doc.add_page_break()

            texts = self._get_diploma_work_texts(language)
            footnote_num = 1

            # Introduction
            intro_para = doc.add_paragraph()
            intro_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            intro_run = intro_para.add_run(texts['introduction'])
            intro_run.font.size = Pt(14)
            intro_run.font.bold = True
            intro_run.font.name = 'Times New Roman'

            intro_text = content.get('introduction', '')
            sentences = intro_text.split('. ')
            if len(sentences) > 4:
                mid = len(sentences) // 2
                paragraphs = ['. '.join(sentences[:mid]) + '.', '. '.join(sentences[mid:])]
            else:
                paragraphs = [intro_text]

            for p_text in paragraphs:
                if not p_text.strip():
                    continue
                p = doc.add_paragraph()
                p.paragraph_format.first_line_indent = Inches(0.5)
                p.paragraph_format.line_spacing = 1.5
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                run = p.add_run(p_text.strip())
                run.font.size = Pt(14)
                run.font.name = 'Times New Roman'

            doc.add_page_break()

            # Intro points
            intro_points_data = content.get('intro_points', {})
            for i, point_label in enumerate(texts['intro_points']):
                p = doc.add_paragraph()
                p.paragraph_format.line_spacing = 1.5
                p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                run = p.add_run(point_label)
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.name = 'Times New Roman'

                point_key = f"point_{i+1}"
                point_content = intro_points_data.get(point_key, "")
                if point_content:
                    if point_label.endswith(':'):
                        tasks = point_content.split('\n') if '\n' in point_content else [point_content]
                        for task in tasks:
                            if not task.strip():
                                continue
                            tp = doc.add_paragraph()
                            tp.paragraph_format.left_indent = Inches(0.5)
                            tp.paragraph_format.line_spacing = 1.5
                            tp.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                            tp_run = tp.add_run(f"• {task.strip()}")
                            tp_run.font.size = Pt(14)
                            tp_run.font.name = 'Times New Roman'
                    else:
                        run_content = p.add_run(f" {point_content}")
                        run_content.font.size = Pt(14)
                        run_content.font.name = 'Times New Roman'

            doc.add_page_break()

            references = content.get('references', [])
            footnote_counter = 1
            sub_counter = 0

            # Chapters
            for i, chapter in enumerate(content.get('chapters', []), 1):
                await asyncio.sleep(0)  # yield to event loop between chapters
                roman_num = self._to_roman(i)
                chapter_para = doc.add_paragraph()
                chapter_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                chapter_para.paragraph_format.space_before = Pt(18)
                chapter_para.paragraph_format.space_after = Pt(6)
                clean_ch_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
                chapter_run = chapter_para.add_run(f"{roman_num} {texts['chapter']}. {clean_ch_title.upper()}")
                chapter_run.font.size = Pt(14)
                chapter_run.font.bold = True
                chapter_run.font.name = 'Times New Roman'

                for j, subsection in enumerate(chapter.get('subsections', []), 1):
                    sub_counter += 1
                    # Subsection title
                    sub_para = doc.add_paragraph()
                    sub_para.paragraph_format.space_before = Pt(12)
                    sub_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                    clean_title = subsection['title']
                    clean_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', clean_title)

                    sub_run = sub_para.add_run(f"{subsection['number']} {clean_title}")
                    sub_run.font.size = Pt(14)
                    sub_run.font.bold = True
                    sub_run.font.name = 'Times New Roman'

                    # For subsection 2: add infographic image BEFORE text
                    if j == 2:
                        image_path, image_description = await self._generate_diploma_infographic(
                            topic, clean_title, language
                        )
                        if image_path and os.path.exists(image_path):
                            # Image paragraph (centered)
                            img_para = doc.add_paragraph()
                            img_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            img_run = img_para.add_run()
                            img_run.add_picture(image_path, width=Inches(5.5))

                            # Caption
                            caption_para = doc.add_paragraph()
                            caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            caption_para.paragraph_format.space_after = Pt(6)

                            if language == 'uz':
                                caption_text = f"{i}-rasm. {clean_title}"
                            elif language == 'ru':
                                caption_text = f"Рисунок {i}. {clean_title}"
                            else:
                                caption_text = f"Figure {i}. {clean_title}"

                            cap_run = caption_para.add_run(caption_text)
                            cap_run.font.size = Pt(12)
                            cap_run.font.italic = True
                            cap_run.font.name = 'Times New Roman'

                            # Description after caption
                            if image_description:
                                from docx.enum.text import WD_LINE_SPACING
                                desc_para = doc.add_paragraph()
                                desc_para.paragraph_format.first_line_indent = Inches(0.5)
                                desc_para.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
                                desc_para.paragraph_format.space_after = Pt(12)
                                desc_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                                desc_run = desc_para.add_run(clean_text(image_description))
                                desc_run.font.size = Pt(14)
                                desc_run.font.name = 'Times New Roman'

                    # Subsection text content
                    sub_content = subsection.get('content', '')
                    content_paragraphs = _split_into_paragraphs(sub_content, target_count=2)
                    footnotes_in_sub = 0
                    last_content_para = None

                    for p_idx, p_text in enumerate(content_paragraphs):
                        if not p_text.strip():
                            continue
                        content_para = doc.add_paragraph()
                        content_para.paragraph_format.first_line_indent = Inches(0.5)
                        content_para.paragraph_format.line_spacing = 1.5
                        content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                        content_run = content_para.add_run(p_text.strip())
                        content_run.font.size = Pt(14)
                        content_run.font.name = 'Times New Roman'
                        last_content_para = content_para

                        # Add a footnote at the end of each paragraph (2-3 per subsection)
                        if references:
                            ref_idx = (footnote_counter - 1) % len(references)
                            ref_text = references[ref_idx]
                            self._add_footnote(content_para, ref_text, footnote_counter)
                            footnote_counter += 1
                            footnotes_in_sub += 1

                    # Ensure at least 2 footnotes per subsection
                    while footnotes_in_sub < 2 and references and last_content_para is not None:
                        ref_idx = (footnote_counter - 1) % len(references)
                        ref_text = references[ref_idx]
                        self._add_footnote(last_content_para, ref_text, footnote_counter)
                        footnote_counter += 1
                        footnotes_in_sub += 1

                    # Table after subsection 3 (only when no extras)
                    if j == 3 and not extras:
                        table_data = content.get(f'table_data_{i}', {})
                        if table_data:
                            self._add_info_table(doc, topic, table_data, language, chapter_num=i)

                    # Add extras per subsection using cycle pattern
                    if extras:
                        cycle_extras = _extras_for_cycle(extras, sub_counter)
                        if cycle_extras:
                            await self._add_section_extras(doc, clean_title, topic, language, cycle_extras, section_idx=sub_counter)

                doc.add_page_break()

            # Conclusion
            conclusion_para = doc.add_paragraph()
            conclusion_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            conclusion_run = conclusion_para.add_run(texts['conclusion'])
            conclusion_run.font.size = Pt(14)
            conclusion_run.font.bold = True
            conclusion_run.font.name = 'Times New Roman'

            conclusion_text = content.get('conclusion', '')
            sentences = conclusion_text.split('. ')
            if len(sentences) > 4:
                mid = len(sentences) // 2
                paragraphs = ['. '.join(sentences[:mid]) + '.', '. '.join(sentences[mid:])]
            else:
                paragraphs = [conclusion_text]

            for p_text in paragraphs:
                if not p_text.strip():
                    continue
                cp = doc.add_paragraph()
                cp.paragraph_format.first_line_indent = Inches(0.5)
                cp.paragraph_format.line_spacing = 1.5
                cp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                cr = cp.add_run(p_text.strip())
                cr.font.size = Pt(14)
                cr.font.name = 'Times New Roman'

            # References
            doc.add_page_break()
            refs_para = doc.add_paragraph()
            refs_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            refs_run = refs_para.add_run(texts['references'])
            refs_run.font.size = Pt(14)
            refs_run.font.bold = True
            refs_run.font.name = 'Times New Roman'

            for idx, ref in enumerate(references, 1):
                ref_para = doc.add_paragraph()
                ref_para.paragraph_format.first_line_indent = Inches(0.5)
                ref_para.paragraph_format.line_spacing = 1.5
                ref_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                ref_run = ref_para.add_run(f"{idx}. {ref}")
                ref_run.font.size = Pt(14)
                ref_run.font.name = 'Times New Roman'

            if extras and "glossary" in extras:
                await self._add_glossary_section(doc, topic, language)

            # Page numbers
            for section in doc.sections:
                section.different_first_page_header_footer = True
                self._add_page_number(section)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"diplom_ishi_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(doc.save, file_path)
            logger.info(f"Diploma work saved: {file_path}")
            return file_path

        except Exception as e:
            logger.error(f"Error creating diploma work: {e}")
            raise

    def _get_diploma_work_texts(self, language: str) -> Dict[str, str]:
        """Get language-specific texts for diploma work"""
        if language == 'ru':
            return {
                'diploma_work': 'ДИПЛОМНАЯ РАБОТА',
                'faculty': 'факультет',
                'topic': 'Тема',
                'prepared_by': 'Выполнил(а)',
                'accepted_by': 'Принял(а)',
                'city': 'Ташкент',
                'contents': 'СОДЕРЖАНИЕ',
                'introduction': 'ВВЕДЕНИЕ',
                'chapter': 'ГЛАВА',
                'conclusion': 'ЗАКЛЮЧЕНИЕ',
                'references': 'СПИСОК ИСПОЛЬЗОВАННОЙ ЛИТЕРАТУРЫ',
                'intro_points': [
                    '1. Предмет дипломной работы.',
                    '2. Объект дипломной работы.',
                    '3. Степень изученности темы.',
                    '4. Цель дипломной работы.',
                    '5. Задачи дипломной работы:',
                    '6. Структура дипломной работы.'
                ]
            }
        elif language == 'en':
            return {
                'diploma_work': 'DIPLOMA WORK',
                'faculty': 'faculty',
                'topic': 'Topic',
                'prepared_by': 'Prepared by',
                'accepted_by': 'Accepted by',
                'city': 'Tashkent',
                'contents': 'CONTENTS',
                'introduction': 'INTRODUCTION',
                'chapter': 'CHAPTER',
                'conclusion': 'CONCLUSION',
                'references': 'REFERENCES',
                'intro_points': [
                    '1. Subject of the diploma work.',
                    '2. Object of the diploma work.',
                    '3. Degree of study of the topic.',
                    '4. Goal of the diploma work.',
                    '5. Tasks of the diploma work:',
                    '6. Structure of the diploma work.'
                ]
            }
        else:  # uz
            return {
                'diploma_work': 'DIPLOM ISHI',
                'faculty': 'fakulteti',
                'topic': 'Mavzu',
                'prepared_by': 'Bajardi',
                'accepted_by': 'Qabul qildi',
                'city': 'Toshkent',
                'contents': 'MUNDARIJA',
                'introduction': 'KIRISH',
                'chapter': 'BO\'LIM',
                'conclusion': 'XULOSA',
                'references': 'FOYDALANILGAN ADABIYOTLAR',
                'intro_points': [
                    '1. Diplom ishining predmeti.',
                    '2. Diplom ishining obyekti.',
                    '3. Mavzuning o\'rganilganlik darajasi.',
                    '4. Diplom ishining maqsadi.',
                    '5. Diplom ishining vazifalari:',
                    '6. Diplom ishining tarkibiy tuzilishi.'
                ]
            }

    async def _create_diploma_work_title_page(self, doc, topic: str, language: str = 'uz', author_name: str = ''):
        """Create diploma work title page"""
        try:
            texts = self._get_diploma_work_texts(language)

            uni_para = doc.add_paragraph()
            uni_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            uni_run = uni_para.add_run("_" * 50)
            uni_run.font.size = Pt(14)
            uni_run.font.name = 'Times New Roman'

            faculty_para = doc.add_paragraph()
            faculty_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            faculty_run = faculty_para.add_run("_" * 30 + f" {texts['faculty']}")
            faculty_run.font.size = Pt(14)
            faculty_run.font.name = 'Times New Roman'

            for _ in range(4):
                doc.add_paragraph()

            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(texts['diploma_work'])
            title_run.font.size = Pt(32)
            title_run.font.bold = True
            title_run.font.name = 'Times New Roman'

            for _ in range(2):
                doc.add_paragraph()

            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_run = topic_para.add_run(f"{texts['topic']}: {topic}")
            topic_run.font.size = Pt(14)
            topic_run.font.name = 'Times New Roman'

            for _ in range(4):
                doc.add_paragraph()

            author_para = doc.add_paragraph()
            author_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            bajardi_run = author_para.add_run(f"{texts['prepared_by']}: ")
            bajardi_run.font.size = Pt(14)
            bajardi_run.font.name = 'Times New Roman'

            if author_name:
                author_run = author_para.add_run(f"{author_name}")
                author_run.font.size = Pt(14)
                author_run.font.name = 'Times New Roman'
                author_run.font.bold = True
            else:
                line_run = author_para.add_run("_" * 20)
                line_run.font.size = Pt(14)
                line_run.font.name = 'Times New Roman'

            author_para.add_run("         ")
            qabul_run = author_para.add_run(f"{texts['accepted_by']}: ")
            qabul_run.font.size = Pt(14)
            qabul_run.font.name = 'Times New Roman'
            qabul_line_run = author_para.add_run("_" * 15)
            qabul_line_run.font.size = Pt(14)
            qabul_line_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            city_para = doc.add_paragraph()
            city_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            current_year = datetime.now().year
            city_run = city_para.add_run(f"{texts['city']} - {current_year}")
            city_run.font.size = Pt(14)
            city_run.font.name = 'Times New Roman'

        except Exception as e:
            logger.error(f"Error creating diploma work title page: {e}")

    def _create_diploma_work_toc(self, doc, content: Dict, language: str):
        """Create table of contents for diploma work"""
        texts = self._get_diploma_work_texts(language)

        toc_para = doc.add_paragraph()
        toc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        toc_run = toc_para.add_run(texts['contents'])
        toc_run.font.size = Pt(14)
        toc_run.font.bold = True
        toc_run.font.name = 'Times New Roman'

        doc.add_paragraph()

        intro_toc = doc.add_paragraph()
        intro_toc.paragraph_format.line_spacing = 1.5
        intro_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        intro_run = intro_toc.add_run(texts['introduction'])
        intro_run.font.size = Pt(14)
        intro_run.font.name = 'Times New Roman'

        for i, chapter in enumerate(content.get('chapters', []), 1):
            roman_num = self._to_roman(i)
            chapter_toc = doc.add_paragraph()
            chapter_toc.paragraph_format.line_spacing = 1.5
            chapter_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            clean_ch_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
            chapter_run = chapter_toc.add_run(f"{roman_num} {texts['chapter']}. {clean_ch_title.upper()}")
            chapter_run.font.size = Pt(14)
            chapter_run.font.bold = True
            chapter_run.font.name = 'Times New Roman'

            for subsection in chapter.get('subsections', []):
                sub_toc = doc.add_paragraph()
                sub_toc.paragraph_format.left_indent = Inches(0.5)
                sub_toc.paragraph_format.line_spacing = 1.5
                sub_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                clean_sub_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', subsection['title'])
                sub_run = sub_toc.add_run(f"{subsection['number']} {clean_sub_title}")
                sub_run.font.size = Pt(14)
                sub_run.font.name = 'Times New Roman'

        conclusion_toc = doc.add_paragraph()
        conclusion_toc.paragraph_format.line_spacing = 1.5
        conclusion_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        conclusion_run = conclusion_toc.add_run(texts['conclusion'])
        conclusion_run.font.size = Pt(14)
        conclusion_run.font.name = 'Times New Roman'

        refs_toc = doc.add_paragraph()
        refs_toc.paragraph_format.line_spacing = 1.5
        refs_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        refs_run = refs_toc.add_run(texts['references'])
        refs_run.font.size = Pt(14)
        refs_run.font.name = 'Times New Roman'

    async def create_course_work(self, topic: str, content: Dict, author_name: str, language: str = 'uz', extras: list = None) -> str:
        """Create course work document with chapters, subsections and footnotes
        
        Structure:
        - Title page
        - Table of contents
        - Introduction (Kirish)
        - Chapters with subsections (Bo'limlar)
        - Conclusion (Xulosa)
        - References (Adabiyotlar)
        
        Each page has a footnote at the bottom
        """
        try:
            doc = Document()
            
            # Set document margins
            for section in doc.sections:
                section.top_margin = Inches(0.79)  # 2 cm
                section.bottom_margin = Inches(0.79)  # 2 cm
                section.left_margin = Inches(1.18)  # 3 cm
                section.right_margin = Inches(0.59)  # 1.5 cm
            
            # Create title page
            await self._create_course_work_title_page(doc, topic, language, author_name)
            doc.add_page_break()
            
            # Create table of contents
            self._create_course_work_toc(doc, content, language)
            doc.add_page_break()
            
            # Footnote counter
            footnote_num = 1
            
            # Introduction
            texts = self._get_course_work_texts(language)
            intro_para = doc.add_paragraph()
            intro_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            intro_run = intro_para.add_run(texts['introduction'])
            intro_run.font.size = Pt(14)
            intro_run.font.bold = True
            intro_run.font.name = 'Times New Roman'
            
            # Intro Part 1: General Info
            intro_text = content.get('introduction', '')
            sentences = intro_text.split('. ')
            if len(sentences) > 4:
                mid = len(sentences) // 2
                paragraphs = ['. '.join(sentences[:mid]) + '.', '. '.join(sentences[mid:])]
            else:
                paragraphs = [intro_text]

            for p_text in paragraphs:
                if not p_text.strip(): continue
                intro_content_para = doc.add_paragraph()
                intro_content_para.paragraph_format.first_line_indent = Inches(0.5)
                intro_content_para.paragraph_format.line_spacing = 1.5
                intro_content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                intro_run = intro_content_para.add_run(p_text.strip())
                intro_run.font.size = Pt(14)
                intro_run.font.name = 'Times New Roman'
            
            doc.add_page_break()
            
            # Intro Part 2: Specific Points
            intro_points_data = content.get('intro_points', {})
            for i, point_label in enumerate(texts['intro_points']):
                p = doc.add_paragraph()
                p.paragraph_format.line_spacing = 1.5
                p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                run = p.add_run(point_label)
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.name = 'Times New Roman'
                
                point_key = f"point_{i+1}"
                point_content = intro_points_data.get(point_key, "")
                if point_content:
                    if point_label.endswith(':'):
                        # List format for tasks
                        tasks = point_content.split('\n') if '\n' in point_content else [point_content]
                        for task in tasks:
                            if not task.strip(): continue
                            tp = doc.add_paragraph()
                            tp.paragraph_format.left_indent = Inches(0.5)
                            tp.paragraph_format.line_spacing = 1.5
                            tp.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                            tp_run = tp.add_run(f"• {task.strip()}")
                            tp_run.font.size = Pt(14)
                            tp_run.font.name = 'Times New Roman'
                    else:
                        run_content = p.add_run(f" {point_content}")
                        run_content.font.size = Pt(14)
                        run_content.font.name = 'Times New Roman'
            
            doc.add_page_break()
            
            # Get references for footnotes
            references = content.get('references', [])
            footnote_counter = 1
            sub_counter = 0

            # Chapters
            for i, chapter in enumerate(content.get('chapters', []), 1):
                await asyncio.sleep(0)  # yield to event loop between chapters
                # Chapter title - Roman numerals
                roman_num = self._to_roman(i)
                chapter_para = doc.add_paragraph()
                chapter_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                chapter_para.paragraph_format.space_before = Pt(18)
                chapter_para.paragraph_format.space_after = Pt(6)
                clean_ch_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
                chapter_run = chapter_para.add_run(f"{roman_num} {texts['chapter']}. {clean_ch_title.upper()}")
                chapter_run.font.size = Pt(14)
                chapter_run.font.bold = True
                chapter_run.font.name = 'Times New Roman'

                # Subsections
                for j, subsection in enumerate(chapter.get('subsections', []), 1):
                    sub_counter += 1
                    # Subsection title - Arabic numerals
                    sub_para = doc.add_paragraph()
                    sub_para.paragraph_format.space_before = Pt(12)
                    sub_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    
                    # Clean title from existing numbering to avoid "1.1 1.1 Title"
                    clean_title = subsection['title']
                    clean_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', clean_title)
                    
                    sub_run = sub_para.add_run(f"{subsection['number']} {clean_title}")
                    sub_run.font.size = Pt(14)
                    sub_run.font.bold = True
                    sub_run.font.name = 'Times New Roman'
                    
                    # Subsection content
                    sub_content = subsection.get('content', '')
                    paragraphs = _split_into_paragraphs(sub_content, target_count=2)
                    footnotes_in_sub = 0
                    last_content_para = None

                    for p_idx, p_text in enumerate(paragraphs):
                        if not p_text.strip(): continue
                        content_para = doc.add_paragraph()
                        content_para.paragraph_format.first_line_indent = Inches(0.5)
                        content_para.paragraph_format.line_spacing = 1.5
                        content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        
                        content_run = content_para.add_run(p_text.strip())
                        content_run.font.size = Pt(14)
                        content_run.font.name = 'Times New Roman'
                        last_content_para = content_para
                        
                        # Add footnote at end of each paragraph (2-3 per subsection)
                        if references:
                            ref_idx = (footnote_counter - 1) % len(references)
                            ref_text = references[ref_idx]
                            self._add_footnote(content_para, ref_text, footnote_counter)
                            footnote_counter += 1
                            footnotes_in_sub += 1

                    # Ensure at least 2 footnotes per subsection
                    while footnotes_in_sub < 2 and references and last_content_para is not None:
                        ref_idx = (footnote_counter - 1) % len(references)
                        ref_text = references[ref_idx]
                        self._add_footnote(last_content_para, ref_text, footnote_counter)
                        footnote_counter += 1
                        footnotes_in_sub += 1
                    
                    # Add extras per subsection using cycle pattern
                    if extras:
                        cycle_extras = _extras_for_cycle(extras, sub_counter)
                        if cycle_extras:
                            await self._add_section_extras(doc, clean_title, topic, language, cycle_extras, section_idx=sub_counter)
                    else:
                        # Old Flux Pro image after subsection 2 (only when no extras)
                        if j == 2 and self.together:
                            try:
                                image_path, image_prompt = await self.together.generate_flux_pro_image(
                                    topic, clean_title, language
                                )
                                if image_path and os.path.exists(image_path):
                                    doc.add_page_break()
                                    img_para = doc.add_paragraph()
                                    img_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                    img_run = img_para.add_run()
                                    img_run.add_picture(image_path, width=Inches(5.0))
                                    caption_para = doc.add_paragraph()
                                    caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                    caption_para.paragraph_format.space_after = Pt(6)
                                    if language == 'uz':
                                        caption_text = f"{i}-rasm. {clean_title}"
                                    elif language == 'ru':
                                        caption_text = f"Рисунок {i}. {clean_title}"
                                    else:
                                        caption_text = f"Figure {i}. {clean_title}"
                                    cap_run = caption_para.add_run(caption_text)
                                    cap_run.font.size = Pt(12)
                                    cap_run.font.italic = True
                                    cap_run.font.name = 'Times New Roman'
                                    image_description = await self.together.generate_image_description(
                                        topic, clean_title, language, image_path
                                    )
                                    if image_description:
                                        from docx.enum.text import WD_LINE_SPACING
                                        desc_para = doc.add_paragraph()
                                        desc_para.paragraph_format.first_line_indent = Inches(0.5)
                                        desc_para.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
                                        desc_para.paragraph_format.space_after = Pt(12)
                                        desc_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                                        desc_run = desc_para.add_run(clean_text(image_description))
                                        desc_run.font.size = Pt(14)
                                        desc_run.font.name = 'Times New Roman'
                                    doc.add_page_break()
                                    logger.info(f"Added Flux Pro image with description for chapter {i} subsection 2")
                                    try:
                                        os.remove(image_path)
                                    except Exception:
                                        pass
                            except Exception as img_error:
                                logger.warning(f"Could not add image for chapter {i}.2: {img_error}")

                        # Add informational table after each chapter's subsection 3 (only without extras)
                        if j == 3:
                            table_data = content.get(f'table_data_{i}', {})
                            if table_data:
                                self._add_info_table(doc, topic, table_data, language, chapter_num=i)

                doc.add_page_break()
            
            # Conclusion
            conclusion_para = doc.add_paragraph()
            conclusion_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            conclusion_run = conclusion_para.add_run(texts['conclusion'])
            conclusion_run.font.size = Pt(14)
            conclusion_run.font.bold = True
            conclusion_run.font.name = 'Times New Roman'
            
            conclusion_text = content.get('conclusion', '')
            sentences = conclusion_text.split('. ')
            if len(sentences) > 4:
                mid = len(sentences) // 2
                paragraphs = ['. '.join(sentences[:mid]) + '.', '. '.join(sentences[mid:])]
            else:
                paragraphs = [conclusion_text]

            for p_text in paragraphs:
                if not p_text.strip(): continue
                conclusion_content_para = doc.add_paragraph()
                conclusion_content_para.paragraph_format.first_line_indent = Inches(0.5)
                conclusion_content_para.paragraph_format.line_spacing = 1.5
                conclusion_content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                conclusion_run = conclusion_content_para.add_run(p_text.strip())
                conclusion_run.font.size = Pt(14)
                conclusion_run.font.name = 'Times New Roman'
            
            # References - New Page
            doc.add_page_break()
            refs_para = doc.add_paragraph()
            refs_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            refs_run = refs_para.add_run(texts['references'])
            refs_run.font.size = Pt(14)
            refs_run.font.bold = True
            refs_run.font.name = 'Times New Roman'
            
            references = content.get('references', [])
            for idx, ref in enumerate(references, 1):
                ref_para = doc.add_paragraph()
                ref_para.paragraph_format.first_line_indent = Inches(0.5)
                ref_para.paragraph_format.line_spacing = 1.5
                ref_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                ref_run = ref_para.add_run(f"{idx}. {ref}")
                ref_run.font.size = Pt(14)
                ref_run.font.name = 'Times New Roman'
            
            if extras and "glossary" in extras:
                await self._add_glossary_section(doc, topic, language)

            # Add page numbers
            for section in doc.sections:
                section.different_first_page_header_footer = True
                self._add_page_number(section)
            
            # Save document
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"kurs_ishi_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(doc.save, file_path)
            logger.info(f"Course work saved: {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"Error creating course work: {e}")
            raise

    async def _create_course_work_title_page(self, doc, topic: str, language: str = 'uz', author_name: str = ''):
        """Create course work title page"""
        try:
            texts = self._get_course_work_texts(language)
            
            # University placeholder
            uni_para = doc.add_paragraph()
            uni_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            uni_run = uni_para.add_run("_" * 50)
            uni_run.font.size = Pt(14)
            uni_run.font.name = 'Times New Roman'
            
            # Faculty placeholder
            faculty_para = doc.add_paragraph()
            faculty_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            faculty_run = faculty_para.add_run("_" * 30 + f" {texts['faculty']}")
            faculty_run.font.size = Pt(14)
            faculty_run.font.name = 'Times New Roman'
            
            for _ in range(4):
                doc.add_paragraph()
            
            # Title
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(texts['course_work'])
            title_run.font.size = Pt(32)
            title_run.font.bold = True
            title_run.font.name = 'Times New Roman'
            
            for _ in range(2):
                doc.add_paragraph()
            
            # Topic
            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_run = topic_para.add_run(f"{texts['topic']}: {topic}")
            topic_run.font.size = Pt(14)
            topic_run.font.name = 'Times New Roman'
            
            for _ in range(4):
                doc.add_paragraph()
            
            # Author
            author_para = doc.add_paragraph()
            author_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
            bajardi_run = author_para.add_run(f"{texts['prepared_by']}: ")
            bajardi_run.font.size = Pt(14)
            bajardi_run.font.name = 'Times New Roman'
            
            if author_name:
                author_run = author_para.add_run(f"{author_name}")
                author_run.font.size = Pt(14)
                author_run.font.name = 'Times New Roman'
                author_run.font.bold = True
            else:
                line_run = author_para.add_run("_" * 20)
                line_run.font.size = Pt(14)
                line_run.font.name = 'Times New Roman'
            
            author_para.add_run("         ")
            
            qabul_run = author_para.add_run(f"{texts['accepted_by']}: ")
            qabul_run.font.size = Pt(14)
            qabul_run.font.name = 'Times New Roman'
            
            qabul_line_run = author_para.add_run("_" * 15)
            qabul_line_run.font.size = Pt(14)
            qabul_line_run.font.name = 'Times New Roman'
            
            for _ in range(3):
                doc.add_paragraph()
            
            # Title Page City
            city_para = doc.add_paragraph()
            city_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Extract year from timestamp or current year
            current_year = datetime.now().year
            city_text = f"{texts['city']} - {current_year}"
            
            city_run = city_para.add_run(city_text)
            city_run.font.size = Pt(14)
            city_run.font.name = 'Times New Roman'
            
        except Exception as e:
            logger.error(f"Error creating course work title page: {e}")

    def _create_course_work_toc(self, doc, content: Dict, language: str):
        """Create table of contents for course work"""
        texts = self._get_course_work_texts(language)
        
        # TOC title
        toc_para = doc.add_paragraph()
        toc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        toc_run = toc_para.add_run(texts['contents'])
        toc_run.font.size = Pt(14)
        toc_run.font.bold = True
        toc_run.font.name = 'Times New Roman'
        
        doc.add_paragraph()
        
        # Introduction
        intro_toc = doc.add_paragraph()
        intro_toc.paragraph_format.line_spacing = 1.5
        intro_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        intro_run = intro_toc.add_run(texts['introduction'])
        intro_run.font.size = Pt(14)
        intro_run.font.name = 'Times New Roman'
        
        # Chapters
        for i, chapter in enumerate(content.get('chapters', []), 1):
            # Chapter entry
            roman_num = self._to_roman(i)
            chapter_toc = doc.add_paragraph()
            chapter_toc.paragraph_format.line_spacing = 1.5
            chapter_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            clean_ch_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
            chapter_run = chapter_toc.add_run(f"{roman_num} {texts['chapter']}. {clean_ch_title.upper()}")
            chapter_run.font.size = Pt(14)
            chapter_run.font.bold = True
            chapter_run.font.name = 'Times New Roman'
            
            # Subsection entries
            for subsection in chapter.get('subsections', []):
                sub_toc = doc.add_paragraph()
                sub_toc.paragraph_format.left_indent = Inches(0.5)
                sub_toc.paragraph_format.line_spacing = 1.5
                sub_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                
                # Clean title from existing numbering for TOC
                clean_sub_title = subsection['title']
                clean_sub_title = re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', clean_sub_title)
                
                sub_run = sub_toc.add_run(f"{subsection['number']} {clean_sub_title}")
                sub_run.font.size = Pt(14)
                sub_run.font.name = 'Times New Roman'
        
        # Conclusion
        conclusion_toc = doc.add_paragraph()
        conclusion_toc.paragraph_format.line_spacing = 1.5
        conclusion_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        conclusion_run = conclusion_toc.add_run(texts['conclusion'])
        conclusion_run.font.size = Pt(14)
        conclusion_run.font.name = 'Times New Roman'
        
        # References
        refs_toc = doc.add_paragraph()
        refs_toc.paragraph_format.line_spacing = 1.5
        refs_toc.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        refs_run = refs_toc.add_run(texts['references'])
        refs_run.font.size = Pt(14)
        refs_run.font.name = 'Times New Roman'

    def _get_course_work_texts(self, language: str) -> Dict[str, str]:
        """Get language-specific texts for course work"""
        if language == 'ru':
            return {
                'course_work': 'КУРСОВАЯ РАБОТА',
                'faculty': 'факультет',
                'topic': 'Тема',
                'prepared_by': 'Выполнил(а)',
                'accepted_by': 'Принял(а)',
                'city': 'Ташкент',
                'contents': 'СОДЕРЖАНИЕ',
                'introduction': 'ВВЕДЕНИЕ',
                'chapter': 'ГЛАВА',
                'conclusion': 'ЗАКЛЮЧЕНИЕ',
                'references': 'СПИСОК ИСПОЛЬЗОВАННОЙ ЛИТЕРАТУРЫ',
                'intro_points': [
                    '1. Предмет курсовой работы.',
                    '2. Объект курсовой работы.',
                    '3. Степень изученности темы.',
                    '4. Цель курсовой работы.',
                    '5. Задачи курсовой работы:',
                    '6. Структура курсовой работы.'
                ]
            }
        elif language == 'en':
            return {
                'course_work': 'COURSE WORK',
                'faculty': 'faculty',
                'topic': 'Topic',
                'prepared_by': 'Prepared by',
                'accepted_by': 'Accepted by',
                'city': 'Tashkent',
                'contents': 'CONTENTS',
                'introduction': 'INTRODUCTION',
                'chapter': 'CHAPTER',
                'conclusion': 'CONCLUSION',
                'references': 'REFERENCES',
                'intro_points': [
                    '1. Subject of the course work.',
                    '2. Object of the course work.',
                    '3. Degree of study of the topic.',
                    '4. Goal of the course work.',
                    '5. Tasks of the course work:',
                    '6. Structure of the course work.'
                ]
            }
        else:  # uz
            return {
                'course_work': 'KURS ISHI',
                'faculty': 'fakulteti',
                'topic': 'Mavzu',
                'prepared_by': 'Bajardi',
                'accepted_by': 'Qabul qildi',
                'city': 'Toshkent',
                'contents': 'MUNDARIJA',
                'introduction': 'KIRISH',
                'chapter': 'BO\'LIM',
                'conclusion': 'XULOSA',
                'references': 'FOYDALANILGAN ADABIYOTLAR',
                'intro_points': [
                    '1. Kurs ishining predmeti.',
                    '2. Kurs ishining obyekti.',
                    '3. Mavzuning o‘rganilganlik darajasi.',
                    '4. Kurs ishining maqsadi.',
                    '5. Kurs ishining vazifalari:',
                    '6. Kurs ishining tarkibiy tuzilishi.'
                ]
            }

    async def create_article(self, topic: str, content: dict, author_name: str, language: str = "uz") -> str:
        """Create an IMRAD-structured academic article as a DOCX file"""
        doc = Document()

        for section in doc.sections:
            section.top_margin = Cm(2)
            section.bottom_margin = Cm(2)
            section.left_margin = Cm(2.5)
            section.right_margin = Cm(1.5)

        def _set_run(run, bold=False, size=14, italic=False):
            run.font.name = 'Times New Roman'
            run.font.size = Pt(size)
            run.bold = bold
            run.italic = italic

        def _add_body(text: str, indent: bool = True):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            run = p.add_run(text)
            _set_run(run)
            p.paragraph_format.line_spacing = Pt(21)
            if indent:
                p.paragraph_format.first_line_indent = Cm(1.25)
            p.paragraph_format.space_after = Pt(0)
            return p

        def _section_header(text: str):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            run = p.add_run(text.upper())
            _set_run(run, bold=True, size=14)
            p.paragraph_format.space_before = Pt(12)
            p.paragraph_format.space_after = Pt(6)

        p_title = doc.add_paragraph()
        p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_t = p_title.add_run(content.get("title", topic))
        _set_run(run_t, bold=True, size=14)

        if author_name:
            p_author = doc.add_paragraph()
            p_author.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run_a = p_author.add_run(author_name)
            _set_run(run_a, italic=True)

        doc.add_paragraph()

        if language == "uz":
            labels = {
                "abstract": "ANNOTATSIYA",
                "keywords": "Kalit so'zlar",
                "introduction": "KIRISH",
                "literature_review": "ADABIYOTLAR SHARHI",
                "methodology": "METODOLOGIYA",
                "results": "NATIJALAR VA MUHOKAMA",
                "conclusion": "XULOSA",
                "recommendations": "AMALIY TAKLIFLAR",
                "references": "FOYDALANILGAN ADABIYOTLAR",
            }
        elif language == "ru":
            labels = {
                "abstract": "АННОТАЦИЯ",
                "keywords": "Ключевые слова",
                "introduction": "ВВЕДЕНИЕ",
                "literature_review": "ОБЗОР ЛИТЕРАТУРЫ",
                "methodology": "МЕТОДОЛОГИЯ",
                "results": "РЕЗУЛЬТАТЫ И ОБСУЖДЕНИЕ",
                "conclusion": "ЗАКЛЮЧЕНИЕ",
                "recommendations": "ПРАКТИЧЕСКИЕ РЕКОМЕНДАЦИИ",
                "references": "СПИСОК ЛИТЕРАТУРЫ",
            }
        else:
            labels = {
                "abstract": "ABSTRACT",
                "keywords": "Keywords",
                "introduction": "INTRODUCTION",
                "literature_review": "LITERATURE REVIEW",
                "methodology": "METHODOLOGY",
                "results": "RESULTS AND DISCUSSION",
                "conclusion": "CONCLUSION",
                "recommendations": "PRACTICAL RECOMMENDATIONS",
                "references": "REFERENCES",
            }

        _section_header(labels["abstract"])
        _add_body(content.get("abstract", ""), indent=False)

        kw_list = content.get("keywords", [])
        if kw_list:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            run_bold = p.add_run(labels["keywords"] + ": ")
            _set_run(run_bold, bold=True)
            run_kw = p.add_run(", ".join(kw_list))
            _set_run(run_kw, italic=True)
            p.paragraph_format.space_after = Pt(6)

        doc.add_paragraph()

        imrad_sections = [
            ("introduction", labels["introduction"]),
            ("literature_review", labels["literature_review"]),
            ("methodology", labels["methodology"]),
            ("results_and_discussion", labels["results"]),
        ]

        for key, label in imrad_sections:
            _section_header(label)
            text = content.get(key, "")
            if text:
                for para in text.split("\n"):
                    para = para.strip()
                    if para:
                        _add_body(para)
            if key == "results_and_discussion":
                table_data = content.get("table", {})
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                caption = table_data.get("caption", "")
                if headers and rows:
                    doc.add_paragraph()
                    tbl = doc.add_table(rows=1, cols=len(headers))
                    tbl.style = 'Table Grid'
                    hdr_cells = tbl.rows[0].cells
                    for i, h in enumerate(headers):
                        hdr_cells[i].text = h
                        for run in hdr_cells[i].paragraphs[0].runs:
                            _set_run(run, bold=True, size=12)
                    for row in rows:
                        row_cells = tbl.add_row().cells
                        for i, val in enumerate(row[:len(headers)]):
                            row_cells[i].text = str(val)
                            for run in row_cells[i].paragraphs[0].runs:
                                _set_run(run, size=12)
                    if caption:
                        p_cap = doc.add_paragraph()
                        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run_cap = p_cap.add_run(caption)
                        _set_run(run_cap, italic=True, size=12)

        _section_header(labels["conclusion"])
        for para in content.get("conclusion", "").split("\n"):
            para = para.strip()
            if para:
                _add_body(para)

        recommendations = content.get("recommendations", "")
        if recommendations:
            _section_header(labels["recommendations"])
            for para in recommendations.split("\n"):
                para = para.strip()
                if para:
                    _add_body(para)

        _section_header(labels["references"])
        for ref in content.get("references", []):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            run = p.add_run(ref)
            _set_run(run, size=12)
            p.paragraph_format.space_after = Pt(2)

        os.makedirs(TEMP_DIR, exist_ok=True)
        safe_topic = re.sub(r'[^\w\s-]', '', topic)[:40].strip()
        filename = f"maqola_{safe_topic}_{author_name[:15] if author_name else 'anon'}.docx"
        file_path = os.path.join(TEMP_DIR, filename)
        await asyncio.to_thread(doc.save, file_path)
        logger.info(f"Article DOCX saved: {file_path}")
        return file_path

    async def create_mahsus_ishlanma(self, topic: str, content: dict, extras: list = None) -> str:
        """Create mahsus ishlanma (special project) document.

        Each content page has:
        - Top half: grey placeholder box with instruction text for image
        - Bottom half: actual section text
        """
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        from docx.shared import RGBColor as DocxRGBColor

        language = content.get('language', 'uz')
        author_name = content.get('author_name', '')

        # ── Labels ────────────────────────────────────────────────────────────
        if language == 'uz':
            labels = {
                'doc_type':      "MAHSUS ISHLANMA",
                'maqsad':        "ISHDAN MAQSAD",
                'tushunchalar':  "MAVZU BO'YICHA QISQACHA TUSHUNCHALAR",
                'amaliy':        "AMALIY QISM",
                'xulosa':        "XULOSA",
                'adabiyotlar':   "FOYDALANILGAN ADABIYOTLAR",
                'toc':           "MUNDARIJA",
                'qadam':         "Qadam",
            }
            placeholder_generic  = "Bu yerga «{topic}» mavzusiga oid tushuntiruvchi rasm qo'yilsin"
            placeholder_step     = "Bu yerga {qadam_nomi} amaliy qadamining screenshot rasmi qo'yilsin"
        elif language == 'ru':
            labels = {
                'doc_type':      "СПЕЦИАЛЬНАЯ РАЗРАБОТКА",
                'maqsad':        "ЦЕЛЬ РАБОТЫ",
                'tushunchalar':  "КРАТКИЕ ПОНЯТИЯ ПО ТЕМЕ",
                'amaliy':        "ПРАКТИЧЕСКАЯ ЧАСТЬ",
                'xulosa':        "ЗАКЛЮЧЕНИЕ",
                'adabiyotlar':   "ИСПОЛЬЗОВАННАЯ ЛИТЕРАТУРА",
                'toc':           "СОДЕРЖАНИЕ",
                'qadam':         "Шаг",
            }
            placeholder_generic  = "Здесь разместите поясняющий рисунок по теме «{topic}»"
            placeholder_step     = "Здесь разместите скриншот практического шага: {qadam_nomi}"
        else:
            labels = {
                'doc_type':      "SPECIAL PROJECT",
                'maqsad':        "PURPOSE OF THE WORK",
                'tushunchalar':  "BRIEF CONCEPTS",
                'amaliy':        "PRACTICAL PART",
                'xulosa':        "CONCLUSION",
                'adabiyotlar':   "REFERENCES",
                'toc':           "TABLE OF CONTENTS",
                'qadam':         "Step",
            }
            placeholder_generic  = "Place an explanatory image related to «{topic}» here"
            placeholder_step     = "Place the screenshot of practical step: {qadam_nomi} here"

        # ── Helper: add grey placeholder box ──────────────────────────────────
        def _add_placeholder(doc, hint_text: str):
            """Add a grey bordered box (≈half-page height) with hint text."""
            tbl = doc.add_table(rows=1, cols=1)
            tbl.style = 'Table Grid'
            cell = tbl.cell(0, 0)

            # Row height ≈ 9 cm
            tr = tbl.rows[0]._tr
            trPr = tr.get_or_add_trPr()
            trHeight = OxmlElement('w:trHeight')
            trHeight.set(qn('w:val'), str(int(9 * 567)))   # 1 cm = 567 twips
            trHeight.set(qn('w:hRule'), 'exact')
            trPr.append(trHeight)

            # Grey fill
            tcPr = cell._tc.get_or_add_tcPr()
            shd = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear')
            shd.set(qn('w:color'), 'auto')
            shd.set(qn('w:fill'), 'D9D9D9')
            tcPr.append(shd)

            # Centered italic hint text
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_before = Pt(120)  # push text to middle
            run = p.add_run(hint_text)
            run.font.size = Pt(12)
            run.font.italic = True
            run.font.name = 'Times New Roman'
            run.font.color.rgb = DocxRGBColor(0x80, 0x80, 0x80)

            # Space after table
            after = doc.add_paragraph()
            after.paragraph_format.space_after = Pt(6)

        # ── Helper: add section heading ────────────────────────────────────────
        def _heading(doc, text: str):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(text)
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.name = 'Times New Roman'
            p.paragraph_format.space_before = Pt(6)
            p.paragraph_format.space_after = Pt(6)

        # ── Helper: add body text ──────────────────────────────────────────────
        def _body(doc, text: str):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            p.paragraph_format.first_line_indent = Inches(0.5)
            p.paragraph_format.line_spacing = 1.5
            run = p.add_run(text)
            run.font.size = Pt(14)
            run.font.name = 'Times New Roman'

        # ── Build document ─────────────────────────────────────────────────────
        doc = Document()

        # Page margins (like independent work)
        for section in doc.sections:
            section.top_margin    = Inches(0.79)
            section.bottom_margin = Inches(0.79)
            section.left_margin   = Inches(1.18)
            section.right_margin  = Inches(0.39)

        style = doc.styles['Normal']
        style.font.name = 'Times New Roman'
        style.font.size = Pt(14)
        style.paragraph_format.line_spacing = 1.5

        # ── Title page ────────────────────────────────────────────────────────
        title_p = doc.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_p.add_run(labels['doc_type'])
        title_run.font.bold = True
        title_run.font.size = Pt(16)
        title_run.font.name = 'Times New Roman'

        doc.add_paragraph()

        topic_p = doc.add_paragraph()
        topic_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        topic_run = topic_p.add_run(f"Mavzu: {topic}" if language == 'uz' else
                                     f"Тема: {topic}" if language == 'ru' else
                                     f"Topic: {topic}")
        topic_run.font.bold = True
        topic_run.font.size = Pt(14)
        topic_run.font.name = 'Times New Roman'

        doc.add_paragraph()

        if author_name:
            author_p = doc.add_paragraph()
            author_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            label_text = ("Bajardi:" if language == 'uz' else
                          "Выполнил(а):" if language == 'ru' else "Prepared by:")
            author_run = author_p.add_run(f"{label_text} {author_name}")
            author_run.font.size = Pt(14)
            author_run.font.name = 'Times New Roman'

        doc.add_page_break()

        # ── Table of contents ─────────────────────────────────────────────────
        toc_h = doc.add_paragraph()
        toc_h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        toc_run = toc_h.add_run(labels['toc'])
        toc_run.font.bold = True
        toc_run.font.size = Pt(14)
        toc_run.font.name = 'Times New Roman'

        toc_items = [
            labels['maqsad'],
            labels['tushunchalar'],
            f"{labels['amaliy']}",
        ]
        steps = content.get('amaliy_qadamlar', [])
        for i, step in enumerate(steps, 1):
            toc_items.append(f"  {labels['qadam']} {i}. {step.get('qadam_nomi', '')}")
        toc_items += [labels['xulosa'], labels['adabiyotlar']]

        for item in toc_items:
            p = doc.add_paragraph()
            run = p.add_run(item)
            run.font.size = Pt(13)
            run.font.name = 'Times New Roman'

        doc.add_page_break()

        # ── Page numbers ──────────────────────────────────────────────────────
        for section in doc.sections:
            self._add_page_number(section)

        # ── Maqsad ────────────────────────────────────────────────────────────
        _heading(doc, labels['maqsad'])
        _add_placeholder(doc, placeholder_generic.format(topic=topic))
        _body(doc, content.get('maqsad', ''))
        doc.add_page_break()

        # ── Tushunchalar ──────────────────────────────────────────────────────
        _heading(doc, labels['tushunchalar'])
        _add_placeholder(doc, placeholder_generic.format(topic=topic))
        _body(doc, content.get('tushunchalar', ''))
        doc.add_page_break()

        # ── Amaliy qadamlar ───────────────────────────────────────────────────
        _heading(doc, labels['amaliy'])

        for i, step in enumerate(steps, 1):
            qadam_nomi = step.get('qadam_nomi', f"{labels['qadam']} {i}")
            qadam_tavsifi = step.get('qadam_tavsifi', '')

            step_h = doc.add_paragraph()
            step_run = step_h.add_run(f"{labels['qadam']} {i}: {qadam_nomi}")
            step_run.font.bold = True
            step_run.font.size = Pt(14)
            step_run.font.name = 'Times New Roman'

            hint = placeholder_step.format(qadam_nomi=qadam_nomi)
            _add_placeholder(doc, hint)
            _body(doc, qadam_tavsifi)

            if extras:
                await self._add_section_extras(doc, qadam_nomi, topic, language, extras, section_idx=i)

            if i < len(steps):
                doc.add_page_break()

        doc.add_page_break()

        # ── Xulosa ────────────────────────────────────────────────────────────
        _heading(doc, labels['xulosa'])
        _add_placeholder(doc, placeholder_generic.format(topic=topic))
        _body(doc, content.get('xulosa', ''))
        doc.add_page_break()

        # ── Adabiyotlar ───────────────────────────────────────────────────────
        _heading(doc, labels['adabiyotlar'])
        for ref in content.get('adabiyotlar', []):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            run = p.add_run(str(ref))
            run.font.size = Pt(13)
            run.font.name = 'Times New Roman'
            p.paragraph_format.space_after = Pt(4)

        if extras and "glossary" in extras:
            await self._add_glossary_section(doc, topic, language)

        # ── Save ──────────────────────────────────────────────────────────────
        os.makedirs(self.documents_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"mahsus_ishlanma_{timestamp}.docx"
        file_path = os.path.join(self.documents_dir, filename)
        await asyncio.to_thread(doc.save, file_path)
        logger.info(f"Mahsus ishlanma saved: {file_path}")
        return file_path


    async def create_dissertation(self, topic: str, content: Dict, author_name: str, language: str = 'uz', extras: list = None) -> str:
        """Create master's dissertation (Magistrlik dissertatsiyasi) with sections:
        1.Title page, 2.Bilingual annotation, 3.TOC, 4.Introduction (10 points),
        5.Main body (>=3 chapters), 6.Conclusion, 7.References, 8.Glossary, 9.Appendices.
        Margins per spec: top/bottom 2cm, left 3cm, right 2cm; line spacing 1.5; Times New Roman.
        """
        try:
            import re as _re
            doc = Document()
            for section in doc.sections:
                section.top_margin = Inches(0.79)      # ~2cm
                section.bottom_margin = Inches(0.79)   # ~2cm
                section.left_margin = Inches(1.18)     # ~3cm
                section.right_margin = Inches(0.79)    # ~2cm

            texts = self._get_dissertation_texts(language)

            # ── 1. TITLE PAGE ──────────────────────────────────────────────
            ministry_para = doc.add_paragraph()
            ministry_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            min_run = ministry_para.add_run(texts['ministry'])
            min_run.font.size = Pt(12)
            min_run.font.name = 'Times New Roman'

            uni_para = doc.add_paragraph()
            uni_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            uni_run = uni_para.add_run("_" * 50)
            uni_run.font.size = Pt(14)
            uni_run.font.bold = True
            uni_run.font.name = 'Times New Roman'

            faculty_para = doc.add_paragraph()
            faculty_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            fac_run = faculty_para.add_run("_" * 30 + f" {texts['faculty']}")
            fac_run.font.size = Pt(14)
            fac_run.font.name = 'Times New Roman'

            for _ in range(4):
                doc.add_paragraph()

            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(texts['dissertation'])
            title_run.font.size = Pt(28)
            title_run.font.bold = True
            title_run.font.name = 'Times New Roman'

            for _ in range(2):
                doc.add_paragraph()

            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_run = topic_para.add_run(f"{texts['topic']}: {topic}")
            topic_run.font.size = Pt(14)
            topic_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            author_para = doc.add_paragraph()
            author_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            author_para.add_run(f"{texts['prepared_by']}: ").font.size = Pt(14)
            if author_name:
                ar = author_para.add_run(author_name)
                ar.font.size = Pt(14)
                ar.font.bold = True
                ar.font.name = 'Times New Roman'
            else:
                author_para.add_run("_" * 20).font.size = Pt(14)

            supervisor_para = doc.add_paragraph()
            supervisor_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            sup_run = supervisor_para.add_run(f"{texts['supervisor']}: " + "_" * 20)
            sup_run.font.size = Pt(14)
            sup_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            city_para = doc.add_paragraph()
            city_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            city_run = city_para.add_run(f"{texts['city']} — {datetime.now().year}")
            city_run.font.size = Pt(14)
            city_run.font.name = 'Times New Roman'

            doc.add_page_break()

            # ── 2. TOC ─────────────────────────────────────────────────────
            toc_heading = doc.add_paragraph()
            toc_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            th_run = toc_heading.add_run(texts['contents'])
            th_run.font.size = Pt(14)
            th_run.font.bold = True
            th_run.font.name = 'Times New Roman'
            doc.add_paragraph()

            tp = doc.add_paragraph()
            tp.paragraph_format.line_spacing = 1.5
            tp.paragraph_format.space_after = Pt(0)
            tr = tp.add_run(texts['introduction'])
            tr.font.size = Pt(14)
            tr.font.name = 'Times New Roman'

            for i, chapter in enumerate(content.get('chapters', []), 1):
                roman_num_toc = self._to_roman(i)
                ch_toc = doc.add_paragraph()
                ch_toc.paragraph_format.line_spacing = 1.5
                ch_toc.paragraph_format.space_after = Pt(0)
                clean_ch_title_toc = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
                ch_tr = ch_toc.add_run(f"{roman_num_toc} {texts['chapter']}. {clean_ch_title_toc.upper()}")
                ch_tr.font.size = Pt(14)
                ch_tr.font.bold = True
                ch_tr.font.name = 'Times New Roman'

                for subsection in chapter.get('subsections', []):
                    sp = doc.add_paragraph()
                    sp.paragraph_format.left_indent = Inches(0.5)
                    sp.paragraph_format.line_spacing = 1.5
                    sp.paragraph_format.space_after = Pt(0)
                    clean_t = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', subsection['title'])
                    sr = sp.add_run(f"{subsection['number']} {clean_t}")
                    sr.font.size = Pt(14)
                    sr.font.name = 'Times New Roman'

            for toc_item in [texts['conclusion'], texts['references'], texts['glossary'], texts['appendices']]:
                tp2 = doc.add_paragraph()
                tp2.paragraph_format.line_spacing = 1.5
                tp2.paragraph_format.space_after = Pt(0)
                tr2 = tp2.add_run(toc_item)
                tr2.font.size = Pt(14)
                tr2.font.name = 'Times New Roman'

            doc.add_page_break()

            # ── 4. INTRODUCTION (10 points) ────────────────────────────────
            intro_heading = doc.add_paragraph()
            intro_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            ih_run = intro_heading.add_run(texts['introduction'])
            ih_run.font.size = Pt(14)
            ih_run.font.bold = True
            ih_run.font.name = 'Times New Roman'

            intro_text = content.get('introduction', '')
            sentences = intro_text.split('. ')
            if len(sentences) > 4:
                mid = len(sentences) // 2
                intro_paras = ['. '.join(sentences[:mid]) + '.', '. '.join(sentences[mid:])]
            else:
                intro_paras = [intro_text]
            for p_text in intro_paras:
                if not p_text.strip(): continue
                ip = doc.add_paragraph()
                ip.paragraph_format.first_line_indent = Inches(0.5)
                ip.paragraph_format.line_spacing = 1.5
                ip.paragraph_format.space_after = Pt(0)
                ip.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                ir = ip.add_run(p_text.strip())
                ir.font.size = Pt(14)
                ir.font.name = 'Times New Roman'

            intro_points_data = content.get('intro_points', {})
            intro_point_labels = texts['intro_points']
            for idx, label in enumerate(intro_point_labels):
                pp = doc.add_paragraph()
                pp.paragraph_format.line_spacing = 1.5
                pp.paragraph_format.space_after = Pt(0)
                pp.paragraph_format.space_before = Pt(4)
                pp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                pl = pp.add_run(label)
                pl.font.bold = True
                pl.font.size = Pt(14)
                pl.font.name = 'Times New Roman'
                point_key = f"point_{idx + 1}"
                point_content = intro_points_data.get(point_key, "")
                if point_content:
                    if label.rstrip().endswith(':') and '\n' in point_content:
                        for task in point_content.split('\n'):
                            if task.strip():
                                tp2 = doc.add_paragraph()
                                tp2.paragraph_format.left_indent = Inches(0.5)
                                tp2.paragraph_format.line_spacing = 1.5
                                tp2.paragraph_format.space_after = Pt(0)
                                tp2.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                                tr2 = tp2.add_run(f"• {task.strip()}")
                                tr2.font.size = Pt(14)
                                tr2.font.name = 'Times New Roman'
                    else:
                        pp.add_run(f" {point_content}").font.size = Pt(14)

            doc.add_page_break()

            # ── 5. MAIN BODY ───────────────────────────────────────────────
            references = content.get('references', [])
            clean_refs = [r for r in references if not r.startswith('__CATEGORY__')]
            footnote_counter = 1
            sub_counter = 0

            for i, chapter in enumerate(content.get('chapters', []), 1):
                await asyncio.sleep(0)  # yield to event loop between chapters
                roman_num = self._to_roman(i)
                ch_para = doc.add_paragraph()
                ch_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                ch_para.paragraph_format.space_before = Pt(18)
                ch_para.paragraph_format.space_after = Pt(6)
                clean_ch_title_body = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
                ch_run = ch_para.add_run(f"{roman_num} {texts['chapter']}. {clean_ch_title_body.upper()}")
                ch_run.font.size = Pt(14)
                ch_run.font.bold = True
                ch_run.font.name = 'Times New Roman'

                for j, subsection in enumerate(chapter.get('subsections', []), 1):
                    sub_counter += 1
                    sub_para = doc.add_paragraph()
                    sub_para.paragraph_format.space_before = Pt(12)
                    sub_para.paragraph_format.space_after = Pt(4)
                    sub_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    clean_title = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', subsection['title'])
                    sub_run = sub_para.add_run(f"{subsection['number']} {clean_title}")
                    sub_run.font.size = Pt(14)
                    sub_run.font.bold = True
                    sub_run.font.name = 'Times New Roman'

                    sub_content = subsection.get('content', '')
                    paragraphs = _split_into_paragraphs(sub_content, target_count=2)
                    footnotes_in_sub = 0
                    last_cp = None

                    for p_idx, p_text in enumerate(paragraphs):
                        if not p_text.strip(): continue
                        cp = doc.add_paragraph()
                        cp.paragraph_format.first_line_indent = Inches(0.5)
                        cp.paragraph_format.line_spacing = 1.5
                        cp.paragraph_format.space_after = Pt(0)
                        cp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        cr = cp.add_run(p_text.strip())
                        cr.font.size = Pt(14)
                        cr.font.name = 'Times New Roman'
                        last_cp = cp
                        if clean_refs:
                            ref_idx = (footnote_counter - 1) % len(clean_refs)
                            ref_text = clean_refs[ref_idx]
                            self._add_word_footnote_xml(doc, cp, ref_text, footnote_counter)
                            footnote_counter += 1
                            footnotes_in_sub += 1

                    # Ensure at least 2 footnotes per subsection
                    while footnotes_in_sub < 2 and clean_refs and last_cp is not None:
                        ref_idx = (footnote_counter - 1) % len(clean_refs)
                        ref_text = clean_refs[ref_idx]
                        self._add_word_footnote_xml(doc, last_cp, ref_text, footnote_counter)
                        footnote_counter += 1
                        footnotes_in_sub += 1

                    if j == 3:
                        table_data = content.get(f'table_data_{i}', {})
                        if table_data:
                            self._add_info_table(doc, topic, table_data, language, chapter_num=i)

                    # Add extras per subsection using cycle pattern
                    if extras:
                        cycle_extras = _extras_for_cycle(extras, sub_counter)
                        if cycle_extras:
                            await self._add_section_extras(doc, clean_title, topic, language, cycle_extras, section_idx=sub_counter)

                doc.add_page_break()

            # ── 6. CONCLUSION ──────────────────────────────────────────────
            conc_heading = doc.add_paragraph()
            conc_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            ch_run = conc_heading.add_run(texts['conclusion'])
            ch_run.font.size = Pt(14)
            ch_run.font.bold = True
            ch_run.font.name = 'Times New Roman'

            conclusion_text = content.get('conclusion', '')
            if conclusion_text:
                sents = conclusion_text.split('. ')
                if len(sents) > 4:
                    mid = len(sents) // 2
                    conc_paras = ['. '.join(sents[:mid]) + '.', '. '.join(sents[mid:])]
                else:
                    conc_paras = [conclusion_text]
                for pt in conc_paras:
                    if not pt.strip(): continue
                    cp2 = doc.add_paragraph()
                    cp2.paragraph_format.first_line_indent = Inches(0.5)
                    cp2.paragraph_format.line_spacing = 1.5
                    cp2.paragraph_format.space_after = Pt(0)
                    cp2.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    cr2 = cp2.add_run(pt.strip())
                    cr2.font.size = Pt(14)
                    cr2.font.name = 'Times New Roman'

            conclusion_points = content.get('conclusion_points', [])
            if conclusion_points:
                for pt in conclusion_points:
                    if not isinstance(pt, str): continue
                    if not pt.strip(): continue
                    bp = doc.add_paragraph()
                    bp.paragraph_format.left_indent = Inches(0.25)
                    bp.paragraph_format.line_spacing = 1.5
                    bp.paragraph_format.space_after = Pt(0)
                    bp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    br = bp.add_run(pt)
                    br.font.size = Pt(14)
                    br.font.name = 'Times New Roman'

            doc.add_page_break()

            # ── 7. REFERENCES ──────────────────────────────────────────────
            refs_heading = doc.add_paragraph()
            refs_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            rh_run = refs_heading.add_run(texts['references'])
            rh_run.font.size = Pt(14)
            rh_run.font.bold = True
            rh_run.font.name = 'Times New Roman'

            ref_num = 1
            for ref in references:
                if ref.startswith('__CATEGORY__'):
                    cat_name = ref[len('__CATEGORY__'):]
                    cat_p = doc.add_paragraph()
                    cat_p.paragraph_format.space_before = Pt(10)
                    cat_p.paragraph_format.space_after = Pt(2)
                    cat_run = cat_p.add_run(cat_name)
                    cat_run.font.bold = True
                    cat_run.font.size = Pt(14)
                    cat_run.font.name = 'Times New Roman'
                else:
                    rp = doc.add_paragraph()
                    rp.paragraph_format.first_line_indent = Inches(0.5)
                    rp.paragraph_format.line_spacing = 1.5
                    rp.paragraph_format.space_after = Pt(0)
                    rp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    rr = rp.add_run(f"{ref_num}. {ref}")
                    rr.font.size = Pt(14)
                    rr.font.name = 'Times New Roman'
                    ref_num += 1

            # ── 8. GLOSSARY ────────────────────────────────────────────────
            doc.add_page_break()
            glos_heading = doc.add_paragraph()
            glos_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            gh_run = glos_heading.add_run(texts['glossary'])
            gh_run.font.size = Pt(14)
            gh_run.font.bold = True
            gh_run.font.name = 'Times New Roman'
            doc.add_paragraph()

            for term_item in content.get('glossary_terms', []):
                if isinstance(term_item, dict):
                    term_val = term_item.get('term', '')
                    def_val = term_item.get('definition', '')
                    gp = doc.add_paragraph()
                    gp.paragraph_format.line_spacing = 1.5
                    gp.paragraph_format.space_after = Pt(0)
                    gp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    gt = gp.add_run(f"{term_val} — ")
                    gt.font.bold = True
                    gt.font.size = Pt(14)
                    gt.font.name = 'Times New Roman'
                    gd = gp.add_run(def_val)
                    gd.font.size = Pt(14)
                    gd.font.name = 'Times New Roman'

            # ── 9. APPENDICES ──────────────────────────────────────────────
            doc.add_page_break()
            app_heading = doc.add_paragraph()
            app_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            ah_run = app_heading.add_run(texts['appendices'])
            ah_run.font.size = Pt(14)
            ah_run.font.bold = True
            ah_run.font.name = 'Times New Roman'
            doc.add_paragraph()

            for app_item in content.get('appendices', []):
                if isinstance(app_item, dict):
                    app_title = app_item.get('title', '')
                    app_desc = app_item.get('description', '')
                    atp = doc.add_paragraph()
                    atp.paragraph_format.space_before = Pt(12)
                    atp.paragraph_format.space_after = Pt(4)
                    atr = atp.add_run(app_title)
                    atr.font.bold = True
                    atr.font.size = Pt(14)
                    atr.font.name = 'Times New Roman'
                    if app_desc:
                        adp = doc.add_paragraph()
                        adp.paragraph_format.first_line_indent = Inches(0.5)
                        adp.paragraph_format.line_spacing = 1.5
                        adp.paragraph_format.space_after = Pt(0)
                        adp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        adr = adp.add_run(app_desc)
                        adr.font.size = Pt(14)
                        adr.font.name = 'Times New Roman'

            for section in doc.sections:
                section.different_first_page_header_footer = True
                self._add_page_number(section)

            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            safe_topic = "".join(c if c.isalnum() else "_" for c in topic[:30])
            filename = f"dissertatsiya_{safe_topic}_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(doc.save, file_path)
            logger.info(f"Dissertation saved: {file_path}")
            return file_path

        except Exception as e:
            logger.error(f"Error creating dissertation: {e}")
            raise

    def _get_dissertation_texts(self, language: str) -> Dict:
        """Language-specific texts for master's dissertation."""
        if language == 'ru':
            return {
                'dissertation': 'МАГИСТЕРСКАЯ ДИССЕРТАЦИЯ',
                'ministry': 'МИНИСТЕРСТВО ВЫСШЕГО ОБРАЗОВАНИЯ, НАУКИ И ИННОВАЦИЙ РЕСПУБЛИКИ УЗБЕКИСТАН',
                'faculty': 'факультет',
                'topic': 'Тема',
                'prepared_by': 'Выполнил(а)',
                'supervisor': 'Научный руководитель',
                'city': 'Ташкент',
                'annotation_native': 'АННОТАЦИЯ',
                'annotation_english': 'ANNOTATION (in English)',
                'contents': 'СОДЕРЖАНИЕ',
                'introduction': 'ВВЕДЕНИЕ',
                'chapter': 'ГЛАВА',
                'conclusion': 'ЗАКЛЮЧЕНИЕ',
                'references': 'СПИСОК ИСПОЛЬЗОВАННОЙ ЛИТЕРАТУРЫ',
                'glossary': 'ГЛОССАРИЙ (СЛОВАРЬ ТЕРМИНОВ)',
                'appendices': 'ПРИЛОЖЕНИЯ',
                'intro_points': [
                    '1. Обоснование темы и её актуальность:',
                    '2. Объект исследования:',
                    '3. Предмет исследования:',
                    '4. Цель и задачи исследования:',
                    '5. Научная новизна:',
                    '6. Основные вопросы и гипотезы исследования:',
                    '7. Обзор литературы по теме исследования:',
                    '8. Описание применённой методики:',
                    '9. Теоретическая и практическая значимость результатов:',
                    '10. Описание структуры работы:',
                ]
            }
        elif language == 'en':
            return {
                'dissertation': "MASTER'S DISSERTATION",
                'ministry': 'MINISTRY OF HIGHER EDUCATION, SCIENCE AND INNOVATIONS OF THE REPUBLIC OF UZBEKISTAN',
                'faculty': 'faculty',
                'topic': 'Topic',
                'prepared_by': 'Prepared by',
                'supervisor': 'Scientific supervisor',
                'city': 'Tashkent',
                'annotation_native': 'ANNOTATION',
                'annotation_english': 'ANNOTATION (in English)',
                'contents': 'CONTENTS',
                'introduction': 'INTRODUCTION',
                'chapter': 'CHAPTER',
                'conclusion': 'CONCLUSION',
                'references': 'REFERENCES',
                'glossary': 'GLOSSARY',
                'appendices': 'APPENDICES',
                'intro_points': [
                    '1. Justification and relevance of the topic:',
                    '2. Object of research:',
                    '3. Subject of research:',
                    '4. Goal and tasks of research:',
                    '5. Scientific novelty:',
                    '6. Main research questions and hypotheses:',
                    '7. Literature review on the research topic:',
                    '8. Description of methodology used:',
                    '9. Theoretical and practical significance of results:',
                    '10. Description of work structure:',
                ]
            }
        else:  # uz
            return {
                'dissertation': 'MAGISTRLIK DISSERTATSIYASI',
                'ministry': "O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM, FAN VA INNOVATSIYALAR VAZIRLIGI",
                'faculty': 'fakulteti',
                'topic': 'Mavzu',
                'prepared_by': 'Bajardi',
                'supervisor': 'Ilmiy rahbar',
                'city': 'Toshkent',
                'annotation_native': 'ANNOTATSIYA',
                'annotation_english': 'ANNOTATION (in English)',
                'contents': 'MUNDARIJA',
                'introduction': 'KIRISH',
                'chapter': 'BOB',
                'conclusion': 'XULOSA',
                'references': 'FOYDALANILGAN ADABIYOTLAR',
                'glossary': "GLOSSARIY (ATAMALAR LUG'ATI)",
                'appendices': 'ILOVALAR',
                'intro_points': [
                    "1. Mavzuning asoslanishi va dolzarbligi:",
                    "2. Tadqiqot obyekti:",
                    "3. Tadqiqot predmeti:",
                    "4. Tadqiqotning maqsadi va vazifalari:",
                    "5. Ilmiy yangiligi:",
                    "6. Tadqiqotning asosiy masalalari va farazlari:",
                    "7. Tadqiqot mavzusi bo'yicha adabiyotlar sharhi:",
                    "8. Qo'llanilgan metodikaning tavsifi:",
                    "9. Tadqiqot natijalarining nazariy va amaliy ahamiyati:",
                    "10. Ish tuzilmasining tavsifi:",
                ]
            }

    async def create_graduation_work(self, topic: str, content: Dict, author_name: str, language: str = 'uz', extras: list = None) -> str:
        """Create graduation qualifying work (bitiruv malakaviy ishi) with 8 sections:
        1.Title page, 2.TOC, 3.Introduction, 4.Main body (chapters), 5.Conclusion+recommendations,
        6.References (categorized), 7.Glossary, 8.Appendices
        """
        try:
            import re as _re
            doc = Document()
            for section in doc.sections:
                section.top_margin = Inches(0.79)
                section.bottom_margin = Inches(0.79)
                section.left_margin = Inches(1.18)
                section.right_margin = Inches(0.59)

            texts = self._get_graduation_work_texts(language)

            # ── 1. TITLE PAGE ──────────────────────────────────────────────
            ministry_para = doc.add_paragraph()
            ministry_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            min_run = ministry_para.add_run(texts['ministry'])
            min_run.font.size = Pt(12)
            min_run.font.name = 'Times New Roman'

            uni_para = doc.add_paragraph()
            uni_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            uni_run = uni_para.add_run("_" * 50)
            uni_run.font.size = Pt(14)
            uni_run.font.bold = True
            uni_run.font.name = 'Times New Roman'

            faculty_para = doc.add_paragraph()
            faculty_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            fac_run = faculty_para.add_run("_" * 30 + f" {texts['faculty']}")
            fac_run.font.size = Pt(14)
            fac_run.font.name = 'Times New Roman'

            for _ in range(4):
                doc.add_paragraph()

            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(texts['graduation_work'])
            title_run.font.size = Pt(28)
            title_run.font.bold = True
            title_run.font.name = 'Times New Roman'

            for _ in range(2):
                doc.add_paragraph()

            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_run = topic_para.add_run(f"{texts['topic']}: {topic}")
            topic_run.font.size = Pt(14)
            topic_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            author_para = doc.add_paragraph()
            author_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            author_para.add_run(f"{texts['prepared_by']}: ").font.size = Pt(14)
            if author_name:
                ar = author_para.add_run(author_name)
                ar.font.size = Pt(14)
                ar.font.bold = True
                ar.font.name = 'Times New Roman'
            else:
                author_para.add_run("_" * 20).font.size = Pt(14)

            supervisor_para = doc.add_paragraph()
            supervisor_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            sup_run = supervisor_para.add_run(f"{texts['supervisor']}: " + "_" * 20)
            sup_run.font.size = Pt(14)
            sup_run.font.name = 'Times New Roman'

            for _ in range(3):
                doc.add_paragraph()

            city_para = doc.add_paragraph()
            city_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            city_run = city_para.add_run(f"{texts['city']} — {datetime.now().year}")
            city_run.font.size = Pt(14)
            city_run.font.name = 'Times New Roman'

            doc.add_page_break()

            # ── 2. TOC ─────────────────────────────────────────────────────
            toc_heading = doc.add_paragraph()
            toc_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            th_run = toc_heading.add_run(texts['contents'])
            th_run.font.size = Pt(14)
            th_run.font.bold = True
            th_run.font.name = 'Times New Roman'
            doc.add_paragraph()

            for toc_item, bold in [(texts['introduction'], False)]:
                tp = doc.add_paragraph()
                tp.paragraph_format.line_spacing = 1.5
                tp.paragraph_format.space_after = Pt(0)
                tr = tp.add_run(toc_item)
                tr.font.size = Pt(14)
                tr.font.name = 'Times New Roman'

            for i, chapter in enumerate(content.get('chapters', []), 1):
                roman_num_toc = self._to_roman(i)
                ch_toc = doc.add_paragraph()
                ch_toc.paragraph_format.line_spacing = 1.5
                ch_toc.paragraph_format.space_after = Pt(0)
                clean_ch_title_toc = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
                ch_tr = ch_toc.add_run(f"{roman_num_toc} {texts['chapter']}. {clean_ch_title_toc.upper()}")
                ch_tr.font.size = Pt(14)
                ch_tr.font.bold = True
                ch_tr.font.name = 'Times New Roman'

                for subsection in chapter.get('subsections', []):
                    sp = doc.add_paragraph()
                    sp.paragraph_format.left_indent = Inches(0.5)
                    sp.paragraph_format.line_spacing = 1.5
                    sp.paragraph_format.space_after = Pt(0)
                    clean_t = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', subsection['title'])
                    sr = sp.add_run(f"{subsection['number']} {clean_t}")
                    sr.font.size = Pt(14)
                    sr.font.name = 'Times New Roman'

            for toc_item in [texts['conclusion_toc'], texts['references'], texts['glossary'], texts['appendices']]:
                tp2 = doc.add_paragraph()
                tp2.paragraph_format.line_spacing = 1.5
                tp2.paragraph_format.space_after = Pt(0)
                tr2 = tp2.add_run(toc_item)
                tr2.font.size = Pt(14)
                tr2.font.name = 'Times New Roman'

            doc.add_page_break()

            # ── 3. INTRODUCTION ────────────────────────────────────────────
            intro_heading = doc.add_paragraph()
            intro_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            ih_run = intro_heading.add_run(texts['introduction'])
            ih_run.font.size = Pt(14)
            ih_run.font.bold = True
            ih_run.font.name = 'Times New Roman'

            intro_text = content.get('introduction', '')
            sentences = intro_text.split('. ')
            if len(sentences) > 4:
                mid = len(sentences) // 2
                intro_paras = ['. '.join(sentences[:mid]) + '.', '. '.join(sentences[mid:])]
            else:
                intro_paras = [intro_text]
            for p_text in intro_paras:
                if not p_text.strip(): continue
                ip = doc.add_paragraph()
                ip.paragraph_format.first_line_indent = Inches(0.5)
                ip.paragraph_format.line_spacing = 1.5
                ip.paragraph_format.space_after = Pt(0)
                ip.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                ir = ip.add_run(p_text.strip())
                ir.font.size = Pt(14)
                ir.font.name = 'Times New Roman'

            intro_points_data = content.get('intro_points', {})
            intro_point_labels = texts['intro_points']
            for idx, label in enumerate(intro_point_labels):
                pp = doc.add_paragraph()
                pp.paragraph_format.line_spacing = 1.5
                pp.paragraph_format.space_after = Pt(0)
                pp.paragraph_format.space_before = Pt(4)
                pp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                pl = pp.add_run(label)
                pl.font.bold = True
                pl.font.size = Pt(14)
                pl.font.name = 'Times New Roman'
                point_key = f"point_{idx + 1}"
                point_content = intro_points_data.get(point_key, "")
                if point_content:
                    if label.rstrip().endswith(':') and '\n' in point_content:
                        for task in point_content.split('\n'):
                            if task.strip():
                                tp2 = doc.add_paragraph()
                                tp2.paragraph_format.left_indent = Inches(0.5)
                                tp2.paragraph_format.line_spacing = 1.5
                                tp2.paragraph_format.space_after = Pt(0)
                                tp2.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                                tr2 = tp2.add_run(f"• {task.strip()}")
                                tr2.font.size = Pt(14)
                                tr2.font.name = 'Times New Roman'
                    else:
                        pp.add_run(f" {point_content}").font.size = Pt(14)

            doc.add_page_break()

            # ── 4. MAIN BODY ───────────────────────────────────────────────
            references = content.get('references', [])
            clean_refs = [r for r in references if not r.startswith('__CATEGORY__')]
            footnote_counter = 1

            for i, chapter in enumerate(content.get('chapters', []), 1):
                await asyncio.sleep(0)  # yield to event loop between chapters
                roman_num = self._to_roman(i)
                ch_para = doc.add_paragraph()
                ch_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                ch_para.paragraph_format.space_before = Pt(18)
                ch_para.paragraph_format.space_after = Pt(6)
                clean_ch_title_body = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', chapter['title'])
                ch_run = ch_para.add_run(f"{roman_num} {texts['chapter']}. {clean_ch_title_body.upper()}")
                ch_run.font.size = Pt(14)
                ch_run.font.bold = True
                ch_run.font.name = 'Times New Roman'

                for j, subsection in enumerate(chapter.get('subsections', []), 1):
                    sub_para = doc.add_paragraph()
                    sub_para.paragraph_format.space_before = Pt(12)
                    sub_para.paragraph_format.space_after = Pt(4)
                    sub_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    clean_title = _re.sub(r'^(?:\d+\.\s+|\d+(?:\.\d+)+\s+)', '', subsection['title'])
                    sub_run = sub_para.add_run(f"{subsection['number']} {clean_title}")
                    sub_run.font.size = Pt(14)
                    sub_run.font.bold = True
                    sub_run.font.name = 'Times New Roman'

                    sub_content = subsection.get('content', '')
                    paragraphs = _split_into_paragraphs(sub_content, target_count=2)
                    footnotes_in_sub = 0
                    last_cp = None

                    for p_idx, p_text in enumerate(paragraphs):
                        if not p_text.strip(): continue
                        cp = doc.add_paragraph()
                        cp.paragraph_format.first_line_indent = Inches(0.5)
                        cp.paragraph_format.line_spacing = 1.5
                        cp.paragraph_format.space_after = Pt(0)
                        cp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        cr = cp.add_run(p_text.strip())
                        cr.font.size = Pt(14)
                        cr.font.name = 'Times New Roman'
                        last_cp = cp
                        if clean_refs:
                            ref_idx = (footnote_counter - 1) % len(clean_refs)
                            ref_text = clean_refs[ref_idx]
                            self._add_word_footnote_xml(doc, cp, ref_text, footnote_counter)
                            footnote_counter += 1
                            footnotes_in_sub += 1

                    # Ensure at least 2 footnotes per subsection
                    while footnotes_in_sub < 2 and clean_refs and last_cp is not None:
                        ref_idx = (footnote_counter - 1) % len(clean_refs)
                        ref_text = clean_refs[ref_idx]
                        self._add_word_footnote_xml(doc, last_cp, ref_text, footnote_counter)
                        footnote_counter += 1
                        footnotes_in_sub += 1

                    if j == 3:
                        table_data = content.get(f'table_data_{i}', {})
                        if table_data:
                            self._add_info_table(doc, topic, table_data, language, chapter_num=i)

                    if extras:
                        await self._add_section_extras(doc, clean_title, topic, language, extras, section_idx=i + j)

                doc.add_page_break()

            # ── 5. CONCLUSION & RECOMMENDATIONS ───────────────────────────
            conc_heading = doc.add_paragraph()
            conc_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            ch_run = conc_heading.add_run(texts['conclusion'])
            ch_run.font.size = Pt(14)
            ch_run.font.bold = True
            ch_run.font.name = 'Times New Roman'

            conclusion_text = content.get('conclusion', '')
            if conclusion_text:
                sents = conclusion_text.split('. ')
                if len(sents) > 4:
                    mid = len(sents) // 2
                    conc_paras = ['. '.join(sents[:mid]) + '.', '. '.join(sents[mid:])]
                else:
                    conc_paras = [conclusion_text]
                for pt in conc_paras:
                    if not pt.strip(): continue
                    cp2 = doc.add_paragraph()
                    cp2.paragraph_format.first_line_indent = Inches(0.5)
                    cp2.paragraph_format.line_spacing = 1.5
                    cp2.paragraph_format.space_after = Pt(0)
                    cp2.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    cr2 = cp2.add_run(pt.strip())
                    cr2.font.size = Pt(14)
                    cr2.font.name = 'Times New Roman'

            conclusion_points = content.get('conclusion_points', [])
            if conclusion_points:
                for pt in conclusion_points:
                    if not isinstance(pt, str): continue
                    if not pt.strip(): continue
                    bp = doc.add_paragraph()
                    bp.paragraph_format.left_indent = Inches(0.25)
                    bp.paragraph_format.line_spacing = 1.5
                    bp.paragraph_format.space_after = Pt(0)
                    bp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    br = bp.add_run(pt)
                    br.font.size = Pt(14)
                    br.font.name = 'Times New Roman'

            doc.add_page_break()

            # ── 6. REFERENCES (CATEGORIZED) ────────────────────────────────
            refs_heading = doc.add_paragraph()
            refs_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            rh_run = refs_heading.add_run(texts['references'])
            rh_run.font.size = Pt(14)
            rh_run.font.bold = True
            rh_run.font.name = 'Times New Roman'

            ref_num = 1
            for ref in references:
                if ref.startswith('__CATEGORY__'):
                    cat_name = ref[len('__CATEGORY__'):]
                    cat_p = doc.add_paragraph()
                    cat_p.paragraph_format.space_before = Pt(10)
                    cat_p.paragraph_format.space_after = Pt(2)
                    cat_run = cat_p.add_run(cat_name)
                    cat_run.font.bold = True
                    cat_run.font.size = Pt(14)
                    cat_run.font.name = 'Times New Roman'
                else:
                    rp = doc.add_paragraph()
                    rp.paragraph_format.first_line_indent = Inches(0.5)
                    rp.paragraph_format.line_spacing = 1.5
                    rp.paragraph_format.space_after = Pt(0)
                    rp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    rr = rp.add_run(f"{ref_num}. {ref}")
                    rr.font.size = Pt(14)
                    rr.font.name = 'Times New Roman'
                    ref_num += 1

            # ── 7. GLOSSARY ────────────────────────────────────────────────
            doc.add_page_break()
            glos_heading = doc.add_paragraph()
            glos_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            gh_run = glos_heading.add_run(texts['glossary'])
            gh_run.font.size = Pt(14)
            gh_run.font.bold = True
            gh_run.font.name = 'Times New Roman'
            doc.add_paragraph()

            glossary_terms = content.get('glossary_terms', [])
            for term_item in glossary_terms:
                if isinstance(term_item, dict):
                    term_val = term_item.get('term', '')
                    def_val = term_item.get('definition', '')
                    gp = doc.add_paragraph()
                    gp.paragraph_format.line_spacing = 1.5
                    gp.paragraph_format.space_after = Pt(0)
                    gp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    gt = gp.add_run(f"{term_val} — ")
                    gt.font.bold = True
                    gt.font.size = Pt(14)
                    gt.font.name = 'Times New Roman'
                    gd = gp.add_run(def_val)
                    gd.font.size = Pt(14)
                    gd.font.name = 'Times New Roman'

            # ── 8. APPENDICES ──────────────────────────────────────────────
            doc.add_page_break()
            app_heading = doc.add_paragraph()
            app_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            ah_run = app_heading.add_run(texts['appendices'])
            ah_run.font.size = Pt(14)
            ah_run.font.bold = True
            ah_run.font.name = 'Times New Roman'
            doc.add_paragraph()

            appendices = content.get('appendices', [])
            for app_item in appendices:
                if isinstance(app_item, dict):
                    app_title = app_item.get('title', '')
                    app_desc = app_item.get('description', '')
                    atp = doc.add_paragraph()
                    atp.paragraph_format.space_before = Pt(12)
                    atp.paragraph_format.space_after = Pt(4)
                    atr = atp.add_run(app_title)
                    atr.font.bold = True
                    atr.font.size = Pt(14)
                    atr.font.name = 'Times New Roman'
                    if app_desc:
                        adp = doc.add_paragraph()
                        adp.paragraph_format.first_line_indent = Inches(0.5)
                        adp.paragraph_format.line_spacing = 1.5
                        adp.paragraph_format.space_after = Pt(0)
                        adp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                        adr = adp.add_run(app_desc)
                        adr.font.size = Pt(14)
                        adr.font.name = 'Times New Roman'

            # Page numbers
            for section in doc.sections:
                section.different_first_page_header_footer = True
                self._add_page_number(section)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"bitiruv_ishi_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            await asyncio.to_thread(doc.save, file_path)
            logger.info(f"Graduation work saved: {file_path}")
            return file_path

        except Exception as e:
            logger.error(f"Error creating graduation work: {e}")
            raise

    def _get_graduation_work_texts(self, language: str) -> Dict[str, str]:
        """Get language-specific texts for graduation qualifying work."""
        if language == 'ru':
            return {
                'graduation_work': 'ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА',
                'ministry': 'МИНИСТЕРСТВО ВЫСШЕГО ОБРАЗОВАНИЯ, НАУКИ И ИННОВАЦИЙ РЕСПУБЛИКИ УЗБЕКИСТАН',
                'faculty': 'факультет',
                'topic': 'Тема',
                'prepared_by': 'Выполнил(а)',
                'supervisor': 'Научный руководитель',
                'city': 'Ташкент',
                'contents': 'СОДЕРЖАНИЕ',
                'introduction': 'ВВЕДЕНИЕ',
                'chapter': 'ГЛАВА',
                'conclusion': 'ЗАКЛЮЧЕНИЕ И РЕКОМЕНДАЦИИ',
                'conclusion_toc': 'ЗАКЛЮЧЕНИЕ',
                'references': 'СПИСОК ИСПОЛЬЗОВАННОЙ ЛИТЕРАТУРЫ',
                'glossary': 'ГЛОССАРИЙ (СЛОВАРЬ ТЕРМИНОВ)',
                'appendices': 'ПРИЛОЖЕНИЯ',
                'intro_points': [
                    '1. Актуальность исследования:',
                    '2. Объект исследования:',
                    '3. Предмет исследования:',
                    '4. Цель исследования:',
                    '5. Задачи исследования:',
                    '6. Методы исследования:',
                    '7. Научная новизна:',
                    '8. Структура работы:',
                ]
            }
        elif language == 'en':
            return {
                'graduation_work': 'GRADUATION QUALIFYING WORK',
                'ministry': 'MINISTRY OF HIGHER EDUCATION, SCIENCE AND INNOVATIONS OF THE REPUBLIC OF UZBEKISTAN',
                'faculty': 'faculty',
                'topic': 'Topic',
                'prepared_by': 'Prepared by',
                'supervisor': 'Scientific supervisor',
                'city': 'Tashkent',
                'contents': 'CONTENTS',
                'introduction': 'INTRODUCTION',
                'chapter': 'CHAPTER',
                'conclusion': 'CONCLUSION AND RECOMMENDATIONS',
                'conclusion_toc': 'CONCLUSION',
                'references': 'REFERENCES',
                'glossary': 'GLOSSARY',
                'appendices': 'APPENDICES',
                'intro_points': [
                    '1. Relevance of research:',
                    '2. Object of research:',
                    '3. Subject of research:',
                    '4. Goal of research:',
                    '5. Tasks of research:',
                    '6. Research methods:',
                    '7. Scientific novelty:',
                    '8. Structure of work:',
                ]
            }
        else:  # uz
            return {
                'graduation_work': 'BITIRUV MALAKAVIY ISHI',
                'ministry': "O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM, FAN VA INNOVATSIYALAR VAZIRLIGI",
                'faculty': 'fakulteti',
                'topic': 'Mavzu',
                'prepared_by': 'Bajardi',
                'supervisor': 'Ilmiy rahbar',
                'city': 'Toshkent',
                'contents': 'MUNDARIJA',
                'introduction': 'KIRISH',
                'chapter': "BO'LIM",
                'conclusion': 'XULOSA VA TAKLIFLAR',
                'conclusion_toc': 'XULOSA',
                'references': 'FOYDALANILGAN ADABIYOTLAR',
                'glossary': "GLOSSARIY (ATAMALAR LUG'ATI)",
                'appendices': 'ILOVALAR',
                'intro_points': [
                    '1. Tadqiqotning dolzarbligi:',
                    '2. Tadqiqot obyekti:',
                    '3. Tadqiqot predmeti:',
                    '4. Tadqiqotning maqsadi:',
                    '5. Tadqiqot vazifalari:',
                    '6. Tadqiqot metodlari:',
                    '7. Ishning ilmiy yangiligi:',
                    '8. Ishning tarkibiy tuzilishi:',
                ]
            }


_document_service_instance: "DocumentService | None" = None


def get_document_service() -> "DocumentService":
    """Return the shared DocumentService singleton."""
    global _document_service_instance
    if _document_service_instance is None:
        _document_service_instance = DocumentService()
        logger.info("DocumentService singleton created")
    return _document_service_instance
