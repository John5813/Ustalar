"""
Kanvas renderer — AI tomonidan tasvirlangan elementlarni python-pptx orqali chizadi.
"""
import logging
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

log = logging.getLogger("layouts")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

CHART_TYPE_MAP = {
    "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":   XL_CHART_TYPE.LINE,
    "pie":    XL_CHART_TYPE.PIE,
    "donut":  XL_CHART_TYPE.DOUGHNUT,
}


def _hex(hex_str: str) -> RGBColor:
    h = (hex_str or "000000").lstrip("#").strip()
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _clamp(val, lo, hi):
    return max(lo, min(hi, val))


def render_canvas(slide, s, image_paths: dict):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex(s.canvas.background)

    for el in s.canvas.elements:
        try:
            if el.type == "rect":
                _draw_rect(slide, el)
            elif el.type == "text":
                _draw_text(slide, el)
            elif el.type == "circle":
                _draw_circle(slide, el)
            elif el.type == "image":
                img_path = image_paths.get(id(el))
                if img_path:
                    _draw_image(slide, el, img_path)
            elif el.type == "chart":
                _draw_chart(slide, el)
        except Exception as exc:
            log.warning("Element chizishda xato (%s, slayd %s): %s", el.type, s.index, exc)


# ─────────────────────────────────────────────────────────── rect

def _draw_rect(slide, el):
    w = _clamp(el.w or 1.0, 0.05, 13.333)
    h = _clamp(el.h or 1.0, 0.05, 7.5)
    x = _clamp(el.x, -0.5, 13.333)
    y = _clamp(el.y, -0.5, 7.5)

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if el.radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if el.fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex(el.fill)
    else:
        shp.fill.background()
    shp.line.fill.background()


# ─────────────────────────────────────────────────────────── text

def _draw_text(slide, el):
    w = _clamp(el.w or 5.0, 0.5, 13.333)
    h = _clamp(el.h or 1.0, 0.2, 7.5)
    x = _clamp(el.x, 0.0, 13.0)
    y = _clamp(el.y, 0.0, 7.3)

    # Minimal font o'lchami: sarlavha bo'lmasa kamida 13pt
    font_size = el.size or 14
    if not el.bold and font_size < 13:
        font_size = 13

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # Ko'p qatorli matnni \n bo'yicha ajratib, har birini alohida paragraf qilamiz
    lines = (el.text or "").split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = ALIGN_MAP.get(el.align or "left", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(_clamp(font_size, 8, 72))
        run.font.bold = el.bold
        run.font.italic = el.italic
        run.font.name = el.font or "Calibri"
        if el.color:
            run.font.color.rgb = _hex(el.color)


# ─────────────────────────────────────────────────────────── circle

def _draw_circle(slide, el):
    d = _clamp(el.d or 1.0, 0.1, 4.0)
    x = _clamp(el.x, 0.0, 13.0)
    y = _clamp(el.y, 0.0, 7.0)

    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.shadow.inherit = False
    if el.fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex(el.fill)
    else:
        shp.fill.background()
    shp.line.fill.background()


# ─────────────────────────────────────────────────────────── image

def _draw_image(slide, el, img_path: str):
    w = _clamp(el.w or 5.0, 0.5, 13.333)
    h = _clamp(el.h or 4.0, 0.5, 7.5)
    x = _clamp(el.x, 0.0, 13.0)
    y = _clamp(el.y, 0.0, 7.0)
    slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))


# ─────────────────────────────────────────────────────────── chart

def _draw_chart(slide, el):
    """python-pptx native chart + caption (izoh matni) chizadi."""
    w = _clamp(el.w or 7.0, 2.0, 13.0)
    h = _clamp(el.h or 4.0, 1.5, 6.5)   # caption uchun joy qoldiramiz
    x = _clamp(el.x, 0.0, 11.0)
    y = _clamp(el.y, 0.0, 6.0)

    chart_data = ChartData()

    categories = el.categories or ["A", "B", "C"]
    chart_data.categories = categories

    series_list = el.series or [{"name": "Ma'lumot", "values": [1] * len(categories)}]
    for s in series_list:
        if not isinstance(s, dict):
            continue
        name = s.get("name", "Series")
        values = s.get("values", [0] * len(categories))
        if len(values) < len(categories):
            values = list(values) + [0] * (len(categories) - len(values))
        chart_data.add_series(name, tuple(values[:len(categories)]))

    chart_type = CHART_TYPE_MAP.get(el.chart_type or "column", XL_CHART_TYPE.COLUMN_CLUSTERED)

    try:
        chart_frame = slide.shapes.add_chart(
            chart_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
        )
        chart = chart_frame.chart

        if el.chart_title:
            chart.has_title = True
            chart.chart_title.text_frame.text = el.chart_title
        else:
            chart.has_title = False

        if el.chart_type in ("pie", "donut"):
            plot = chart.plots[0]
            plot.has_data_labels = True

        if len(series_list) <= 1:
            chart.has_legend = False

        # ── Caption: diagramma ostida izoh matni ──────────────────
        caption_text = el.caption or el.chart_title or ""
        if caption_text:
            cap_y = _clamp(y + h + 0.08, 0.0, 7.3)
            # caption slayd ichida bo'lishini tekshir
            if cap_y + 0.35 <= 7.5:
                cap_tb = slide.shapes.add_textbox(
                    Inches(x), Inches(cap_y), Inches(w), Inches(0.35)
                )
                cap_tf = cap_tb.text_frame
                cap_tf.word_wrap = True
                cap_p = cap_tf.paragraphs[0]
                cap_p.alignment = PP_ALIGN.CENTER
                cap_run = cap_p.add_run()
                cap_run.text = f"📊 {caption_text}"
                cap_run.font.size = Pt(11)
                cap_run.font.italic = True
                cap_run.font.name = "Calibri"
                cap_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    except Exception as exc:
        log.error("Chart chizishda xato: %s", exc)
