"""
Template Service for Presentation Backgrounds
Manages background template selection and application
"""

import os
import logging
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

class TemplateService:
    """Manages presentation background templates"""
    
    def __init__(self):
        self.templates = {
            'template_1': {
                'name': {'uz': 'Oltin Naqsh', 'ru': 'Золотой Узор', 'en': 'Gold Ornamental'},
                'file': '1777532829029_1779186936497.png',
                'colors': {'title': RGBColor(102, 51, 0), 'text': RGBColor(51, 25, 0)}
            },
            'template_2': {
                'name': {'uz': 'Yashil To\'lqin', 'ru': 'Зелёная Волна', 'en': 'Green Wave'},
                'file': '1777532637847_1779186936522.png',
                'colors': {'title': RGBColor(0, 102, 51), 'text': RGBColor(0, 77, 38)}
            },
            'template_3': {
                'name': {'uz': 'Yashil Gradient', 'ru': 'Зелёный Градиент', 'en': 'Green Gradient'},
                'file': '1777532497212_1779186936543.png',
                'colors': {'title': RGBColor(51, 102, 51), 'text': RGBColor(51, 51, 51)}
            },
            'template_4': {
                'name': {'uz': 'Moviy To\'lqin', 'ru': 'Бирюзовая Волна', 'en': 'Teal Wave'},
                'file': '20260430_115315_0_io_thread_1777531658261_1779186936558.jpg',
                'colors': {'title': RGBColor(0, 102, 102), 'text': RGBColor(0, 77, 77)}
            },
            'template_5': {
                'name': {'uz': 'Binafsha To\'lqin', 'ru': 'Фиолетовая Волна', 'en': 'Purple Wave'},
                'file': '1777531410805_1779186936575.png',
                'colors': {'title': RGBColor(255, 255, 255), 'text': RGBColor(230, 200, 255)}
            },
            'template_6': {
                'name': {'uz': 'Moviy Bokeh', 'ru': 'Голубой Боке', 'en': 'Blue Bokeh'},
                'file': '1777531148413_1779186936597.png',
                'colors': {'title': RGBColor(0, 102, 153), 'text': RGBColor(0, 77, 128)}
            },
            'template_7': {
                'name': {'uz': 'Pushti Naqsh', 'ru': 'Розовый Узор', 'en': 'Pink Floral'},
                'file': '1777530910334_1779186936619.png',
                'colors': {'title': RGBColor(153, 51, 102), 'text': RGBColor(102, 51, 77)}
            },
            'template_8': {
                'name': {'uz': 'Pushti Naqsh 2', 'ru': 'Розовый Узор 2', 'en': 'Pink Floral 2'},
                'file': '1777530854578_1779186936643.png',
                'colors': {'title': RGBColor(153, 51, 102), 'text': RGBColor(102, 51, 77)}
            },
            'template_9': {
                'name': {'uz': 'Marmar', 'ru': 'Мрамор', 'en': 'Marble'},
                'file': '1777454132249_1779186936667.png',
                'colors': {'title': RGBColor(51, 77, 51), 'text': RGBColor(51, 51, 51)}
            },
            'template_10': {
                'name': {'uz': 'Ko\'k Mandala', 'ru': 'Голубая Мандала', 'en': 'Teal Mandala'},
                'file': '1777453990082_1779186936691.png',
                'colors': {'title': RGBColor(51, 102, 102), 'text': RGBColor(51, 51, 51)}
            },
            'template_11': {
                'name': {'uz': 'Oltin Geometrik', 'ru': 'Золотая Геометрия', 'en': 'Gold Geometric'},
                'file': '1777453725241_1779186936718.png',
                'colors': {'title': RGBColor(102, 77, 0), 'text': RGBColor(77, 51, 0)}
            },
            'template_12': {
                'name': {'uz': 'Kulrang Abstrakt', 'ru': 'Серый Абстракт', 'en': 'Gray Abstract'},
                'file': '1777453595870_1779186936741.png',
                'colors': {'title': RGBColor(51, 51, 51), 'text': RGBColor(77, 77, 77)}
            },
            'template_13': {
                'name': {'uz': 'Pushti-Ko\'k', 'ru': 'Розово-Голубой', 'en': 'Pink-Blue Soft'},
                'file': '1777453518093_1779186936762.png',
                'colors': {'title': RGBColor(153, 51, 102), 'text': RGBColor(51, 77, 153)}
            },
            'template_14': {
                'name': {'uz': 'Texno Olti Burchak', 'ru': 'Тех Шестиугольник', 'en': 'Tech Hexagon'},
                'file': '1777453465580_1779186936783.png',
                'colors': {'title': RGBColor(0, 102, 153), 'text': RGBColor(0, 77, 102)}
            },
            'template_15': {
                'name': {'uz': 'Moviy Texno', 'ru': 'Голубое Техно', 'en': 'Blue Tech Wave'},
                'file': '1777453324575_1779186936804.png',
                'colors': {'title': RGBColor(0, 102, 153), 'text': RGBColor(51, 77, 102)}
            },
            'template_16': {
                'name': {'uz': 'Ko\'k Olti Burchak', 'ru': 'Синий Шестиугольник', 'en': 'Blue Hexagon'},
                'file': '1777453197355_1779186936825.png',
                'colors': {'title': RGBColor(0, 102, 153), 'text': RGBColor(51, 102, 102)}
            },
            'template_17': {
                'name': {'uz': 'Oq Texno', 'ru': 'Белое Техно', 'en': 'White Tech'},
                'file': '1777453038132_1779186936849.png',
                'colors': {'title': RGBColor(0, 102, 153), 'text': RGBColor(51, 77, 102)}
            },
            'template_18': {
                'name': {'uz': 'Tabiat To\'lqin', 'ru': 'Природная Волна', 'en': 'Nature Wave'},
                'file': '1777452836463_1779186936871.png',
                'colors': {'title': RGBColor(51, 102, 51), 'text': RGBColor(51, 77, 51)}
            },
            'template_19': {
                'name': {'uz': 'Ko\'k To\'lqin', 'ru': 'Синяя Волна', 'en': 'Blue Wave Tech'},
                'file': '1777452798417_1779186936895.png',
                'colors': {'title': RGBColor(0, 77, 153), 'text': RGBColor(0, 51, 102)}
            },
            'template_20': {
                'name': {'uz': 'Minimalist Ko\'k', 'ru': 'Минималист Синий', 'en': 'Minimal Blue'},
                'file': 'copilot_image_1777451906343_1779186936915.jpeg',
                'colors': {'title': RGBColor(51, 102, 153), 'text': RGBColor(51, 77, 102)}
            }
        }
    
    def get_template_groups(self) -> List[List[Dict]]:
        """Get templates grouped by 5"""
        templates_list = list(self.templates.items())
        groups = []
        
        for i in range(0, len(templates_list), 5):
            group = []
            for j in range(i, min(i + 5, len(templates_list))):
                template_id, template_data = templates_list[j]
                group.append({
                    'id': template_id,
                    'name': template_data['name'],
                    'file': template_data['file']
                })
            groups.append(group)
        
        return groups
    
    def apply_template_to_slide(self, slide, template_id: str):
        """Apply template background to a slide"""
        try:
            if template_id not in self.templates:
                template_id = 'template_20'  # Default
                
            template = self.templates[template_id]
            
            # Add background image if specified
            if template['file']:
                # Use absolute path based on this file's location
                base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                bg_path = os.path.join(base_dir, 'attached_assets', template['file'])
                if os.path.exists(bg_path):
                    self._set_slide_background(slide, bg_path)
                else:
                    logger.warning(f"Background image not found: {bg_path}")
            
            return template
            
        except Exception as e:
            logger.error(f"Error applying template: {e}")
            return self.templates['template_20']  # Default
    
    def _set_slide_background(self, slide, image_path: str):
        """Set background image for a slide"""
        try:
            from pptx.util import Inches
            
            # Use standard slide dimensions (16:9)
            slide_width = Inches(13.33)
            slide_height = Inches(7.5)
            
            # Add background image at position 0,0 filling entire slide
            pic = slide.shapes.add_picture(
                image_path,
                0, 0,
                width=slide_width,
                height=slide_height
            )
            
            # Move background behind all real shapes. The shape tree begins with
            # framework elements (<nvGrpSpPr>, <grpSpPr>) followed by drawable
            # shapes (sp/pic/grpSp/graphicFrame). Hard-coding "insert at index 2"
            # is fragile — the framework prefix is sometimes 1 element, sometimes
            # 3 depending on the layout. Walk the tree and insert just before the
            # first drawable shape so the picture always lands at the back.
            shapes_tree = slide.shapes._spTree
            shapes_tree.remove(pic._element)

            DRAWABLE = {
                '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}sp',
                '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}pic',
            }
            # Use suffix matching to be namespace-agnostic.
            insert_idx = len(shapes_tree)
            for i, child in enumerate(shapes_tree):
                tag = child.tag.rsplit('}', 1)[-1] if '}' in child.tag else child.tag
                if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
                    insert_idx = i
                    break
            shapes_tree.insert(insert_idx, pic._element)
            
            logger.info(f"Successfully applied background image: {image_path}")
            
        except Exception as e:
            logger.error(f"Error setting slide background: {e}")
            # Alternative approach if first method fails
            try:
                from pptx.util import Inches
                slide.shapes.add_picture(
                    image_path,
                    0, 0, 
                    Inches(13.33), Inches(7.5)
                )
                logger.info(f"Applied background with alternative method: {image_path}")
            except Exception as e2:
                logger.error(f"Alternative background method also failed: {e2}")
    
    def get_template_colors(self, template_id: str) -> Dict:
        """Get color scheme for a template"""
        template = self.templates.get(template_id, self.templates['template_20'])
        return template['colors']
    
    def get_template_name(self, template_id: str, language: str = 'uz') -> str:
        """Get template name in specified language"""
        template = self.templates.get(template_id, self.templates['template_20'])
        name_dict = template['name']
        
        # Return name based on language, fallback to uzbek
        if isinstance(name_dict, dict):
            return name_dict.get(language, name_dict.get('uz', 'Standart'))
        else:
            # For backward compatibility if name is still string
            return name_dict