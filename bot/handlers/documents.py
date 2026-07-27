import logging
import json
import asyncio
import os
import uuid as uuid_mod
from aiogram import Router, F
from aiogram.types import Message, CallbackQuery, FSInputFile, InlineKeyboardMarkup, InlineKeyboardButton, WebAppInfo
from aiogram.fsm.context import FSMContext
from aiogram.filters import StateFilter

from bot.states import DocumentStates
import re as _re_plan
from bot.keyboards import get_slide_count_keyboard, get_page_count_keyboard, get_main_keyboard, get_template_keyboard, get_manual_input_keyboard, get_outline_review_keyboard, get_references_choice_keyboard, get_doc_language_keyboard, get_plan_slide_keyboard, get_icon_choice_keyboard, get_course_work_page_keyboard, get_diploma_work_page_keyboard, get_graduation_work_page_keyboard, get_dissertation_page_keyboard, get_payment_choice_keyboard, get_insufficient_balance_keyboard, get_back_inline_keyboard, get_article_page_keyboard, get_source_selection_keyboard, get_other_services_keyboard, get_extras_keyboard, get_gw_outline_choice_keyboard
from database.database import Database
from utils.security import sanitize_user_input, validate_topic_length
from services.ai_service import AIService, get_ai_service
from services.document_service import DocumentService, get_document_service
from services.template_service import TemplateService
from services.channel_service import ChannelService
from translations import get_text
from config import PRESENTATION_PRICES, DOCUMENT_PRICES, COURSE_WORK_PRICES, DIPLOMA_WORK_PRICES, GRADUATION_WORK_PRICES, DISSERTATION_PRICES, ARTICLE_PRICES, EXTRAS_PRICES, som_to_stars, TEMP_DIR

def _friendly_error(e: Exception) -> str:
    msg = str(e)
    if any(k in msg for k in ("JSON", "json", "parse", "char 0", "Expecting value", "empty response", "no JSON")):
        return "AI javob bermadi yoki noto'g'ri javob qaytardi. Qayta urinib ko'ring."
    if any(k in msg for k in ("timeout", "Timeout", "Connection", "connection")):
        return "Server bilan bog'lanishda muammo. Qayta urinib ko'ring."
    if any(k in msg for k in ("429", "rate limit", "quota")):
        return "AI cheklov tufayli to'xtatildi. Bir daqiqadan so'ng urinib ko'ring."
    return msg[:120]
import webapp
import time as _rl_time

router = Router()
logger = logging.getLogger(__name__)

# ── Per-user generation rate limit ────────────────────────────────────
# Document generation is expensive (AI tokens + Together image credits).
# Block a user from kicking off a second generation while their first one
# is still running, and throttle back-to-back successful generations.
_GEN_INFLIGHT: set[int] = set()
_GEN_LAST_AT: dict[int, float] = {}
_GEN_COOLDOWN_SEC = 8  # minimum gap between two finished generations


def _rl_messages(lang: str, kind: str) -> str:
    if kind == "inflight":
        return {
            "uz": "⏳ Avvalgi hujjat hali tayyorlanmoqda. Iltimos, kuting.",
            "ru": "⏳ Предыдущий документ ещё готовится. Пожалуйста, подождите.",
            "en": "⏳ Your previous document is still being generated. Please wait.",
        }.get(lang, "⏳ Avvalgi hujjat hali tayyorlanmoqda. Iltimos, kuting.")
    return {
        "uz": "🚦 Juda tez. Bir necha soniyadan so'ng qayta urinib ko'ring.",
        "ru": "🚦 Слишком часто. Попробуйте через несколько секунд.",
        "en": "🚦 Too fast. Try again in a few seconds.",
    }.get(lang, "🚦 Juda tez. Bir necha soniyadan so'ng qayta urinib ko'ring.")


def _rate_limit_check(user_id: int, lang: str = "uz") -> str | None:
    """Return None if the user may start a generation, or a localized error string."""
    if not user_id:
        return None
    if user_id in _GEN_INFLIGHT:
        return _rl_messages(lang, "inflight")
    last = _GEN_LAST_AT.get(user_id, 0.0)
    gap = _rl_time.time() - last
    if gap < _GEN_COOLDOWN_SEC:
        return _rl_messages(lang, "cooldown")
    return None


class _RateLimitSlot:
    """Async context manager that holds a generation slot for one user."""
    def __init__(self, user_id: int):
        self.user_id = user_id
    def __enter__(self):
        if self.user_id:
            _GEN_INFLIGHT.add(self.user_id)
        return self
    def __exit__(self, exc_type, exc, tb):
        if self.user_id:
            _GEN_INFLIGHT.discard(self.user_id)
            _GEN_LAST_AT[self.user_id] = _rl_time.time()
        return False

_BOOK_MODE_INSTRUCTIONS = {
    "uz": (
        "[KITOB ASOSIDA] Quyidagi kitob mazmuni ASOSIY MANBA hisoblanadi. "
        "MUHIM: Mavzu bir nechta tomonni qamrashi mumkin (masalan: 'O'zbekiston va Germaniya qonunchiligi taqqoslamasi'). "
        "Kitob faqat bir tomonni yoritishi mumkin. Siz BARCHA tomonlarni to'liq yozishingiz SHART:\n"
        "- Kitobda bor bo'lgan mavzular: kitob ma'lumotiga tayaning.\n"
        "- Kitobda YO'Q bo'lgan mavzular (masalan, boshqa davlat, boshqa soha): "
        "o'z bilimingizdan TO'LIQ va BATAFSIL yozing — bu bo'limlarni qoldirib bo'lmaydi.\n"
        "Kitob ma'lumotiga zid keluvchi hech narsa yozmang. "
        "AI detektorlarga ilinmaydigan, talaba o'zi yozgandek tabiiy uslubda bo'lsin. "
        "Barcha matn O'ZBEK TILIDA bo'lishi shart."
    ),
    "ru": (
        "[НА ОСНОВЕ КНИГИ] Содержание книги ниже является ОСНОВНЫМ ИСТОЧНИКОМ. "
        "ВАЖНО: Тема может охватывать несколько сторон (например: 'Сравнение законодательства Узбекистана и Германии'). "
        "Книга может освещать только одну сторону. Вы ОБЯЗАНЫ написать ВСЕ части темы полностью:\n"
        "- Темы, которые ЕСТЬ в книге: опирайтесь на содержание книги.\n"
        "- Темы, которых НЕТ в книге (например, другая страна, другая область): "
        "напишите ПОЛНЫЙ, ПОДРОБНЫЙ раздел из общих знаний — эти разделы нельзя пропускать.\n"
        "Не добавляйте ничего, что противоречит книге. "
        "Текст должен быть естественным, как будто студент сам прочитал книгу и сделал выводы. "
        "Весь текст должен быть на РУССКОМ ЯЗЫКЕ."
    ),
    "en": (
        "[BOOK-BASED] The book content below is the PRIMARY SOURCE. "
        "IMPORTANT: The topic may cover multiple subjects (e.g. 'Comparison of Uzbekistan and Germany law'). "
        "The book may only cover one of them. You MUST write ALL parts of the topic completely:\n"
        "- Subjects FOUND in the book: rely on the book content.\n"
        "- Subjects NOT in the book (e.g. a different country, a different field): "
        "write a COMPLETE, DETAILED section from your general knowledge — do NOT skip these.\n"
        "Never add anything that contradicts the book. "
        "Write in a natural student-like style. "
        "All text MUST be in ENGLISH."
    ),
}

_COMPARISON_MARKERS = (
    " va ", " и ", " vs ", " vs. ", " versus ",
    " taqqos", " solishti", " сравн", " сопостав",
    " comparison", " compare",
)

_COMPARATIVE_EXTRA = {
    "uz": (
        "\n[TAQQOSLAMA MAVZU ANIQLANADI] Ushbu mavzu ikki yoki undan ortiq tomonni taqqoslashni talab qiladi. "
        "Hujjatning HAR BIR bo'limi to'liq yozilishi SHART. "
        "Kitobda ko'rsatilmagan har qanday tomon (davlat, tashkilot, tizim va h.k.) uchun "
        "o'z bilimingizdan to'liq va batafsil yozing. Bo'sh qoldirmang, 'ma'lumot yo'q' demang."
    ),
    "ru": (
        "\n[СРАВНИТЕЛЬНАЯ ТЕМА ОБНАРУЖЕНА] Эта тема требует сравнения двух или более сторон. "
        "КАЖДЫЙ раздел документа должен быть написан полностью. "
        "Для любой стороны, отсутствующей в книге (страна, организация, система и т.д.), "
        "напишите полный подробный раздел из общих знаний. Не оставляйте пустых мест."
    ),
    "en": (
        "\n[COMPARATIVE TOPIC DETECTED] This topic requires comparing two or more sides. "
        "EVERY section of the document MUST be written in full. "
        "For any side not covered in the book (country, system, organization, etc.), "
        "write a complete detailed section from general knowledge. Do NOT leave any part empty."
    ),
}


def _is_comparative_topic(topic: str) -> bool:
    """Return True if the topic appears to compare two or more subjects."""
    lower = topic.lower()
    return any(marker in lower for marker in _COMPARISON_MARKERS)


def _build_book_topic(topic: str, book_content: str, doc_lang: str) -> str:
    instructions = _BOOK_MODE_INSTRUCTIONS.get(doc_lang, _BOOK_MODE_INSTRUCTIONS["uz"])
    extra = ""
    if _is_comparative_topic(topic):
        extra = _COMPARATIVE_EXTRA.get(doc_lang, _COMPARATIVE_EXTRA["uz"])
    return f"{topic}\n\n{instructions}{extra}\n\nBOOK CONTENT:\n{book_content}"

_EDIT_BTN_TEXT = {"uz": "✏️ Faylni tahrirlash", "ru": "✏️ Редактировать файл", "en": "✏️ Edit file"}


# Edit-token TTL — used ONLY by the standalone edit-file menu flow.
# Generated documents no longer carry an edit button (anonymity policy).
_EDIT_TOKEN_TTL = 1800  # 30 minutes


def _make_edit_keyboard(file_path: str, topic: str, doc_lang: str, user_lang: str, user_id: int, chat_id: int):
    """Create inline keyboard with Mini App edit button, register token. Returns (keyboard, token)."""
    import time as _time
    token = str(uuid_mod.uuid4())
    webapp.DOC_TOKENS[token] = {
        "file_path": file_path,
        "topic": topic,
        "doc_lang": doc_lang,
        "user_lang": user_lang,
        "user_id": user_id,
        "chat_id": chat_id,
        "_expires": _time.time() + _EDIT_TOKEN_TTL,
    }
    webapp.save_tokens_to_disk()
    domain = webapp.WEBAPP_DOMAIN
    url = f"https://{domain}/edit?token={token}"
    btn_text = _EDIT_BTN_TEXT.get(user_lang, _EDIT_BTN_TEXT["uz"])
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=btn_text, web_app=WebAppInfo(url=url))]
    ])
    return kb, token


def _attach_and_schedule(token: str, message_id: int, delay: int = _EDIT_TOKEN_TTL):
    """Store message_id in token and schedule auto-removal of button after delay."""
    if token in webapp.DOC_TOKENS:
        webapp.DOC_TOKENS[token]["message_id"] = message_id
        webapp.save_tokens_to_disk()
    asyncio.create_task(_expire_token(token, delay=delay))


def _safe_remove_file(file_path: str) -> None:
    """Best-effort delete of a generated file after delivery. Never raises."""
    if not file_path:
        return
    try:
        if os.path.exists(file_path):
            os.remove(file_path)
            logger.info(f"Deleted delivered file: {file_path}")
    except Exception as e:
        logger.warning(f"Failed to delete {file_path}: {e}")


async def _expire_token(token: str, delay: int = _EDIT_TOKEN_TTL):
    """Remove token after delay seconds, delete the edit button and the generated file."""
    await asyncio.sleep(delay)
    data = webapp.DOC_TOKENS.pop(token, None)
    webapp.save_tokens_to_disk()
    if data and webapp.BOT:
        try:
            chat_id = data.get("chat_id")
            message_id = data.get("message_id")
            if chat_id and message_id:
                await webapp.BOT.edit_message_reply_markup(
                    chat_id=chat_id,
                    message_id=message_id,
                    reply_markup=None
                )
        except Exception:
            pass
    # Delete the generated file to free disk space
    if data:
        file_path = data.get("file_path", "")
        if file_path and os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"Expired token cleanup: removed {file_path}")
            except Exception:
                pass


def _create_edited_docx(content: str, topic: str) -> str:
    """Create a new DOCX from edited plain text, return file path"""
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    lines = content.split("\n")
    first_non_empty = True
    for line in lines:
        if first_non_empty and line.strip():
            p = doc.add_heading(line.strip(), level=1)
            first_non_empty = False
        else:
            para = doc.add_paragraph(line)
            para.style = doc.styles["Normal"]
    fname = f"edited_{uuid_mod.uuid4().hex[:8]}.docx"
    out_path = os.path.join(TEMP_DIR, fname)
    os.makedirs(TEMP_DIR, exist_ok=True)
    doc.save(out_path)
    return out_path


_NS_RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def _create_edited_pptx(content: str, original_path: str) -> str:
    """Replace text in original PPTX template with edited content, return new file path"""
    from pptx import Presentation
    prs = Presentation(original_path)
    sld_ids = list(prs.slides._sldIdLst)

    import re
    sections = re.split(r"=== Slayd \d+ ===\s*\n?", content)
    sections = [s.strip() for s in sections if s.strip()]

    for i, sld_id in enumerate(sld_ids):
        if i >= len(sections):
            break
        rId = sld_id.get(_NS_RID)
        try:
            slide = prs.slides.part.related_slide(rId)
        except Exception:
            continue
        new_text = sections[i]
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            if not tf.text.strip():
                continue
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            if tf.paragraphs:
                tf.paragraphs[0].runs[0].text = new_text if tf.paragraphs[0].runs else ""
                if not tf.paragraphs[0].runs:
                    from pptx.util import Pt
                    run = tf.paragraphs[0].add_run()
                    run.text = new_text
            break

    fname = f"edited_{uuid_mod.uuid4().hex[:8]}.pptx"
    out_path = os.path.join(TEMP_DIR, fname)
    os.makedirs(TEMP_DIR, exist_ok=True)
    prs.save(out_path)
    return out_path

# Promokod handlers moved to settings

# Document type mapping
DOCUMENT_TYPES = {
    "🌟 Taqdimot": "presentation",
    "🌟 Презентация": "presentation",
    "🌟 Presentation": "presentation",
    "💥 Mustaqil ish": "independent_work",
    "💥 Самостоятельная работа": "independent_work",
    "💥 Independent Work": "independent_work",
    "⚡ Referat": "referat",
    "⚡ Реферат": "referat",
    "⚡ Research Paper": "referat",
    "✨ Kurs ishi": "course_work",
    "✨ Курсовая работа": "course_work",
    "✨ Course Work": "course_work",
    "🏆 Diplom ishi": "bitiruv_ishi",
    "🏆 Дипломная работа": "bitiruv_ishi",
    "🏆 Diploma Work": "bitiruv_ishi",
    "📝 Tezis": "tezis",
    "📝 Тезис": "tezis",
    "📝 Thesis": "tezis",
    "📰 Maqola": "maqola",
    "📰 Статья": "maqola",
    "📰 Article": "maqola",
    "🔬 Mahsus ishlanma": "mahsus_ishlanma",
    "🔬 Специальная разработка": "mahsus_ishlanma",
    "🔬 Special Project": "mahsus_ishlanma",
}

@router.message(F.text.in_(list(DOCUMENT_TYPES.keys())))
async def handle_document_type_selection(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Handle document type selection from main menu"""
    import traceback
    try:
        # IMPORTANT: Clear any active state first
        await state.clear()

        # Auto-register user if not in DB yet
        if user is None:
            await db.create_user(
                telegram_id=message.from_user.id,
                username=message.from_user.username,
                first_name=message.from_user.first_name,
                language="uz",
            )
            user = await db.get_user(message.from_user.id)

        # Check channel subscription (use Telegram user id — safe even if user is still None)
        channels = await db.get_active_channels()
        if channels:
            channel_service = ChannelService(message.bot)
            is_subscribed = await channel_service.check_user_subscription(message.from_user.id, channels)

            if not is_subscribed:
                from bot.keyboards import get_subscription_check_keyboard
                await message.answer(
                    get_text(user_lang, "subscription_required"),
                    reply_markup=get_subscription_check_keyboard(user_lang, channels)
                )
                return

        doc_type = DOCUMENT_TYPES[message.text]

        # Feature gate: check if this doc type is currently enabled
        if doc_type == "mahsus_ishlanma":
            enabled = await db.get_feature_status("mahsus_ishlanma")
            if not enabled:
                await message.answer(
                    get_text(user_lang, "media_disabled"),
                    reply_markup=get_main_keyboard(user_lang, mahsus_ishlanma_enabled=False)
                )
                return

        await state.update_data(document_type=doc_type, source_step_visited=True)

        await message.answer(
            get_text(user_lang, "select_source"),
            reply_markup=get_source_selection_keyboard(user_lang)
        )
        await state.set_state(DocumentStates.waiting_for_source_selection)

    except Exception as e:
        logger.error(f"Error in document type selection: {e}\n{traceback.format_exc()}")
        await message.answer("❌ Xatolik yuz berdi. Qayta urinib ko'ring.")

@router.callback_query(F.data == "doc_source_topic", DocumentStates.waiting_for_source_selection)
async def handle_source_topic(callback: CallbackQuery, state: FSMContext, user_lang: str):
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass
    await state.update_data(book_content='', book_context=False, topic='')
    await callback.message.answer(
        get_text(user_lang, "select_doc_language"),
        reply_markup=get_doc_language_keyboard(user_lang, back_callback="back_to_source_selection")
    )
    await state.set_state(DocumentStates.waiting_for_doc_language)

@router.callback_query(F.data == "doc_source_book", DocumentStates.waiting_for_source_selection)
async def handle_source_book(callback: CallbackQuery, state: FSMContext, user_lang: str):
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass
    from services.url_book_service import INSTRUCTIONS
    instructions = INSTRUCTIONS.get(user_lang, INSTRUCTIONS["uz"])
    await callback.message.answer(
        instructions,
        parse_mode="HTML",
        reply_markup=get_back_inline_keyboard(user_lang, "back_to_source_selection")
    )
    await state.set_state(DocumentStates.waiting_for_book_url)

@router.callback_query(F.data == "doc_source_url", DocumentStates.waiting_for_source_selection)
async def handle_source_url(callback: CallbackQuery, state: FSMContext, user_lang: str):
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass
    from services.url_book_service import INSTRUCTIONS
    instructions = INSTRUCTIONS.get(user_lang, INSTRUCTIONS["uz"])
    await callback.message.answer(
        instructions,
        parse_mode="HTML",
        reply_markup=get_back_inline_keyboard(user_lang, "back_to_source_selection")
    )
    await state.set_state(DocumentStates.waiting_for_book_url)


@router.message(DocumentStates.waiting_for_book_url)
async def handle_book_url_input(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """URL(lar)ni tekshiradi va saqlaydi — matn mavzu kiritilgandan keyin AI orqali olinadi"""
    from services.url_book_service import extract_urls_from_text, validate_url, get_error_message

    text = (message.text or "").strip()
    if not text:
        return

    urls = extract_urls_from_text(text)

    if not urls:
        no_url_texts = {
            "uz": "❌ URL topilmadi. Manzil <code>https://</code> bilan boshlanishi kerak.\n\nQayta yuboring 👇",
            "ru": "❌ Ссылка не найдена. Адрес должен начинаться с <code>https://</code>.\n\nОтправьте ещё раз 👇",
            "en": "❌ No URL found. The address must start with <code>https://</code>.\n\nTry again 👇",
        }
        await message.answer(no_url_texts.get(user_lang, no_url_texts["uz"]), parse_mode="HTML")
        return

    invalid = []
    for url in urls:
        ok, reason = validate_url(url)
        if not ok:
            invalid.append((url, reason))

    if invalid and len(invalid) == len(urls):
        err_msg = get_error_message(invalid[0][1], user_lang)
        await message.answer(err_msg, parse_mode="HTML")
        return

    await state.update_data(
        book_urls=urls,
        book_content='',
        book_context=True,
        topic='',
    )

    n = len(urls)
    ask_topic_texts = {
        "uz": (
            f"✅ <b>{n} ta manzil qabul qilindi.</b>\n\n"
            f"Endi mavzuni kiriting — AI shu manzillardan aynan mavzuga oid ma'lumotlarni oladi:"
        ),
        "ru": (
            f"✅ <b>Принято {n} ссылки(-ок).</b>\n\n"
            f"Теперь введите тему — AI извлечёт из этих страниц только нужную информацию:"
        ),
        "en": (
            f"✅ <b>{n} link(s) accepted.</b>\n\n"
            f"Now enter the topic — AI will extract only relevant information from those pages:"
        ),
    }
    await message.answer(ask_topic_texts.get(user_lang, ask_topic_texts["uz"]), parse_mode="HTML")
    await state.set_state(DocumentStates.waiting_for_book_topic)


@router.callback_query(F.data == "back_to_source_selection")
async def back_to_source_selection_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    await callback.message.answer(
        get_text(user_lang, "select_source"),
        reply_markup=get_source_selection_keyboard(user_lang)
    )
    await state.set_state(DocumentStates.waiting_for_source_selection)

@router.message(DocumentStates.waiting_for_book_file)
async def handle_book_file_upload(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    doc = message.document
    if not doc:
        await message.answer(get_text(user_lang, "book_source_not_docx"))
        return

    file_name = doc.file_name or ""
    is_pdf = file_name.lower().endswith(".pdf")
    is_docx = file_name.lower().endswith(".docx")

    if not is_pdf and not is_docx:
        await message.answer(get_text(user_lang, "book_source_not_docx"))
        return

    wait_msg = await message.answer(get_text(user_lang, "book_source_processing"))
    try:
        os.makedirs(TEMP_DIR, exist_ok=True)
        ext = ".pdf" if is_pdf else ".docx"
        local_path = os.path.join(TEMP_DIR, f"bk_src_{message.from_user.id}_{doc.file_id[-8:]}{ext}")
        tg_file = await message.bot.get_file(doc.file_id)
        await message.bot.download_file(tg_file.file_path, local_path)

        if is_pdf:
            from services.book_translate_service import auto_convert_pdf_to_docx
            docx_path = await auto_convert_pdf_to_docx(local_path)
            try:
                os.remove(local_path)
            except Exception:
                pass
        else:
            docx_path = local_path

        from docx import Document as DocxDocument
        docx_doc = DocxDocument(docx_path)
        paragraphs_text = []
        word_count = 0
        topic_candidate = ""
        for para in docx_doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if not topic_candidate and len(text) > 3:
                topic_candidate = text[:150]
            paragraphs_text.append(text)
            word_count += len(text.split())
            if word_count > 15000:
                break

        book_content = "\n\n".join(paragraphs_text)
        if len(book_content) > 60000:
            book_content = book_content[:60000]

        try:
            os.remove(docx_path)
        except Exception:
            pass

        if not book_content or word_count < 10:
            await wait_msg.delete()
            await message.answer(get_text(user_lang, "book_source_error"))
            return

        if not topic_candidate:
            topic_candidate = os.path.splitext(file_name)[0][:100]

        await state.update_data(
            book_content=book_content,
            book_context=True,
            topic='',
        )

        await wait_msg.delete()
        await message.answer(
            get_text(user_lang, "book_source_ask_topic", word_count=word_count)
        )
        await state.set_state(DocumentStates.waiting_for_book_topic)

    except Exception as e:
        logger.error(f"Error processing book file for doc: {e}")
        try:
            await wait_msg.delete()
        except Exception:
            pass
        await message.answer(get_text(user_lang, "book_source_error"))

@router.message(DocumentStates.waiting_for_book_topic)
async def handle_book_topic_input(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Kitob/URL asosida mavzu kiritish — AI faqat mavzuga oid qismlarni oladi"""
    topic = (message.text or "").strip()
    if len(topic) < 3:
        await message.answer(get_text(user_lang, "topic_too_short"))
        return

    data = await state.get_data()
    raw_book_content = data.get("book_content", "")
    book_urls = data.get("book_urls", [])

    loading_texts = {
        "uz": "⏳ AI ma'lumotlarni tahlil qilmoqda, iltimos kuting...",
        "ru": "⏳ AI анализирует данные, подождите...",
        "en": "⏳ AI is analysing the data, please wait...",
    }
    wait_msg = await message.answer(loading_texts.get(user_lang, loading_texts["uz"]))

    focused_content = raw_book_content

    # URL manbalari: avval yuklab, keyin AI orqali mavzuga oid qismlarni ajrat
    if book_urls:
        try:
            from services.url_book_service import fetch_multiple_urls, get_error_message
            multi = await fetch_multiple_urls(book_urls)
            combined_raw = multi["combined_content"]
            ok_results = [r for r in multi["results"] if r["ok"]]

            if not ok_results:
                await wait_msg.delete()
                fail_texts = {
                    "uz": "❌ Hech qaysi manzildan matn olib bo'lmadi. Boshqa manzillar bilan qayta urinib ko'ring.",
                    "ru": "❌ Не удалось загрузить ни одну страницу. Попробуйте другие ссылки.",
                    "en": "❌ Could not load any page. Try other links.",
                }
                await message.answer(fail_texts.get(user_lang, fail_texts["uz"]))
                return

            # AI orqali faqat mavzuga oid qismlarni ajrat
            from services.book_translate_service import extract_relevant_content_for_topic
            focused_content = await extract_relevant_content_for_topic(combined_raw, topic, user_lang)

        except Exception as e:
            logger.error(f"URL fetch/extract error in topic handler: {e}")
            await wait_msg.delete()
            error_texts = {
                "uz": "❌ Manzillardan ma'lumot olishda xatolik. Qayta urinib ko'ring.",
                "ru": "❌ Ошибка при получении данных. Попробуйте ещё раз.",
                "en": "❌ Error fetching data from links. Please try again.",
            }
            await message.answer(error_texts.get(user_lang, error_texts["uz"]))
            return

    # Fayl manbasi: avvalgi xatti-harakat
    elif raw_book_content:
        try:
            from services.book_translate_service import extract_relevant_content_for_topic
            focused_content = await extract_relevant_content_for_topic(raw_book_content, topic, user_lang)
        except Exception as e:
            logger.warning(f"Topic extraction failed, using full content: {e}")

    try:
        await wait_msg.delete()
    except Exception:
        pass

    await state.update_data(topic=topic, book_content=focused_content, book_urls=[])
    await message.answer(
        get_text(user_lang, "select_doc_language"),
        reply_markup=get_doc_language_keyboard(user_lang, back_callback="back_to_source_selection")
    )
    await state.set_state(DocumentStates.waiting_for_doc_language)


@router.callback_query(F.data.startswith("doc_lang_"), DocumentStates.waiting_for_doc_language)
async def handle_doc_language_selection(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Handle document language selection"""
    try:
        await callback.answer()
    except Exception:
        pass

    doc_lang = callback.data.split("_")[-1]
    await state.update_data(doc_language=doc_lang)

    try:
        await callback.message.delete()
    except Exception:
        pass

    data = await state.get_data()
    if data.get("topic") and data.get("book_context"):
        name_prompts = {
            "uz": "👤 Ism va Familiyangizni to'liq kiriting:\n\n(Masalan: Aliyev Jasur)",
            "ru": "👤 Введите ваше полное имя и фамилию:\n\n(Например: Иванов Иван)",
            "en": "👤 Enter your full name:\n\n(Example: John Smith)"
        }
        await callback.message.answer(
            name_prompts.get(doc_lang, name_prompts["uz"]),
            reply_markup=get_back_inline_keyboard(user_lang, "back_to_doc_lang")
        )
        await state.set_state(DocumentStates.waiting_for_author_name)
        return

    topic_prompts = {
        "uz": "📝 Mavzuni kiriting:",
        "ru": "📝 Введите тему:",
        "en": "📝 Enter the topic:"
    }
    await callback.message.answer(
        topic_prompts.get(doc_lang, topic_prompts["uz"]),
        reply_markup=get_back_inline_keyboard(user_lang, "back_to_doc_lang")
    )
    await state.set_state(DocumentStates.waiting_for_topic)

@router.message(DocumentStates.waiting_for_topic)
async def handle_topic_input(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Handle topic input from user"""
    try:
        # Sanitize user input to prevent injection attacks
        topic = sanitize_user_input(message.text, max_length=200)

        if not validate_topic_length(topic, min_length=3, max_length=200):
            await message.answer(get_text(user_lang, "topic_too_short"))
            return

        # Get document language from state
        data = await state.get_data()
        doc_lang = data.get('doc_language', user_lang)

        # Translate topic to document language if they differ
        if doc_lang != user_lang:
            translating_msgs = {
                "uz": "🔄 Mavzu tarjima qilinmoqda...",
                "ru": "🔄 Тема переводится...",
                "en": "🔄 Translating topic..."
            }
            wait_msg = await message.answer(translating_msgs.get(user_lang, translating_msgs["uz"]))
            try:
                ai_service = get_ai_service()
                topic = await ai_service.translate_topic(topic, doc_lang)
            except Exception:
                pass
            finally:
                await wait_msg.delete()

        await state.update_data(topic=topic)

        # Ask for author name in document language
        name_prompts = {
            "uz": "👤 Ism va Familiyangizni to'liq kiriting:\n\n(Masalan: Aliyev Jasur)",
            "ru": "👤 Введите ваше полное имя и фамилию:\n\n(Например: Иванов Иван)",
            "en": "👤 Enter your full name:\n\n(Example: John Smith)"
        }
        await message.answer(
            name_prompts.get(doc_lang, name_prompts["uz"]),
            reply_markup=get_back_inline_keyboard(user_lang, "back_to_topic")
        )
        await state.set_state(DocumentStates.waiting_for_author_name)

    except Exception as e:
        logger.error(f"Error handling topic input: {e}")
        await message.answer("❌ Xatolik yuz berdi. Qayta urinib ko'ring.")

@router.message(DocumentStates.waiting_for_author_name)
async def handle_author_name_input(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Handle author name input from user"""
    try:
        # Sanitize author name
        author_name = sanitize_user_input(message.text, max_length=100)

        if len(author_name.strip()) < 3:
            # ... (existing error logic)
            return

        await state.update_data(author_name=author_name.strip())

        # Get document type from state
        data = await state.get_data()
        doc_type = data.get('document_type')
        doc_lang = data.get('doc_language', user_lang)

        # If thesis, ask for university
        if doc_type == "tezis":
            await message.answer(
                get_text(doc_lang, "enter_university"),
                reply_markup=get_back_inline_keyboard(user_lang, "back_to_author_name")
            )
            await state.set_state(DocumentStates.waiting_for_university)
            return

        # Ask for slide/page count based on document type
        if doc_type == "presentation":
            await message.answer(
                get_text(doc_lang, "select_slide_count"),
                reply_markup=get_slide_count_keyboard(doc_lang)
            )
            await state.set_state(DocumentStates.waiting_for_slide_count)
        elif doc_type == "course_work":
            await message.answer(
                get_text(doc_lang, "select_page_count"),
                reply_markup=get_course_work_page_keyboard(doc_lang)
            )
            await state.set_state(DocumentStates.waiting_for_course_work_pages)
        elif doc_type == "diploma_work":
            await message.answer(
                get_text(doc_lang, "select_page_count"),
                reply_markup=get_diploma_work_page_keyboard(doc_lang)
            )
            await state.set_state(DocumentStates.waiting_for_diploma_work_pages)
        elif doc_type == "bitiruv_ishi":
            await message.answer(
                get_text(doc_lang, "select_page_count"),
                reply_markup=get_graduation_work_page_keyboard(doc_lang)
            )
            await state.set_state(DocumentStates.waiting_for_graduation_work_pages)
        elif doc_type == "dissertatsiya":
            await message.answer(
                get_text(doc_lang, "select_page_count"),
                reply_markup=get_dissertation_page_keyboard(doc_lang)
            )
            await state.set_state(DocumentStates.waiting_for_dissertation_pages)
        elif doc_type == "maqola":
            maqola_page_prompts = {
                "uz": "📄 Maqola hajmini tanlang:",
                "ru": "📄 Выберите объём статьи:",
                "en": "📄 Select article size:"
            }
            await message.answer(
                maqola_page_prompts.get(doc_lang, maqola_page_prompts["uz"]),
                reply_markup=get_article_page_keyboard(doc_lang)
            )
            await state.set_state(DocumentStates.waiting_for_page_count)
        else:  # referat or independent_work
            await message.answer(
                get_text(doc_lang, "select_page_count"),
                reply_markup=get_page_count_keyboard(doc_type, doc_lang)
            )
            await state.set_state(DocumentStates.waiting_for_page_count)

    except Exception as e:
        logger.error(f"Error handling author name input: {e}")
        await message.answer("❌ Xatolik yuz berdi. Qayta urinib ko'ring.")

@router.message(DocumentStates.waiting_for_university)
async def handle_university_input(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Handle university name input for thesis"""
    try:
        university = sanitize_user_input(message.text, max_length=200)
        if len(university.strip()) < 5:
            await message.answer("❌ Universitet nomi juda qisqa.")
            return

        await state.update_data(university=university.strip())
        data = await state.get_data()
        doc_lang = data.get('doc_language', user_lang)
        await message.answer(
            get_text(doc_lang, "enter_faculty"),
            reply_markup=get_back_inline_keyboard(user_lang, "back_to_university")
        )
        await state.set_state(DocumentStates.waiting_for_faculty)

    except Exception as e:
        logger.error(f"Error in university input: {e}")
        await message.answer("❌ Xatolik yuz berdi.")

@router.message(DocumentStates.waiting_for_faculty)
async def handle_faculty_input(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Handle faculty name input for thesis"""
    try:
        faculty = sanitize_user_input(message.text, max_length=200)
        if len(faculty.strip()) < 3:
            await message.answer("❌ Fakultet nomi juda qisqa.")
            return

        await state.update_data(faculty=faculty.strip())
        data = await state.get_data()
        doc_lang = data.get('doc_language', user_lang)
        await message.answer(
            get_text(doc_lang, "enter_group"),
            reply_markup=get_back_inline_keyboard(user_lang, "back_to_faculty")
        )
        await state.set_state(DocumentStates.waiting_for_group)

    except Exception as e:
        logger.error(f"Error in faculty input: {e}")
        await message.answer("❌ Xatolik yuz berdi.")

@router.message(DocumentStates.waiting_for_group)
async def handle_group_input(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Handle group name input for thesis"""
    try:
        group = sanitize_user_input(message.text, max_length=100)
        if len(group.strip()) < 2:
            await message.answer("❌ Guruh raqami juda qisqa.")
            return

        await state.update_data(group=group.strip())

        from config import DOCUMENT_PRICES
        price = DOCUMENT_PRICES.get("tezis", 5000)
        stars = som_to_stars(price)
        balance = user.balance if user else 0
        await state.update_data(price=price, doc_next_step="tezis_gen")
        await state.set_state(DocumentStates.waiting_for_payment)
        await message.answer(
            get_text(user_lang, "payment_choose", price=price, stars=stars, balance=balance),
            reply_markup=get_payment_choice_keyboard(user_lang, price, stars, balance, "pay_balance_doc", back_callback="back_from_doc_payment")
        )

    except Exception as e:
        logger.error(f"Error in group input: {e}")
        await message.answer("❌ Xatolik yuz berdi.")

async def generate_thesis(message: Message, state: FSMContext, db: Database, user_lang: str, user):
    """Generate thesis document"""
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        author_name = data['author_name']
        university = data['university']
        faculty = data.get('faculty', '')
        group = data.get('group', '')
        doc_lang = data.get('doc_language', user_lang)
        price = data['price']

        order_id = await db.create_document_order(user_id=user.id, document_type="tezis", topic=topic, specifications=json.dumps({"university": university, "faculty": faculty, "group": group}))

        ai_service = get_ai_service()
        content = await ai_service.generate_thesis_content(ai_topic, doc_lang)

        doc_service = get_document_service()
        file_path = await doc_service.create_thesis(topic, content, author_name, university, doc_lang, faculty=faculty, group=group)

        document = FSInputFile(file_path)
        await message.answer_document(document=document, caption=f"📝 {topic}")

        await db.update_document_order(order_id, "completed", file_path)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)
        await message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))
        await message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

    except Exception as e:
        logger.error(f"Error generating thesis: {e}")
        await message.answer("❌ Xatolik yuz berdi.")
        if 'order_id' in locals(): await db.update_document_order(order_id, "failed")
        _safe_remove_file(locals().get("file_path"))
    finally:
        await state.clear()

# Dynamic pricing helper function
def get_document_price(document_type: str, count_data: dict) -> int:
    """Get price based on document type and count"""
    if document_type == "presentation":
        slide_count = count_data.get('slide_count', 10)
        return PRESENTATION_PRICES.get(slide_count, 5000)
    elif document_type == "course_work":
        min_pages = count_data.get('min_pages', 15)
        max_pages = count_data.get('max_pages', 20)
        chapters = count_data.get('chapters', 3)
        page_key = f"{min_pages}_{max_pages}_{chapters}"
        return COURSE_WORK_PRICES.get(page_key, 15000)
    elif document_type == "diploma_work":
        min_pages = count_data.get('min_pages', 15)
        max_pages = count_data.get('max_pages', 20)
        chapters = count_data.get('chapters', 2)
        page_key = f"{min_pages}_{max_pages}_{chapters}"
        return DIPLOMA_WORK_PRICES.get(page_key, 15000)
    elif document_type == "bitiruv_ishi":
        min_pages = count_data.get('min_pages', 40)
        max_pages = count_data.get('max_pages', 50)
        chapters = count_data.get('chapters', 3)
        page_key = f"{min_pages}_{max_pages}_{chapters}"
        return GRADUATION_WORK_PRICES.get(page_key, 50000)
    elif document_type == "dissertatsiya":
        min_pages = count_data.get('min_pages', 70)
        max_pages = count_data.get('max_pages', 80)
        chapters = count_data.get('chapters', 3)
        page_key = f"{min_pages}_{max_pages}_{chapters}"
        return DISSERTATION_PRICES.get(page_key, 110000)
    elif document_type == "maqola":
        min_pages = count_data.get('min_pages', 4)
        max_pages = count_data.get('max_pages', 5)
        page_key = f"{min_pages}_{max_pages}"
        return ARTICLE_PRICES.get(page_key, 5000)
    else:  # independent_work or referat
        min_pages = count_data.get('min_pages', 10)
        max_pages = count_data.get('max_pages', 15)
        page_key = f"{min_pages}_{max_pages}"
        return DOCUMENT_PRICES.get(page_key, 5000)

# Subscription check helper function
async def check_user_subscription_required(message: Message, user, db: Database, user_lang: str) -> bool:
    """Check if user is subscribed to required channels"""
    channels = await db.get_active_channels()

    if not channels:
        return True  # No channels required

    channel_service = ChannelService(message.bot)
    is_subscribed = await channel_service.check_user_subscription(user.telegram_id, channels)

    if not is_subscribed:
        # Show subscription requirement
        from bot.keyboards import get_subscription_check_keyboard

        if user_lang == "uz":
            text = "❌ Hujjat yaratish uchun avval kanallarga a'zo bo'lishingiz shart!\n\n👇 Kanalga o'tish uchun tugmani bosing:"
        elif user_lang == "ru":
            text = "❌ Для создания документа сначала подпишитесь на каналы!\n\n👇 Нажмите кнопку для перехода в канал:"
        else:  # en
            text = "❌ To create document, you must subscribe to channels first!\n\n👇 Click the button to go to the channel:"

        await message.answer(
            text,
            reply_markup=get_subscription_check_keyboard(user_lang, channels)
        )
        return False

    return True

async def show_template_selection(message: Message, state: FSMContext, user_lang: str, group: int = 1, edit_message: bool = False):
    """Show all 20 templates in one overview image with numbered buttons"""
    try:
        from aiogram.types import FSInputFile

        # Send the overview image showing all 20 templates
        overview_image_path = "attached_assets/1777534399422_1779186783751.png"

        photo_msg = None
        if os.path.exists(overview_image_path):
            # Use translated text - only title, no description
            title_text = get_text(user_lang, "template_selection_title")

            photo_msg = await message.answer_photo(
                photo=FSInputFile(overview_image_path),
                caption=title_text,
                parse_mode="Markdown"
            )
        else:
            # Fallback if overview image not found - use translated fallback text
            text = get_text(user_lang, "template_selection_fallback")
            photo_msg = await message.answer(text, parse_mode="Markdown")

        # Send compact numbered keyboard with all 20 options
        from bot.keyboards import get_all_templates_keyboard
        keyboard = get_all_templates_keyboard(user_lang)
        keyboard_msg = await message.answer(
            get_text(user_lang, "template_select_number"), 
            reply_markup=keyboard,
            parse_mode="Markdown"
        )

        # Save message IDs to state for later deletion
        await state.update_data(
            template_photo_msg_id=photo_msg.message_id if photo_msg else None,
            template_keyboard_msg_id=keyboard_msg.message_id
        )

    except Exception as e:
        logger.error(f"Error in show_template_selection: {e}")
        await message.answer("❌ Xatolik yuz berdi. Qayta urinib ko'ring.")

@router.callback_query(F.data.startswith("template_group_"))
async def handle_template_group_navigation(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Handle template group navigation"""
    try:
        group = int(callback.data.split('_')[-1])
        await callback.answer()
        # Send new template selection as fresh message instead of editing
        await show_template_selection(callback.message, state, user_lang, group, edit_message=False)
    except Exception as e:
        logger.error(f"Error in template group navigation: {e}")
        await callback.answer("❌ Xatolik yuz berdi")

@router.callback_query(F.data.regexp(r"^template_\d+$"))
async def handle_template_selection(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle template selection and ask about plan slide"""
    try:
        # Extract template number from callback data (template_X)
        template_num = callback.data.split("_")[-1]
        template_id = f"template_{template_num}"
        await callback.answer()

        # Save selected template
        await state.update_data(selected_template=template_id)

        # Delete template selection messages (photo and keyboard)
        data = await state.get_data()
        photo_msg_id = data.get('template_photo_msg_id')
        keyboard_msg_id = data.get('template_keyboard_msg_id')
        doc_lang = data.get('doc_language', user_lang)

        try:
            if photo_msg_id:
                await callback.message.bot.delete_message(
                    chat_id=callback.message.chat.id,
                    message_id=photo_msg_id
                )
            if keyboard_msg_id:
                await callback.message.bot.delete_message(
                    chat_id=callback.message.chat.id,
                    message_id=keyboard_msg_id
                )
        except Exception as del_err:
            logger.warning(f"Could not delete template messages: {del_err}")

        # Auto-set plan slide and references, then ask about icons
        await state.update_data(add_plan_slide=True, add_references=True)

        await callback.message.answer(
            get_text(doc_lang, "add_icons_question"),
            reply_markup=get_icon_choice_keyboard(doc_lang)
        )
        await state.set_state(DocumentStates.waiting_for_icon_choice)

    except Exception as e:
        logger.error(f"Error in template selection: {e}")
        await callback.message.answer("❌ Xatolik yuz berdi")

@router.callback_query(F.data == "icon_yes", DocumentStates.waiting_for_icon_choice)
async def handle_icon_yes(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle user choosing to add icons"""
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass
    await state.update_data(add_icons=True)
    await callback.message.answer("⏳ " + get_text(user_lang, "generating"))
    await generate_presentation_with_template(callback, state, db, user_lang, user)

@router.callback_query(F.data == "icon_no", DocumentStates.waiting_for_icon_choice)
async def handle_icon_no(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle user choosing not to add icons"""
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass
    await state.update_data(add_icons=False)
    await callback.message.answer("⏳ " + get_text(user_lang, "generating"))
    await generate_presentation_with_template(callback, state, db, user_lang, user)

@router.callback_query(F.data == "plan_slide_yes", DocumentStates.waiting_for_plan_slide_choice)
async def handle_plan_slide_yes(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Handle user choosing to add plan slide"""
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass

    await state.update_data(add_plan_slide=True)

    # Now ask about references
    data = await state.get_data()
    doc_lang = data.get('doc_language', user_lang)
    await callback.message.answer(
        get_text(doc_lang, "add_references_question"),
        reply_markup=get_references_choice_keyboard(doc_lang)
    )
    await state.set_state(DocumentStates.waiting_for_references_choice)

@router.callback_query(F.data == "plan_slide_no", DocumentStates.waiting_for_plan_slide_choice)
async def handle_plan_slide_no(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Handle user choosing not to add plan slide"""
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass

    await state.update_data(add_plan_slide=False)

    # Now ask about references
    data = await state.get_data()
    doc_lang = data.get('doc_language', user_lang)
    await callback.message.answer(
        get_text(doc_lang, "add_references_question"),
        reply_markup=get_references_choice_keyboard(doc_lang)
    )
    await state.set_state(DocumentStates.waiting_for_references_choice)

@router.callback_query(F.data == "add_references_yes", DocumentStates.waiting_for_references_choice)
async def handle_add_references_yes(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle user choosing to add references"""
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass

    # Save choice
    await state.update_data(add_references=True)

    # Start generation
    await callback.message.answer("⏳ " + get_text(user_lang, "generating"))
    await generate_presentation_with_template(callback, state, db, user_lang, user)

@router.callback_query(F.data == "add_references_no", DocumentStates.waiting_for_references_choice)
async def handle_add_references_no(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle user choosing not to add references"""
    try:
        await callback.answer()
    except Exception:
        pass
    try:
        await callback.message.delete()
    except Exception:
        pass

    # Save choice
    await state.update_data(add_references=False)

    # Start generation
    await callback.message.answer("⏳ " + get_text(user_lang, "generating"))
    await generate_presentation_with_template(callback, state, db, user_lang, user)

async def generate_presentation_with_template(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate presentation with selected template"""
    # Per-user rate limit guard: prevents one user from spamming concurrent
    # generations (which would burn AI tokens and Together image credits).
    _rl_uid = user.telegram_id if user else 0
    _rl_err = _rate_limit_check(_rl_uid, user_lang)
    if _rl_err:
        try:
            await callback.answer(_rl_err, show_alert=True)
        except Exception:
            await callback.message.answer(_rl_err)
        return
    if _rl_uid:
        _GEN_INFLIGHT.add(_rl_uid)
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        slide_count = data['slide_count']
        template_id = data.get('selected_template', 'template_20')
        price = data.get('price', 0)
        manual_outline = data.get('manual_outline', [])
        add_references = data.get('add_references', False)
        add_plan_slide = data.get('add_plan_slide', False)
        add_icons = data.get('add_icons', True)
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)

        # Create order record
        specifications = json.dumps({
            "slide_count": slide_count,
            "template": template_id,
            "manual_outline": len(manual_outline) > 0,
            "add_plan_slide": add_plan_slide
        })
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="presentation",
            topic=topic,
            specifications=specifications
        )

        # Generate content with NEW AI BATCH SYSTEM using document language
        ai_service = get_ai_service()

        # If manual outline provided, use it
        if manual_outline:
            content = await ai_service.generate_presentation_with_manual_titles(
                ai_topic, manual_outline, doc_lang
            )
        else:
            content = await ai_service.generate_presentation_in_batches(ai_topic, slide_count, doc_lang)

        # Generate references if requested
        references = []
        if add_references:
            references = await ai_service.generate_references(ai_topic, doc_lang)

        # Generate plan items if requested
        plan_items = []
        if add_plan_slide:
            plan_items = await ai_service.generate_plan_items(ai_topic, doc_lang)

        # Validate AI response
        if not content or 'slides' not in content:
            logger.error(f"Invalid AI response from batch generation: {content}")
            content = {
                'slides': [
                    {'title': topic, 'content': f"Bu taqdimot {topic} mavzusida tayyorlangan.", 'layout_type': 'bullet_points', 'slide_number': 1},
                    {'title': 'Kirish', 'content': f"{topic} haqida batafsil ma'lumot va asosiy nuqtalar.", 'layout_type': 'bullet_points', 'slide_number': 2}
                ]
            }

        # Create presentation with selected template background
        doc_service = get_document_service()
        doc_service.use_icons = add_icons
        template_service = TemplateService()

        # Apply template to presentation with author name and plan slide
        file_path = await doc_service.create_presentation_with_template_background(
            topic, content, author_name, template_id, template_service, doc_lang, references, plan_items
        )

        # Verify file was created
        if not file_path or not os.path.exists(file_path):
            logger.error(f"Presentation file not created or not found: {file_path}")
            raise Exception(f"File not created: {file_path}")

        # Get template name for caption
        template_name = template_service.get_template_name(template_id, user_lang)

        # Send file FIRST - only proceed if successful
        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=get_text(user_lang, "document_ready_caption", 
                topic=topic,
                slide_count=slide_count,
                template=template_name
            ),
        )

        # File sent successfully - NOW update database and balance
        await db.update_document_order(order_id, "completed", file_path)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        # Send success message AFTER file is delivered
        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))

        # Send icon usage summary if icons were enabled
        if add_icons:
            used_icons = getattr(doc_service, '_last_used_icons', set())
            if used_icons:
                icon_names = ", ".join(sorted(used_icons))
                summary_msg = (
                    f"🎨 {len(used_icons)} ta unikal ikonka ishlatildi.\n"
                    f"Icons: {icon_names}"
                )
                logger.info(f"[Icon summary] user={user.telegram_id} icons={icon_names}")
                await callback.message.answer(summary_msg)
            else:
                logger.info(f"[Icon summary] user={user.telegram_id} no icons were added")
                await callback.message.answer("ℹ️ Ikonkalar qo'shilmadi (mos ikonka topilmadi).")

        # Send gentle reminder about content review
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Presentation successfully generated and sent: {file_path} for user {user.telegram_id}")
        await state.clear()

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        logger.error(f"Error generating presentation with template: {e}\n{error_details}")
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + _friendly_error(e),
            reply_markup=get_main_keyboard(user_lang)
        )
        try:
            if 'order_id' in locals():
                await db.update_document_order(order_id, "failed")
        except:
            pass
        _safe_remove_file(locals().get("file_path"))
        await state.clear()
    finally:
        if _rl_uid:
            _GEN_INFLIGHT.discard(_rl_uid)
            _GEN_LAST_AT[_rl_uid] = _rl_time.time()

@router.callback_query(F.data.startswith("slides_"), DocumentStates.waiting_for_slide_count)
async def handle_slide_count(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle slide count selection"""
    try:
        await callback.answer()
    except Exception:
        pass

    try:
        # Check if user exists
        if not user:
            await callback.message.answer(
                "❌ Xatolik yuz berdi. Iltimos, /start buyrug'ini bajaring.",
                reply_markup=get_main_keyboard(user_lang)
            )
            await state.clear()
            return

        slide_count = int(callback.data.split("_")[1])
        await state.update_data(slide_count=slide_count)

        # Calculate price based on slide count
        price = get_document_price("presentation", {"slide_count": slide_count})

        # Always show payment choice — balance shown inline
        stars = som_to_stars(price)
        balance = user.balance if user else 0
        await state.update_data(price=price, doc_next_step="presentation_template")
        await state.set_state(DocumentStates.waiting_for_payment)
        try:
            await callback.message.delete()
        except Exception:
            pass
        await callback.message.answer(
            get_text(user_lang, "payment_choose", price=price, stars=stars, balance=balance),
            reply_markup=get_payment_choice_keyboard(user_lang, price, stars, balance, "pay_balance_doc", back_callback="back_from_doc_payment")
        )
    except Exception as e:
        logger.error(f"handle_slide_count error: {e}", exc_info=True)

async def generate_article(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate an academic article (IMRAD structure) as a DOCX file"""
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        min_pages = data['min_pages']
        max_pages = data['max_pages']
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)
        price = data.get('price', 0)

        specifications = json.dumps({"min_pages": min_pages, "max_pages": max_pages})
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="maqola",
            topic=topic,
            specifications=specifications
        )

        ai_service = get_ai_service()
        content = await ai_service.generate_article_content(ai_topic, min_pages, max_pages, doc_lang)

        doc_service = get_document_service()
        file_path = await doc_service.create_article(topic, content, author_name, doc_lang)

        if not file_path or not os.path.exists(file_path):
            raise Exception(f"Article file not created: {file_path}")

        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=f"📰 {topic}\n📄 {min_pages}-{max_pages} varoq | IMRAD tuzilma",
        )

        await db.update_document_order(order_id, "completed", file_path)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Article generated: {file_path} for user {user.telegram_id}")

    except Exception as e:
        import traceback
        logger.error(f"Error generating article: {e}\n{traceback.format_exc()}")
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + _friendly_error(e),
            reply_markup=get_main_keyboard(user_lang)
        )
        if 'order_id' in locals():
            await db.update_document_order(order_id, "failed")
        _safe_remove_file(locals().get("file_path"))
    finally:
        await state.clear()


@router.callback_query(F.data.startswith("pages_"), DocumentStates.waiting_for_page_count)
async def handle_page_count(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle page count selection"""
    # Check if user exists
    if not user:
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, /start buyrug'ini bajaring.",
            reply_markup=get_main_keyboard(user_lang)
        )
        await state.clear()
        return

    page_range = callback.data.split("_")[1:]
    min_pages = int(page_range[0])
    max_pages = int(page_range[1])

    await state.update_data(min_pages=min_pages, max_pages=max_pages)

    # Calculate price based on page count
    data = await state.get_data()
    document_type = data['document_type']
    price = get_document_price(document_type, {"min_pages": min_pages, "max_pages": max_pages})

    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    # Show extras panel for supported document types
    if document_type in ("independent_work", "referat", "mahsus_ishlanma"):
        await state.update_data(base_price=price, doc_next_step="outline_choice")
        await state.set_state(DocumentStates.waiting_for_extras_choice)
        panel_title = get_text(user_lang, "extras_panel_title")
        await callback.message.answer(
            panel_title,
            reply_markup=get_extras_keyboard(user_lang, [], price)
        )
    else:
        # Non-extras types: go straight to payment
        stars = som_to_stars(price)
        balance = user.balance if user else 0
        await state.update_data(price=price, doc_next_step="outline_choice")
        await state.set_state(DocumentStates.waiting_for_payment)
        await callback.message.answer(
            get_text(user_lang, "payment_choose", price=price, stars=stars, balance=balance),
            reply_markup=get_payment_choice_keyboard(user_lang, price, stars, balance, "pay_balance_doc", back_callback="back_from_doc_payment")
        )

@router.callback_query(F.data.startswith("art_pages_"), DocumentStates.waiting_for_page_count)
async def handle_article_pages(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle article page count selection"""
    if not user:
        await callback.message.answer("❌ Xatolik yuz berdi.", reply_markup=get_main_keyboard(user_lang))
        await state.clear()
        return

    parts = callback.data.split("_")[2:]
    min_pages = int(parts[0])
    max_pages = int(parts[1])
    await state.update_data(min_pages=min_pages, max_pages=max_pages)

    price = get_document_price("maqola", {"min_pages": min_pages, "max_pages": max_pages})
    stars = som_to_stars(price)
    balance = user.balance if user else 0
    await state.update_data(price=price, doc_next_step="maqola_gen")
    await state.set_state(DocumentStates.waiting_for_payment)
    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)
    await callback.message.answer(
        get_text(user_lang, "payment_choose", price=price, stars=stars, balance=balance),
        reply_markup=get_payment_choice_keyboard(user_lang, price, stars, balance, "pay_balance_doc", back_callback="back_from_doc_payment")
    )


@router.callback_query(F.data.startswith("cw_pages_"), DocumentStates.waiting_for_course_work_pages)
async def handle_course_work_pages(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle course work page/chapter selection"""
    if not user:
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, /start buyrug'ini bajaring.",
            reply_markup=get_main_keyboard(user_lang)
        )
        await state.clear()
        return

    # Parse callback data: cw_pages_15_20_2 -> min=15, max=20, chapters=2
    parts = callback.data.split("_")[2:]  # ['15', '20', '2']
    min_pages = int(parts[0])
    max_pages = int(parts[1])
    chapters = int(parts[2])

    await state.update_data(min_pages=min_pages, max_pages=max_pages, chapters=chapters)

    # Calculate price
    price = get_document_price("course_work", {"min_pages": min_pages, "max_pages": max_pages, "chapters": chapters})

    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    await state.update_data(base_price=price, doc_next_step="course_work_gen")
    await state.set_state(DocumentStates.waiting_for_extras_choice)
    await callback.message.answer(
        get_text(user_lang, "extras_panel_title"),
        reply_markup=get_extras_keyboard(user_lang, [], price)
    )

@router.callback_query(F.data.startswith("dw_pages_"), DocumentStates.waiting_for_diploma_work_pages)
async def handle_diploma_work_pages(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle diploma work page/chapter selection"""
    if not user:
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, /start buyrug'ini bajaring.",
            reply_markup=get_main_keyboard(user_lang)
        )
        await state.clear()
        return

    # Parse callback data: dw_pages_15_20_2 -> min=15, max=20, chapters=2
    parts = callback.data.split("_")[2:]  # ['15', '20', '2']
    min_pages = int(parts[0])
    max_pages = int(parts[1])
    chapters = int(parts[2])

    await state.update_data(min_pages=min_pages, max_pages=max_pages, chapters=chapters)

    # Calculate price
    price = get_document_price("diploma_work", {"min_pages": min_pages, "max_pages": max_pages, "chapters": chapters})

    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    await state.update_data(base_price=price, doc_next_step="diploma_work_gen")
    await state.set_state(DocumentStates.waiting_for_extras_choice)
    await callback.message.answer(
        get_text(user_lang, "extras_panel_title"),
        reply_markup=get_extras_keyboard(user_lang, [], price)
    )


@router.callback_query(F.data.startswith("ds_pages_"), DocumentStates.waiting_for_dissertation_pages)
async def handle_dissertation_pages(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle master's dissertation page/chapter selection"""
    if not user:
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, /start buyrug'ini bajaring.",
            reply_markup=get_main_keyboard(user_lang)
        )
        await state.clear()
        return

    parts = callback.data.split("_")[2:]
    min_pages = int(parts[0])
    max_pages = int(parts[1])
    chapters = int(parts[2])

    await state.update_data(min_pages=min_pages, max_pages=max_pages, chapters=chapters)

    price = get_document_price("dissertatsiya", {"min_pages": min_pages, "max_pages": max_pages, "chapters": chapters})

    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    await state.update_data(base_price=price, doc_next_step="dissertation_gen")
    await state.set_state(DocumentStates.waiting_for_extras_choice)
    await callback.message.answer(
        get_text(user_lang, "extras_panel_title"),
        reply_markup=get_extras_keyboard(user_lang, [], price)
    )


@router.callback_query(F.data.startswith("gw_pages_"), DocumentStates.waiting_for_graduation_work_pages)
async def handle_graduation_work_pages(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle graduation (bitiruv) work page/chapter selection"""
    if not user:
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, /start buyrug'ini bajaring.",
            reply_markup=get_main_keyboard(user_lang)
        )
        await state.clear()
        return

    parts = callback.data.split("_")[2:]
    min_pages = int(parts[0])
    max_pages = int(parts[1])
    chapters = int(parts[2])

    await state.update_data(min_pages=min_pages, max_pages=max_pages, chapters=chapters)

    price = get_document_price("bitiruv_ishi", {"min_pages": min_pages, "max_pages": max_pages, "chapters": chapters})

    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    await state.update_data(base_price=price, doc_next_step="graduation_work_gen")

    # Show outline choice: Auto (AI) vs Manual (user enters full plan at once)
    await state.set_state(DocumentStates.waiting_for_gw_outline_choice)
    outline_prompt = {
        "uz": "📋 BMI uchun mundarijani qanday tuzamiz?",
        "ru": "📋 Как составим оглавление для ВКР?",
        "en": "📋 How should we create the table of contents for the thesis?",
    }
    await callback.message.answer(
        outline_prompt.get(user_lang, outline_prompt["uz"]),
        reply_markup=get_gw_outline_choice_keyboard(user_lang),
    )


@router.callback_query(F.data == "gw_outline_auto", DocumentStates.waiting_for_gw_outline_choice)
async def handle_gw_outline_auto(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """User chose AI-generated outline — proceed to extras."""
    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)
    data = await state.get_data()
    price = data.get("base_price", 0)
    await state.update_data(gw_manual_plan=None)
    await state.set_state(DocumentStates.waiting_for_extras_choice)
    await callback.message.answer(
        get_text(user_lang, "extras_panel_title"),
        reply_markup=get_extras_keyboard(user_lang, [], price),
    )


@router.callback_query(F.data == "gw_outline_manual", DocumentStates.waiting_for_gw_outline_choice)
async def handle_gw_outline_manual(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """User chose manual outline entry — show sample format and ask for full plan."""
    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)
    data = await state.get_data()
    chapters = data.get("chapters", 3)

    if user_lang == "ru":
        chap_word = "ГЛАВА"
        sub_example = ["Понятие и классификация", "Методы оценки", "Зарубежный опыт"]
        intro_note = (
            "✏️ Введите <b>полное оглавление</b> одним сообщением в формате ниже.\n"
            "Каждая глава — {n} подразделов (3.1, 3.2, 3.3).\n\n"
            "📌 <b>Образец:</b>"
        )
        hint = "\n\n⚠️ Соблюдайте формат. Каждую главу введите с <b>I BOB / II BOB</b> и т.д."
    elif user_lang == "en":
        chap_word = "CHAPTER"
        sub_example = ["Definition and classification", "Assessment methods", "International experience"]
        intro_note = (
            "✏️ Enter the <b>full table of contents</b> in one message using the format below.\n"
            "Each chapter has 3 subsections.\n\n"
            "📌 <b>Sample:</b>"
        )
        hint = "\n\n⚠️ Follow the format. Each chapter starts with <b>I BOB / II BOB</b>, etc."
    else:
        chap_word = "BOB"
        sub_example = ["Tushunchasi va tasnifi", "Baholash usullari", "Xorijiy tajriba"]
        intro_note = (
            "✏️ Quyidagi formatda <b>butun mundarijani</b> bitta xabarda yuboring.\n"
            "Har bir bob 3 ta kichik bo'limdan iborat bo'lishi kerak.\n\n"
            "📌 <b>Namuna:</b>"
        )
        hint = "\n\n⚠️ Formatga rioya qiling. Har bir bob <b>I BOB / II BOB</b> ko'rinishida boshlansin."

    roman_nums = ["I", "II", "III", "IV", "V"]
    lines = [intro_note]
    lines.append("<pre>")
    for i in range(1, chapters + 1):
        lines.append(f"{roman_nums[i-1]} {chap_word}. [Bob {i} sarlavhasi]")
        for j in range(1, 4):
            lines.append(f"  {i}.{j}. [{sub_example[(j-1) % len(sub_example)]}]")
        if i < chapters:
            lines.append("")
    lines.append("</pre>")
    lines.append(hint)

    await state.set_state(DocumentStates.waiting_for_gw_plan_text)
    await callback.message.answer("\n".join(lines), parse_mode="HTML")


@router.message(DocumentStates.waiting_for_gw_plan_text)
async def handle_gw_plan_text(message: Message, state: FSMContext, db: Database, user_lang: str, user):
    """Parse the user-submitted full plan and proceed to extras."""
    text = (message.text or "").strip()
    if len(text) < 20:
        err = {
            "uz": "❌ Reja juda qisqa. Iltimos, yuqoridagi namuna formatida butun mundarijani yuboring.",
            "ru": "❌ Оглавление слишком короткое. Пожалуйста, отправьте полное оглавление по образцу.",
            "en": "❌ The plan is too short. Please send the full table of contents following the sample format.",
        }
        await message.answer(err.get(user_lang, err["uz"]))
        return

    # Parse plan: detect chapter lines and subsection lines
    plan = []
    current_chapter = None
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        # Chapter line: Roman numeral + BOB/ГЛАВА/CHAPTER (case-insensitive)
        m_ch = _re_plan.match(
            r'^([IVXLCDM]+)\s+(?:BOB|ГЛАВА|CHAPTER)[.:\s]+(.+)$', line, _re_plan.IGNORECASE
        )
        if m_ch:
            if current_chapter and current_chapter["subsections"]:
                plan.append(current_chapter)
            current_chapter = {"title": m_ch.group(2).strip(), "subsections": []}
            continue
        # Subsection line: N.M. title
        m_sub = _re_plan.match(r'^\d+\.\d+\.?\s+(.+)$', line)
        if m_sub and current_chapter is not None:
            current_chapter["subsections"].append(m_sub.group(1).strip())

    if current_chapter and current_chapter["subsections"]:
        plan.append(current_chapter)

    data = await state.get_data()
    chapters = data.get("chapters", 3)

    if len(plan) < 1:
        err = {
            "uz": "❌ Reja to'g'ri formatda kiritilmadi. Har bir bob uchun:\n<b>I BOB. Sarlavha</b>\n  1.1. Kichik bo'lim\n  1.2. ...\n\nQaytadan yuboring.",
            "ru": "❌ Неверный формат. Для каждой главы:\n<b>I ГЛАВА. Название</b>\n  1.1. Подраздел\n  1.2. ...\n\nОтправьте заново.",
            "en": "❌ Incorrect format. For each chapter:\n<b>I CHAPTER. Title</b>\n  1.1. Subsection\n  1.2. ...\n\nPlease try again.",
        }
        await message.answer(err.get(user_lang, err["uz"]), parse_mode="HTML")
        return

    # Ensure each chapter has exactly 3 subsections (pad or trim)
    for ch in plan:
        while len(ch["subsections"]) < 3:
            ch["subsections"].append(f"{ch['title']} bo'yicha tahlil")
        ch["subsections"] = ch["subsections"][:3]

    # Trim or extend plan to match chapters count
    if len(plan) > chapters:
        plan = plan[:chapters]

    await state.update_data(gw_manual_plan=plan)

    price = data.get("base_price", 0)
    await state.set_state(DocumentStates.waiting_for_extras_choice)

    ok = {
        "uz": f"✅ Reja qabul qilindi! ({len(plan)} ta bob)\n\nEndi qo'shimcha xizmatlarni tanlang:",
        "ru": f"✅ Оглавление принято! ({len(plan)} гл.)\n\nВыберите дополнительные услуги:",
        "en": f"✅ Plan accepted! ({len(plan)} chapters)\n\nChoose additional services:",
    }
    await message.answer(ok.get(user_lang, ok["uz"]), reply_markup=get_extras_keyboard(user_lang, [], price))


@router.callback_query(F.data.startswith("extras_toggle_"), DocumentStates.waiting_for_extras_choice)
async def handle_extras_toggle(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Toggle an extra on/off in the extras multi-select panel."""
    key = callback.data.split("extras_toggle_", 1)[1]
    data = await state.get_data()
    selected: list = list(data.get("selected_extras", []))
    base_price: int = data.get("base_price", 0)
    if key in selected:
        selected.remove(key)
    else:
        selected.append(key)
    await state.update_data(selected_extras=selected)
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(
            reply_markup=get_extras_keyboard(user_lang, selected, base_price)
        )
    except Exception:
        pass


@router.callback_query(F.data == "extras_confirm", DocumentStates.waiting_for_extras_choice)
async def handle_extras_confirm(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """User confirmed their extras selection; compute final price and proceed to payment."""
    if not user:
        await callback.message.answer("❌ Xatolik yuz berdi.", reply_markup=get_main_keyboard(user_lang))
        await state.clear()
        return

    data = await state.get_data()
    selected: list = list(data.get("selected_extras", []))
    base_price: int = data.get("base_price", 0)
    extras_total = sum(EXTRAS_PRICES.get(k, 0) for k in selected)
    final_price = base_price + extras_total

    await state.update_data(doc_extras=selected, price=final_price)

    stars = som_to_stars(final_price)
    balance = user.balance if user else 0
    await state.set_state(DocumentStates.waiting_for_payment)
    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)
    await callback.message.answer(
        get_text(user_lang, "payment_choose", price=final_price, stars=stars, balance=balance),
        reply_markup=get_payment_choice_keyboard(user_lang, final_price, stars, balance, "pay_balance_doc", back_callback="back_from_doc_payment")
    )


# ─── HEAVY DOCUMENT QUEUE ──────────────────────────────────────────────────
# Heavy documents (kurs ishi, diplom ishi, dissertatsiya, bitiruv ishi) are
# routed through bot.queue_service so only ONE big generation runs at a time.
# This prevents OOM spikes on the production server (1vCPU/2GB).
from bot.queue_service import (
    HeavyDocTask,
    get_doc_queue,
    run_animated_status,
    build_initial_status_text,
)

# Per-user lock so a single user can't queue multiple heavy docs at once
# (would otherwise charge them twice). The asyncio.Lock makes the
# check-and-add atomic across concurrent callbacks from the same user.
_HEAVY_USERS: set[int] = set()
_HEAVY_USERS_LOCK = asyncio.Lock()

_HEAVY_BUSY_MSG = {
    "uz": "⏳ Sizning oldingi katta hujjatingiz hali navbatda yoki tayyorlanmoqda. Iltimos, u tugagunicha kuting.",
    "ru": "⏳ Ваш предыдущий большой документ ещё в очереди или готовится. Пожалуйста, дождитесь его завершения.",
    "en": "⏳ Your previous large document is still queued or being generated. Please wait for it to finish.",
}

# Configuration per heavy doc type: AI method, doc method, caption emoji
_HEAVY_DOC_CFG = {
    "course_work": {
        "ai_method": "generate_course_work_content",
        "doc_method": "create_course_work",
        "emoji": "📚",
        "label": "bo'lim",
    },
    "diploma_work": {
        "ai_method": "generate_diploma_work_content",
        "doc_method": "create_diploma_work",
        "emoji": "🎓",
        "label": "bo'lim",
    },
    "dissertatsiya": {
        "ai_method": "generate_dissertation_content",
        "doc_method": "create_dissertation",
        "emoji": "📕",
        "label": "bob",
    },
    "bitiruv_ishi": {
        "ai_method": "generate_graduation_work_content",
        "doc_method": "create_graduation_work",
        "emoji": "🎓",
        "label": "bob",
    },
}


async def _execute_heavy_generation(
    bot,
    chat_id: int,
    status_msg_id: int,
    user_telegram_id: int,
    user_lang: str,
    doc_type: str,
    topic: str,
    ai_topic: str,
    min_pages: int,
    max_pages: int,
    chapters: int,
    author_name: str,
    doc_lang: str,
    price: int,
    doc_extras: list,
    order_id: int,
    db: Database,
    gw_manual_plan: list = None,
):
    """Runs INSIDE the queue worker. Does AI generation, builds the .docx
    file, sends it to the user, updates DB, and reports errors."""
    cfg = _HEAVY_DOC_CFG[doc_type]
    try:
        ai_service = get_ai_service()
        ai_method = getattr(ai_service, cfg["ai_method"])

        if doc_type == "bitiruv_ishi":
            content = await ai_method(
                ai_topic, chapters, doc_lang,
                min_pages=min_pages, max_pages=max_pages,
                manual_plan=gw_manual_plan,
            )
        else:
            content = await ai_method(
                ai_topic, chapters, doc_lang,
                min_pages=min_pages, max_pages=max_pages,
            )

        doc_service = get_document_service()
        doc_method = getattr(doc_service, cfg["doc_method"])
        file_path = await doc_method(
            topic, content, author_name, doc_lang,
            extras=doc_extras or None,
        )

        if not file_path or not os.path.exists(file_path):
            raise Exception(f"File not created: {file_path}")

        # Delete the animated status message
        try:
            await bot.delete_message(chat_id=chat_id, message_id=status_msg_id)
        except Exception:
            pass

        # ── Charge BEFORE sending document, with refund on failure. ─────────
        # If we sent first and the deduction failed, the user would receive
        # the document for free. So we deduct first; if the file send fails
        # (network issue, Telegram error), we refund the same amount.
        balance_deducted = False
        if price and price > 0:
            try:
                await db.update_user_balance(user_telegram_id, -price)
                balance_deducted = True
            except Exception as charge_err:
                logger.error(f"Queue: failed to charge user {user_telegram_id}: {charge_err}")
                raise

        try:
            # Send the actual document
            caption = f"{cfg['emoji']} {topic}\n📄 {min_pages}-{max_pages} varoq | {chapters} {cfg['label']}"
            document = FSInputFile(file_path)
            await bot.send_document(
                chat_id=chat_id,
                document=document,
                caption=caption,
            )
        except Exception as send_err:
            # Refund the user — they were charged but never received the file.
            if balance_deducted:
                try:
                    await db.update_user_balance(user_telegram_id, price)
                    logger.warning(
                        f"Queue: refunded {price} to user {user_telegram_id} after send failure"
                    )
                except Exception as refund_err:
                    logger.error(
                        f"Queue: REFUND FAILED for user {user_telegram_id}, amount={price}: {refund_err}"
                    )
            raise send_err

        # Mark order completed only after the file was successfully delivered.
        await db.update_document_order(order_id, "completed", file_path)
        _safe_remove_file(file_path)

        # Final messages
        await bot.send_message(
            chat_id=chat_id,
            text=get_text(user_lang, "document_ready"),
        )
        await bot.send_message(
            chat_id=chat_id,
            text=get_text(user_lang, "document_reminder"),
            parse_mode="Markdown",
        )

        logger.info(f"Queue: {doc_type} done for user {user_telegram_id}, file: {file_path}")

    except Exception as e:
        import traceback
        logger.error(
            f"Queue: {doc_type} FAILED for user {user_telegram_id}: {e}\n{traceback.format_exc()}"
        )
        try:
            await bot.delete_message(chat_id=chat_id, message_id=status_msg_id)
        except Exception:
            pass
        try:
            await bot.send_message(
                chat_id=chat_id,
                text="❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: "
                     + _friendly_error(e),
            )
        except Exception:
            pass
        try:
            await db.update_document_order(order_id, "failed")
        except Exception:
            pass
        # Best-effort cleanup of any partial/undelivered file on disk.
        _safe_remove_file(locals().get("file_path"))
    finally:
        _HEAVY_USERS.discard(user_telegram_id)


async def _enqueue_heavy_doc(
    callback: CallbackQuery,
    state: FSMContext,
    db: Database,
    user_lang: str,
    user,
    doc_type: str,
):
    """Common entry point for the 4 heavy generation types.
    Extracts state data, creates the order record, sends the initial status
    message, starts the animated status loop, and submits the actual
    generation coroutine to the document queue."""
    cfg = _HEAVY_DOC_CFG[doc_type]
    bot = callback.bot
    chat_id = callback.message.chat.id
    user_telegram_id = user.telegram_id

    # Atomic per-user duplicate-submission guard: check + add inside the
    # single lock so two near-simultaneous callbacks can't both pass.
    async with _HEAVY_USERS_LOCK:
        if user_telegram_id in _HEAVY_USERS:
            await callback.message.answer(
                _HEAVY_BUSY_MSG.get(user_lang, _HEAVY_BUSY_MSG["uz"]),
                reply_markup=get_main_keyboard(user_lang),
            )
            await state.clear()
            return
        _HEAVY_USERS.add(user_telegram_id)

    order_id = None
    status_msg = None
    anim_task = None
    done_event = None
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        min_pages = data['min_pages']
        max_pages = data['max_pages']
        chapters = data['chapters']
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)
        price = data.get('price', 0)
        doc_extras = data.get('doc_extras') or []

        # Create order record up front
        specifications = json.dumps({
            "min_pages": min_pages,
            "max_pages": max_pages,
            "chapters": chapters,
        })
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type=doc_type,
            topic=topic,
            specifications=specifications,
        )

        queue = get_doc_queue()
        # Position before enqueue = currently pending + 1
        initial_position = queue.pending_count() + 1

        # Send the initial status message
        initial_text = build_initial_status_text(user_lang, doc_type, topic, initial_position)
        status_msg = await bot.send_message(
            chat_id=chat_id,
            text=initial_text,
            parse_mode="HTML",
        )

        task_id = uuid_mod.uuid4().hex
        done_event = asyncio.Event()

        # Start animated status loop in background
        anim_task = asyncio.create_task(run_animated_status(
            bot=bot,
            chat_id=chat_id,
            msg_id=status_msg.message_id,
            task_id=task_id,
            lang=user_lang,
            doc_type=doc_type,
            topic=topic,
            queue=queue,
            done_event=done_event,
        ))

        # (User was already marked busy under _HEAVY_USERS_LOCK above.)

        # Wrapped coroutine that runs in the queue worker
        async def _wrapped():
            try:
                await _execute_heavy_generation(
                    bot=bot,
                    chat_id=chat_id,
                    status_msg_id=status_msg.message_id,
                    user_telegram_id=user_telegram_id,
                    user_lang=user_lang,
                    doc_type=doc_type,
                    topic=topic,
                    ai_topic=ai_topic,
                    min_pages=min_pages,
                    max_pages=max_pages,
                    chapters=chapters,
                    author_name=author_name,
                    doc_lang=doc_lang,
                    price=price,
                    doc_extras=doc_extras,
                    order_id=order_id,
                    db=db,
                    gw_manual_plan=data.get('gw_manual_plan'),
                )
            finally:
                done_event.set()
                # Give the animation loop a moment to notice the event
                await asyncio.sleep(0.3)
                if not anim_task.done():
                    anim_task.cancel()

        task = HeavyDocTask(
            task_id=task_id,
            coro_factory=_wrapped,
            user_telegram_id=user_telegram_id,
            chat_id=chat_id,
            lang=user_lang,
            doc_type=doc_type,
            topic=topic,
            bot=bot,
            done_event=done_event,
        )
        await queue.enqueue(task)
        logger.info(f"Queue: enqueued {doc_type} for user {user_telegram_id} (position={initial_position})")

    except Exception as e:
        import traceback
        logger.error(f"Error enqueueing {doc_type}: {e}\n{traceback.format_exc()}")

        # ── Mandatory cleanup if anything went wrong during enqueue ──────────
        # Without this, the user could be left looking at an animated status
        # message that never updates (because no worker will ever set the
        # done_event), and would also be permanently stuck in _HEAVY_USERS.
        if done_event is not None and not done_event.is_set():
            done_event.set()
        if anim_task is not None and not anim_task.done():
            anim_task.cancel()
        if status_msg is not None:
            try:
                await bot.delete_message(
                    chat_id=chat_id,
                    message_id=status_msg.message_id,
                )
            except Exception:
                pass
        _HEAVY_USERS.discard(user_telegram_id)

        try:
            await callback.message.answer(
                "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: "
                + _friendly_error(e),
                reply_markup=get_main_keyboard(user_lang),
            )
        except Exception:
            pass
        if order_id is not None:
            try:
                await db.update_document_order(order_id, "failed")
            except Exception:
                pass
    finally:
        # State is cleared immediately so user can interact with the bot
        # while their heavy document is processing in the background queue.
        await state.clear()


async def generate_course_work(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Submit course work generation to the heavy-document queue."""
    await _enqueue_heavy_doc(callback, state, db, user_lang, user, "course_work")


async def generate_diploma_work(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Submit diploma work generation to the heavy-document queue."""
    await _enqueue_heavy_doc(callback, state, db, user_lang, user, "diploma_work")


async def generate_dissertation(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Submit dissertation generation to the heavy-document queue."""
    await _enqueue_heavy_doc(callback, state, db, user_lang, user, "dissertatsiya")


async def generate_graduation_work(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Submit graduation work generation to the heavy-document queue."""
    await _enqueue_heavy_doc(callback, state, db, user_lang, user, "bitiruv_ishi")


async def generate_presentation(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate presentation document"""
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        slide_count = data['slide_count']
        price = data.get('price', 0)

        # Create order record
        specifications = json.dumps({"slide_count": slide_count})
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="presentation",
            topic=topic,
            specifications=specifications
        )

        doc_lang = data.get('doc_language', user_lang)
        ai_service = get_ai_service()
        content = await ai_service.generate_presentation_in_batches(ai_topic, slide_count, doc_lang)

        # Validate AI response
        if not content or 'slides' not in content:
            logger.error(f"Invalid AI response from batch generation: {content}")
            # Create fallback content with new layout system
            content = {
                'slides': [
                    {'title': topic, 'content': f"Bu taqdimot {topic} mavzusida tayyorlangan.", 'layout_type': 'bullet_points', 'slide_number': 1},
                    {'title': 'Kirish', 'content': f"{topic} haqida batafsil ma'lumot va asosiy nuqtalar.", 'layout_type': 'bullet_points', 'slide_number': 2},
                    {'title': 'Asosiy qism', 'content': f"{topic}ning asosiy jihatlari va muhim ma'lumotlar.", 'layout_type': 'text_with_image', 'slide_number': 3}
                ]
            }

        # Create presentation file with NEW SYSTEM (DALL-E + 3 layouts)
        doc_service = get_document_service()
        file_path = await doc_service.create_new_presentation_system(topic, content, user.first_name or "", user_lang)

        # Verify file was created
        if not file_path or not os.path.exists(file_path):
            logger.error(f"Presentation file (legacy) not created: {file_path}")
            raise Exception(f"File not created: {file_path}")

        # Send file FIRST - only proceed if successful
        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=f"📊 {topic}",
        )

        # File sent successfully - NOW update database and balance
        await db.update_document_order(order_id, "completed", file_path)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        # Send success message AFTER file is delivered
        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))

        # Send gentle reminder about content review
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Presentation (legacy) generated and sent: {file_path} for user {user.telegram_id}")

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        logger.error(f"Error generating presentation: {e}\n{error_details}")
        await callback.message.edit_text("❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + str(e)[:100])
        await callback.message.answer("Asosiy menyu:", reply_markup=get_main_keyboard(user_lang))
        # Update order status
        if 'order_id' in locals():
            await db.update_document_order(order_id, "failed")
        _safe_remove_file(locals().get("file_path"))

    finally:
        await state.clear()

async def generate_independent_work_manual(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate independent work with manual outline"""
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        min_pages = data['min_pages']
        max_pages = data['max_pages']
        manual_outline = data.get('manual_outline', [])
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)
        doc_extras = data.get('doc_extras') or []

        # Create order record
        specifications = json.dumps({"min_pages": min_pages, "max_pages": max_pages, "manual_outline": True})
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="independent_work",
            topic=topic,
            specifications=specifications
        )

        # Use manual outline instead of AI-generated - use doc_lang
        ai_service = get_ai_service()

        # Generate content for each manually entered section
        sections = []
        for i, section_title in enumerate(manual_outline):
            section_content = await ai_service._generate_section_content(
                ai_topic, section_title, i + 1, len(manual_outline), "independent_work", doc_lang
            )
            sections.append({
                "title": section_title,
                "content": section_content
            })

        # Generate references
        references = await ai_service._generate_references(ai_topic, doc_lang)

        content = {
            "title": topic,
            "sections": sections,
            "references": references,
            "language": doc_lang,
            "author_name": author_name
        }

        # Create document file
        doc_service = get_document_service()
        file_path = await doc_service.create_independent_work(topic, content, extras=doc_extras or None)

        # Verify file was created
        if not file_path or not os.path.exists(file_path):
            logger.error(f"Manual independent work file not created: {file_path}")
            raise Exception(f"File not created: {file_path}")

        # Send file FIRST - only proceed if successful
        from aiogram.types import FSInputFile
        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=f"🎓 {topic}",
        )

        # File sent successfully - NOW update database and balance
        await db.update_document_order(order_id, "completed", file_path)
        price = data.get('price', 0)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        # Send success message AFTER file is delivered
        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Manual independent work generated and sent: {file_path} for user {user.telegram_id}")
        await state.clear()

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        logger.error(f"Error generating manual independent work: {e}\n{error_details}")
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + _friendly_error(e),
            reply_markup=get_main_keyboard(user_lang)
        )
        if 'order_id' in locals():
            await db.update_document_order(order_id, "failed")
        _safe_remove_file(locals().get("file_path"))
        await state.clear()

async def generate_referat_manual(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate referat with manual outline"""
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        min_pages = data['min_pages']
        max_pages = data['max_pages']
        manual_outline = data.get('manual_outline', [])
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)
        doc_extras = data.get('doc_extras') or []

        # Create order record
        specifications = json.dumps({"min_pages": min_pages, "max_pages": max_pages, "manual_outline": True})
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="referat",
            topic=topic,
            specifications=specifications
        )

        # Use manual outline instead of AI-generated - use doc_lang
        ai_service = get_ai_service()

        # Generate content for each manually entered section
        sections = []
        for i, section_title in enumerate(manual_outline):
            section_content = await ai_service._generate_section_content(
                ai_topic, section_title, i + 1, len(manual_outline), "referat", doc_lang
            )
            sections.append({
                "title": section_title,
                "content": section_content
            })

        # Generate references
        references = await ai_service._generate_references(ai_topic, doc_lang)

        content = {
            "title": topic,
            "sections": sections,
            "references": references,
            "language": doc_lang,
            "author_name": author_name
        }

        # Create document file
        doc_service = get_document_service()
        file_path = await doc_service.create_referat(topic, content, extras=doc_extras or None)

        # Verify file was created
        if not file_path or not os.path.exists(file_path):
            logger.error(f"Manual referat file not created: {file_path}")
            raise Exception(f"File not created: {file_path}")

        # Send file FIRST - only proceed if successful
        from aiogram.types import FSInputFile
        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=f"📄 {topic}",
        )

        # File sent successfully - NOW update database and balance
        await db.update_document_order(order_id, "completed", file_path)
        price = data.get('price', 0)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        # Send success message AFTER file is delivered
        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Manual referat generated and sent: {file_path} for user {user.telegram_id}")
        await state.clear()

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        logger.error(f"Error generating manual referat: {e}\n{error_details}")
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + _friendly_error(e),
            reply_markup=get_main_keyboard(user_lang)
        )
        if 'order_id' in locals():
            await db.update_document_order(order_id, "failed")
        _safe_remove_file(locals().get("file_path"))
        await state.clear()

async def generate_presentation_duplicate(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate presentation document (DUPLICATE - REDIRECTS TO MAIN)"""
    await generate_presentation(callback, state, db, user_lang, user)

async def generate_independent_work(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate independent work document"""
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        min_pages = data['min_pages']
        max_pages = data['max_pages']
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)
        doc_extras = data.get('doc_extras') or []

        # Create order record
        specifications = json.dumps({"min_pages": min_pages, "max_pages": max_pages})
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="independent_work",
            topic=topic,
            specifications=specifications
        )

        # Determine section count based on page range
        if max_pages <= 15:
            section_count = 6
        elif max_pages <= 20:
            section_count = 9
        elif max_pages <= 25:
            section_count = 12
        else:
            section_count = 15

        # Generate content with AI using old professional service - use doc_lang
        ai_service = get_ai_service()
        content = await ai_service.generate_document_content(
            ai_topic, section_count, "independent_work", doc_lang
        )

        # Add language and author info to content for template
        content['language'] = doc_lang
        content['author_name'] = author_name

        # Create document file using old professional service
        doc_service = get_document_service()
        file_path = await doc_service.create_independent_work(topic, content, extras=doc_extras or None)

        # Verify file was created
        if not file_path or not os.path.exists(file_path):
            logger.error(f"Independent work file not created: {file_path}")
            raise Exception(f"File not created: {file_path}")

        # Send file FIRST - only proceed if successful
        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=f"🎓 {topic}",
        )

        # File sent successfully - NOW update database and balance
        await db.update_document_order(order_id, "completed", file_path)
        price = data.get('price', 0)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        # Send success message AFTER file is delivered
        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))

        # Send gentle reminder about content review
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Independent work generated and sent: {file_path} for user {user.telegram_id}")

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        logger.error(f"Error generating independent work: {e}\n{error_details}")
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + _friendly_error(e),
            reply_markup=get_main_keyboard(user_lang)
        )
        if 'order_id' in locals():
            await db.update_document_order(order_id, "failed")
        _safe_remove_file(locals().get("file_path"))

    finally:
        await state.clear()

async def generate_referat(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate referat document"""
    try:
        data = await state.get_data()
        topic = data['topic']
        ai_topic = topic
        book_content = data.get('book_content', '')
        if book_content:
            ai_topic = _build_book_topic(topic, book_content, data.get('doc_language', user_lang))
        elif data.get('book_context'):
            ai_topic = f"{topic}\n\n{data['book_context']}"
        min_pages = data['min_pages']
        max_pages = data['max_pages']
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)
        doc_extras = data.get('doc_extras') or []

        # Create order record
        specifications = json.dumps({"min_pages": min_pages, "max_pages": max_pages})
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="referat",
            topic=topic,
            specifications=specifications
        )

        # Determine section count based on new page ranges
        if max_pages <= 15:
            section_count = 6
        elif max_pages <= 20:
            section_count = 9
        elif max_pages <= 25:
            section_count = 12
        else:
            section_count = 15

        # Generate content with AI using old professional service - use doc_lang
        ai_service = get_ai_service()
        content = await ai_service.generate_document_content(
            ai_topic, section_count, "referat", doc_lang
        )

        # Add language and author info to content for template
        content['language'] = doc_lang
        content['author_name'] = author_name

        # Create document file using old professional service  
        doc_service = get_document_service()
        file_path = await doc_service.create_referat(topic, content, extras=doc_extras or None)

        # Verify file was created
        if not file_path or not os.path.exists(file_path):
            logger.error(f"Referat file not created: {file_path}")
            raise Exception(f"File not created: {file_path}")

        # Send file FIRST - only proceed if successful
        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=f"📄 {topic}",
        )

        # File sent successfully - NOW update database and balance
        await db.update_document_order(order_id, "completed", file_path)
        price = data.get('price', 0)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        # Send success message AFTER file is delivered
        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))

        # Send gentle reminder about content review
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Referat generated and sent: {file_path} for user {user.telegram_id}")

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        logger.error(f"Error generating referat: {e}\n{error_details}")
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + _friendly_error(e),
            reply_markup=get_main_keyboard(user_lang)
        )
        if 'order_id' in locals():
            await db.update_document_order(order_id, "failed")
        _safe_remove_file(locals().get("file_path"))

    finally:
        await state.clear()



@router.callback_query(F.data == "outline_auto", DocumentStates.waiting_for_outline_choice)
async def handle_outline_auto(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle automatic outline generation"""
    await callback.answer()

    # Delete outline choice message
    await callback.message.delete()

    data = await state.get_data()
    document_type = data.get('document_type')

    if document_type == "presentation":
        # For presentation, show template selection directly (no manual option)
        await show_template_selection(callback.message, state, user_lang, group=1, edit_message=False)
        await state.set_state(DocumentStates.waiting_for_template)
    else:
        # For documents, start generation
        generation_msg = await callback.message.answer("⏳ " + get_text(user_lang, "generating"))

        if document_type == "independent_work":
            await generate_independent_work(callback, state, db, user_lang, user)
        elif document_type == "mahsus_ishlanma":
            await generate_mahsus_ishlanma(callback, state, db, user_lang, user)
        else:  # referat
            await generate_referat(callback, state, db, user_lang, user)

@router.callback_query(F.data == "cancel_document")
async def handle_cancel_document(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Handle document creation cancellation"""
    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    cancel_texts = {
        "uz": "❌ Hujjat yaratish bekor qilindi.",
        "ru": "❌ Создание документа отменено.",
        "en": "❌ Document creation cancelled."
    }

    await callback.message.answer(
        cancel_texts.get(user_lang, cancel_texts["uz"]),
        reply_markup=get_main_keyboard(user_lang)
    )
    await state.clear()

# ─── BACK NAVIGATION HANDLERS ────────────────────────────────────────────────

@router.callback_query(F.data == "back_to_main")
async def back_to_main_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Back to main menu"""
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await state.clear()
    _texts = {"uz": "🏠 Asosiy menyu", "ru": "🏠 Главное меню", "en": "🏠 Main menu"}
    await callback.message.answer(
        _texts.get(user_lang, _texts["uz"]),
        reply_markup=get_main_keyboard(user_lang)
    )

@router.callback_query(F.data == "back_to_doc_lang")
async def back_to_doc_lang_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Back to document language selection"""
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    data = await state.get_data()
    back_cb = "back_to_source_selection" if data.get("source_step_visited") else "back_to_main"
    await state.set_state(DocumentStates.waiting_for_doc_language)
    await callback.message.answer(
        get_text(user_lang, "select_doc_language"),
        reply_markup=get_doc_language_keyboard(user_lang, back_callback=back_cb)
    )

@router.callback_query(F.data == "back_to_topic")
async def back_to_topic_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Back to topic input"""
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    data = await state.get_data()
    doc_lang = data.get('doc_language', user_lang)
    topic_prompts = {
        "uz": "📝 Mavzuni kiriting:",
        "ru": "📝 Введите тему:",
        "en": "📝 Enter the topic:"
    }
    await state.set_state(DocumentStates.waiting_for_topic)
    await callback.message.answer(
        topic_prompts.get(doc_lang, topic_prompts["uz"]),
        reply_markup=get_back_inline_keyboard(user_lang, "back_to_doc_lang")
    )

@router.callback_query(F.data == "back_to_author_name")
async def back_to_author_name_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Back to author name input"""
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    data = await state.get_data()
    doc_lang = data.get('doc_language', user_lang)
    name_prompts = {
        "uz": "👤 Ism va Familiyangizni to'liq kiriting:\n\n(Masalan: Aliyev Jasur)",
        "ru": "👤 Введите ваше полное имя и фамилию:\n\n(Например: Иванов Иван)",
        "en": "👤 Enter your full name:\n\n(Example: John Smith)"
    }
    back_cb = "back_to_doc_lang" if data.get("book_content") or data.get("book_context") else "back_to_topic"
    await state.set_state(DocumentStates.waiting_for_author_name)
    await callback.message.answer(
        name_prompts.get(doc_lang, name_prompts["uz"]),
        reply_markup=get_back_inline_keyboard(user_lang, back_cb)
    )

@router.callback_query(F.data == "back_to_university")
async def back_to_university_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Back to university input (thesis only)"""
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    data = await state.get_data()
    doc_lang = data.get('doc_language', user_lang)
    await state.set_state(DocumentStates.waiting_for_university)
    await callback.message.answer(
        get_text(doc_lang, "enter_university"),
        reply_markup=get_back_inline_keyboard(user_lang, "back_to_author_name")
    )

@router.callback_query(F.data == "back_to_faculty")
async def back_to_faculty_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Back to faculty input (thesis only)"""
    await callback.answer()
    try:
        await callback.message.delete()
    except Exception:
        pass
    data = await state.get_data()
    doc_lang = data.get('doc_language', user_lang)
    await state.set_state(DocumentStates.waiting_for_faculty)
    await callback.message.answer(
        get_text(doc_lang, "enter_faculty"),
        reply_markup=get_back_inline_keyboard(user_lang, "back_to_university")
    )

@router.callback_query(F.data == "back_from_doc_payment")
async def back_from_doc_payment_handler(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Back from payment screen to count/page keyboard"""
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    data = await state.get_data()
    doc_next_step = data.get('doc_next_step', '')
    doc_type = data.get('document_type', '')
    doc_lang = data.get('doc_language', user_lang)

    if doc_next_step == "presentation_template":
        await state.set_state(DocumentStates.waiting_for_slide_count)
        await callback.message.answer(
            get_text(doc_lang, "select_slide_count"),
            reply_markup=get_slide_count_keyboard(doc_lang)
        )
    elif doc_next_step == "outline_choice":
        await state.set_state(DocumentStates.waiting_for_page_count)
        await callback.message.answer(
            get_text(doc_lang, "select_page_count"),
            reply_markup=get_page_count_keyboard(doc_type, doc_lang)
        )
    elif doc_next_step == "course_work_gen":
        await state.set_state(DocumentStates.waiting_for_course_work_pages)
        await callback.message.answer(
            get_text(doc_lang, "select_page_count"),
            reply_markup=get_course_work_page_keyboard(doc_lang)
        )
    elif doc_next_step == "diploma_work_gen":
        await state.set_state(DocumentStates.waiting_for_diploma_work_pages)
        await callback.message.answer(
            get_text(doc_lang, "select_page_count"),
            reply_markup=get_diploma_work_page_keyboard(doc_lang)
        )
    elif doc_next_step == "graduation_work_gen":
        await state.set_state(DocumentStates.waiting_for_graduation_work_pages)
        await callback.message.answer(
            get_text(doc_lang, "select_page_count"),
            reply_markup=get_graduation_work_page_keyboard(doc_lang)
        )
    elif doc_next_step == "dissertation_gen":
        await state.set_state(DocumentStates.waiting_for_dissertation_pages)
        await callback.message.answer(
            get_text(doc_lang, "select_page_count"),
            reply_markup=get_dissertation_page_keyboard(doc_lang)
        )
    elif doc_next_step == "tezis_gen":
        await state.set_state(DocumentStates.waiting_for_group)
        await callback.message.answer(
            get_text(doc_lang, "enter_group"),
            reply_markup=get_back_inline_keyboard(user_lang, "back_to_faculty")
        )

@router.callback_query(F.data == "back_from_outline")
async def back_from_outline_handler(callback: CallbackQuery, state: FSMContext, user_lang: str, db: Database, user):
    """Back from outline choice to payment screen"""
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    data = await state.get_data()
    price = data.get('price', 0)
    stars = som_to_stars(price)
    balance = user.balance if user else 0
    await state.set_state(DocumentStates.waiting_for_payment)
    await callback.message.answer(
        get_text(user_lang, "payment_choose", price=price, stars=stars, balance=balance),
        reply_markup=get_payment_choice_keyboard(user_lang, price, stars, balance, "pay_balance_doc", back_callback="back_from_doc_payment")
    )

@router.callback_query(F.data == "back_from_template")
async def back_from_template_handler(callback: CallbackQuery, state: FSMContext, user_lang: str, db: Database, user):
    """Back from template selection to payment screen"""
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    data = await state.get_data()
    price = data.get('price', 0)
    stars = som_to_stars(price)
    balance = user.balance if user else 0
    await state.set_state(DocumentStates.waiting_for_payment)
    await callback.message.answer(
        get_text(user_lang, "payment_choose", price=price, stars=stars, balance=balance),
        reply_markup=get_payment_choice_keyboard(user_lang, price, stars, balance, "pay_balance_doc", back_callback="back_from_doc_payment")
    )

# ─── END BACK NAVIGATION ──────────────────────────────────────────────────────

@router.callback_query(F.data == "outline_manual", DocumentStates.waiting_for_outline_choice)
async def handle_outline_manual(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Handle manual outline entry"""
    await callback.answer()

    # Delete outline choice message
    await callback.message.delete()

    data = await state.get_data()
    document_type = data.get('document_type')

    # Calculate how many sections/slides needed
    if document_type == "presentation":
        slide_count = data.get('slide_count', 10)
        await state.update_data(manual_outline=[], current_section=1, total_sections=slide_count)

        # Show instruction with total count
        await callback.message.answer(
            get_text(user_lang, "manual_outline_instruction_presentation", total_slides=slide_count)
        )
        await callback.message.answer(
            get_text(user_lang, "enter_slide_title", slide_num=1, total_slides=slide_count),
            reply_markup=get_manual_input_keyboard(user_lang)
        )
    else:
        # For documents, determine section count based on pages
        max_pages = data.get('max_pages', 15)
        if max_pages <= 15:
            section_count = 6
        elif max_pages <= 20:
            section_count = 9
        elif max_pages <= 25:
            section_count = 12
        else:
            section_count = 15

        await state.update_data(manual_outline=[], current_section=1, total_sections=section_count)

        # Show instruction with total count
        await callback.message.answer(
            get_text(user_lang, "manual_outline_instruction_document", total_sections=section_count)
        )
        await callback.message.answer(
            get_text(user_lang, "enter_section_title", section_num=1, total_sections=section_count),
            reply_markup=get_manual_input_keyboard(user_lang)
        )

    await state.set_state(DocumentStates.waiting_for_manual_outline)

@router.message(DocumentStates.waiting_for_manual_outline)
async def handle_manual_outline_input(message: Message, state: FSMContext, db: Database, user_lang: str, user):
    """Handle manual outline section/slide title input"""

    # Check if user wants to go back
    back_texts = ["🔙 Ortga qaytish", "🔙 Назад", "🔙 Back"]
    if message.text in back_texts:
        # Go back to outline choice
        from bot.keyboards import get_outline_choice_keyboard
        await message.answer(
            get_text(user_lang, "outline_choice"),
            reply_markup=get_outline_choice_keyboard(user_lang)
        )
        await state.set_state(DocumentStates.waiting_for_outline_choice)
        return

    data = await state.get_data()
    manual_outline = data.get('manual_outline', [])
    current_section = data.get('current_section', 1)
    total_sections = data.get('total_sections', 1)
    document_type = data.get('document_type')

    # Sanitize and validate outline input
    outline_text = sanitize_user_input(message.text, max_length=150)

    if not validate_topic_length(outline_text, min_length=2, max_length=150):
        await message.answer("❌ Mavzu juda qisqa yoki uzun. 2-150 belgi oralig'ida kiriting.")
        return

    # Add current title to outline
    manual_outline.append(outline_text)
    current_section += 1

    if current_section <= total_sections:
        # Ask for next section/slide
        await state.update_data(manual_outline=manual_outline, current_section=current_section)

        if document_type == "presentation":
            await message.answer(
                get_text(user_lang, "enter_slide_title", slide_num=current_section, total_slides=total_sections),
                reply_markup=get_manual_input_keyboard(user_lang)
            )
        else:
            await message.answer(
                get_text(user_lang, "enter_section_title", section_num=current_section, total_sections=total_sections),
                reply_markup=get_manual_input_keyboard(user_lang)
            )
    else:
        # All sections/slides collected - show review
        await state.update_data(manual_outline=manual_outline)

        # Format outline for display
        outline_text = ""
        for i, item in enumerate(manual_outline, 1):
            outline_text += f"{i}. {item}\n"

        # Show review with confirm/edit buttons
        await message.answer(
            get_text(user_lang, "outline_review", outline_list=outline_text),
            reply_markup=get_outline_review_keyboard(user_lang),
            parse_mode="Markdown"
        )
        await state.set_state(DocumentStates.waiting_for_outline_confirmation)

@router.callback_query(F.data == "confirm_outline", DocumentStates.waiting_for_outline_confirmation)
async def handle_confirm_outline(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle outline confirmation"""
    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    data = await state.get_data()
    document_type = data.get('document_type')

    await callback.message.answer(get_text(user_lang, "outline_complete"), reply_markup=get_main_keyboard(user_lang))

    if document_type == "presentation":
        # Show template selection for presentation
        await show_template_selection(callback.message, state, user_lang, group=1, edit_message=False)
        await state.set_state(DocumentStates.waiting_for_template)
    else:
        # Start document generation
        await callback.message.answer("⏳ " + get_text(user_lang, "generating"))

        if document_type == "independent_work":
            await generate_independent_work_manual(callback, state, db, user_lang, user)
        elif document_type == "mahsus_ishlanma":
            await generate_mahsus_ishlanma(callback, state, db, user_lang, user)
        else:  # referat
            await generate_referat_manual(callback, state, db, user_lang, user)

@router.callback_query(F.data == "edit_outline", DocumentStates.waiting_for_outline_confirmation)
async def handle_edit_outline(callback: CallbackQuery, state: FSMContext, user_lang: str):
    """Handle outline editing - restart from beginning"""
    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    data = await state.get_data()
    document_type = data.get('document_type')

    # Reset outline and start over
    if document_type == "presentation":
        slide_count = data.get('slide_count', 10)
        total_sections = slide_count
        await state.update_data(manual_outline=[], current_section=1, total_sections=total_sections)

        await callback.message.answer(
            get_text(user_lang, "manual_outline_instruction_presentation", total_slides=slide_count)
        )
        await callback.message.answer(
            get_text(user_lang, "enter_slide_title", slide_num=1, total_slides=slide_count),
            reply_markup=get_manual_input_keyboard(user_lang)
        )
    else:
        max_pages = data.get('max_pages', 15)
        if max_pages <= 15:
            section_count = 6
        elif max_pages <= 20:
            section_count = 9
        elif max_pages <= 25:
            section_count = 12
        else:
            section_count = 15

        await state.update_data(manual_outline=[], current_section=1, total_sections=section_count)

        await callback.message.answer(
            get_text(user_lang, "manual_outline_instruction_document", total_sections=section_count)
        )
        await callback.message.answer(
            get_text(user_lang, "enter_section_title", section_num=1, total_sections=section_count),
            reply_markup=get_manual_input_keyboard(user_lang)
        )

    await state.set_state(DocumentStates.waiting_for_manual_outline)

# Help button texts in different languages
HELP_BUTTON_TEXTS = ["💬 Yordam", "💬 Помощь", "💬 Help"]

@router.callback_query(DocumentStates.waiting_for_payment, F.data == "pay_balance_doc")
async def pay_balance_doc_handler(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Handle 'pay from balance' when user confirms payment for document service"""
    if not user:
        await callback.answer("Xatolik!", show_alert=True)
        return

    data = await state.get_data()
    price = data.get("price", 0)
    doc_next_step = data.get("doc_next_step", "")

    if user.balance < price:
        await callback.answer()
        return

    await callback.answer()
    await callback.message.edit_reply_markup(reply_markup=None)

    if doc_next_step == "presentation_template":
        await state.set_state(DocumentStates.waiting_for_template)
        await show_template_selection(callback.message, state, user_lang, group=1, edit_message=False)

    elif doc_next_step == "outline_choice":
        from bot.keyboards import get_outline_choice_keyboard
        await state.set_state(DocumentStates.waiting_for_outline_choice)
        await callback.message.answer(
            get_text(user_lang, "outline_choice"),
            reply_markup=get_outline_choice_keyboard(user_lang)
        )

    elif doc_next_step == "course_work_gen":
        # Heavy docs handle their own animated status — no pre-message here.
        await generate_course_work(callback, state, db, user_lang, user)

    elif doc_next_step == "diploma_work_gen":
        await generate_diploma_work(callback, state, db, user_lang, user)

    elif doc_next_step == "graduation_work_gen":
        await generate_graduation_work(callback, state, db, user_lang, user)

    elif doc_next_step == "dissertation_gen":
        await generate_dissertation(callback, state, db, user_lang, user)

    elif doc_next_step == "tezis_gen":
        doc_lang = data.get("doc_language", user_lang)
        await callback.message.answer("⏳ " + get_text(doc_lang, "generating"))
        await generate_thesis(callback.message, state, db, user_lang, user)

    elif doc_next_step == "maqola_gen":
        doc_lang = data.get("doc_language", user_lang)
        await callback.message.answer("⏳ " + get_text(doc_lang, "generating"))
        await generate_article(callback, state, db, user_lang, user)


@router.message(F.web_app_data)
async def handle_web_app_data(message: Message, user_lang: str, db: Database, user):
    """Handle data sent from Mini App editor"""
    _WAIT_TEXT = {"uz": "⏳ Tahrirlangan hujjat tayyorlanmoqda...", "ru": "⏳ Готовим отредактированный документ...", "en": "⏳ Preparing edited document..."}
    _DONE_TEXT = {"uz": "✅ Tahrirlangan fayl tayyor!", "ru": "✅ Отредактированный файл готов!", "en": "✅ Edited file is ready!"}
    _TOKEN_ERR = {"uz": "❌ Token muddati o'tgan. Yangi fayl oling.", "ru": "❌ Токен истёк. Получите новый файл.", "en": "❌ Token expired. Please generate a new file."}
    _GEN_ERR = {"uz": "❌ Fayl yaratishda xatolik yuz berdi.", "ru": "❌ Ошибка при создании файла.", "en": "❌ Error creating file."}

    try:
        data = json.loads(message.web_app_data.data)
        token = data.get("token", "")
        content = data.get("content", "")
        file_type = data.get("file_type", "docx")

        info = webapp.DOC_TOKENS.get(token)
        if not info:
            await message.answer(_TOKEN_ERR.get(user_lang, _TOKEN_ERR["uz"]))
            return

        topic = info.get("topic", "")
        doc_lang = info.get("doc_lang", user_lang)
        original_path = info.get("file_path", "")

        wait_msg = await message.answer(_WAIT_TEXT.get(user_lang, _WAIT_TEXT["uz"]))

        if file_type == "pptx":
            new_path = await asyncio.get_event_loop().run_in_executor(
                None, _create_edited_pptx, content, original_path
            )
        else:
            new_path = await asyncio.get_event_loop().run_in_executor(
                None, _create_edited_docx, content, topic
            )

        await wait_msg.delete()
        await message.answer_document(
            document=FSInputFile(new_path),
            caption=f"✏️ {topic}",
            reply_markup=get_main_keyboard(user_lang)
        )
        await message.answer(_DONE_TEXT.get(user_lang, _DONE_TEXT["uz"]))

        # Wipe both the user-uploaded source AND the regenerated edited file
        # immediately, plus pop+persist the token so temp/doc_tokens.json is
        # cleaned without waiting for the scheduled expiry.
        webapp.DOC_TOKENS.pop(token, None)
        webapp.save_tokens_to_disk()
        _safe_remove_file(new_path)
        _safe_remove_file(original_path)

    except json.JSONDecodeError:
        await message.answer(_GEN_ERR.get(user_lang, _GEN_ERR["uz"]))
    except Exception as e:
        logger.error(f"web_app_data handler error: {e}")
        await message.answer(_GEN_ERR.get(user_lang, _GEN_ERR["uz"]))


@router.message(F.text.in_(HELP_BUTTON_TEXTS))
async def help_handler(message: Message, state: FSMContext, user_lang: str):
    """Handles the 'Help' button click."""
    await state.clear()

    help_text = get_text(user_lang, "help_text")

    from bot.keyboards import get_help_keyboard
    await message.answer(
        help_text,
        reply_markup=get_help_keyboard(user_lang),
        parse_mode="Markdown"
    )


OTHER_SERVICES_BUTTON_TEXTS = [
    "🟩 Boshqa professional xizmatlar",
    "🟩 Другие профессиональные услуги",
    "🟩 Other Professional Services",
]


@router.message(F.text.in_(OTHER_SERVICES_BUTTON_TEXTS))
async def other_services_handler(message: Message, state: FSMContext, user_lang: str, db: Database):
    await state.clear()
    media_enabled = await db.get_feature_status("media")
    book_translate_enabled = await db.get_feature_status("book_translate")
    await message.answer(
        get_text(user_lang, "other_services_title"),
        reply_markup=get_other_services_keyboard(
            user_lang,
            media_enabled=media_enabled,
            book_translate_enabled=book_translate_enabled,
        )
    )


OS_DOC_TYPES = {
    "diploma_work": "diploma_work",
    "tezis": "tezis",
    "maqola": "maqola",
    "bitiruv_ishi": "bitiruv_ishi",
    "dissertatsiya": "dissertatsiya",
}


@router.callback_query(F.data.startswith("os:"))
async def other_services_callback_handler(callback: CallbackQuery, state: FSMContext, user_lang: str, db: Database, user):
    action = callback.data.split(":", 1)[1]
    await callback.answer()

    if action == "back":
        media_enabled = await db.get_feature_status("media")
        book_translate_enabled = await db.get_feature_status("book_translate")
        await callback.message.delete()
        await callback.message.answer(
            get_text(user_lang, "language_selected"),
            reply_markup=get_main_keyboard(user_lang, media_enabled=media_enabled, book_translate_enabled=book_translate_enabled)
        )
        return

    if action in OS_DOC_TYPES:
        await state.clear()
        doc_type = OS_DOC_TYPES[action]

        channels = await db.get_active_channels()
        if channels:
            channel_service = ChannelService(callback.message.bot)
            is_subscribed = await channel_service.check_user_subscription(user.telegram_id, channels)
            if not is_subscribed:
                from bot.keyboards import get_subscription_check_keyboard
                await callback.message.answer(
                    get_text(user_lang, "subscription_required"),
                    reply_markup=get_subscription_check_keyboard(user_lang, channels)
                )
                return

        await state.update_data(document_type=doc_type, source_step_visited=True)
        await callback.message.answer(
            get_text(user_lang, "select_source"),
            reply_markup=get_source_selection_keyboard(user_lang)
        )
        await state.set_state(DocumentStates.waiting_for_source_selection)
        return

    if action == "emoji_art":
        await state.clear()
        _sticker_texts = {
            "uz": "🎭 Stiker yaratish\n\n📸 Rasm yuboring → WebP stikerni olasiz\n🎬 Video yuboring → Kadrdan stiker yasaladi\n\nHar yuborilgan fayl shaxsiy to'plamingizga avtomatik qo'shiladi!",
            "ru": "🎭 Создать стикер\n\n📸 Отправьте фото → получите WebP стикер\n🎬 Отправьте видео → из кадра сделается стикер\n\nКаждый файл автоматически добавляется в вашу личную коллекцию!",
            "en": "🎭 Create Sticker\n\n📸 Send a photo → get a WebP sticker\n🎬 Send a video → a frame becomes your sticker\n\nEvery file is automatically added to your personal collection!",
        }
        await callback.message.answer(_sticker_texts.get(user_lang, _sticker_texts["uz"]))
        return

    if action == "pdf_convert":
        await state.clear()
        await state.set_state(None)
        from bot.handlers.converter import handle_pdf_convert_menu
        await handle_pdf_convert_menu(callback.message, state, user_lang, db, user)
        return

    if action == "pptx_to_pdf":
        await state.clear()
        from bot.handlers.pptx_converter import handle_pptx_to_pdf_menu
        await handle_pptx_to_pdf_menu(callback.message, state, user_lang)
        return

    if action == "book_translate":
        await state.clear()
        from bot.handlers.book_translate import handle_book_translate_start
        await handle_book_translate_start(callback.message, state, user_lang, db, user)
        return

    if action == "edit_file":
        await state.clear()
        await state.set_state(DocumentStates.waiting_for_edit_file)
        await callback.message.answer(get_text(user_lang, "edit_file_send_file"))
        return


@router.message(DocumentStates.waiting_for_edit_file, F.document)
async def handle_edit_file_upload(message: Message, state: FSMContext, user_lang: str, db: Database, user):
    """Receive DOCX or PPTX file and open Mini App editor."""
    doc = message.document
    fname = (doc.file_name or "").lower()

    if not (fname.endswith(".docx") or fname.endswith(".pptx")):
        await message.answer(get_text(user_lang, "edit_file_not_supported"))
        return

    try:
        os.makedirs(TEMP_DIR, exist_ok=True)
        ext = ".pptx" if fname.endswith(".pptx") else ".docx"
        local_path = os.path.join(TEMP_DIR, f"edit_upload_{doc.file_id[-12:]}{ext}")
        await message.bot.download(doc, destination=local_path)
    except Exception as e:
        logger.error(f"edit_file download error: {e}")
        await message.answer(get_text(user_lang, "edit_file_error"))
        return

    topic = doc.file_name or "Fayl"
    chat_id = message.chat.id

    kb, token = _make_edit_keyboard(
        file_path=local_path,
        topic=topic,
        doc_lang=user_lang,
        user_lang=user_lang,
        user_id=message.from_user.id,
        chat_id=chat_id,
    )

    sent = await message.answer(
        get_text(user_lang, "edit_file_ready"),
        reply_markup=kb,
    )
    _attach_and_schedule(token, sent.message_id)
    await state.clear()


async def generate_mahsus_ishlanma(callback: CallbackQuery, state: FSMContext, db: Database, user_lang: str, user):
    """Generate mahsus ishlanma document"""
    try:
        data = await state.get_data()
        topic = data['topic']
        max_pages = data.get('max_pages', 15)
        author_name = data.get('author_name', user.first_name or "")
        doc_lang = data.get('doc_language', user_lang)
        doc_extras = data.get('doc_extras') or []

        specifications = json.dumps({"min_pages": data.get('min_pages', 10), "max_pages": max_pages})
        order_id = await db.create_document_order(
            user_id=user.id,
            document_type="mahsus_ishlanma",
            topic=topic,
            specifications=specifications
        )

        # Number of practical steps based on page range
        if max_pages <= 12:
            step_count = 3
        elif max_pages <= 18:
            step_count = 4
        elif max_pages <= 24:
            step_count = 5
        else:
            step_count = 6

        ai_service = get_ai_service()
        content = await ai_service.generate_mahsus_ishlanma_content(topic, step_count, doc_lang)

        content['language'] = doc_lang
        content['author_name'] = author_name

        doc_service = get_document_service()
        file_path = await doc_service.create_mahsus_ishlanma(topic, content, extras=doc_extras or None)

        if not file_path or not os.path.exists(file_path):
            raise Exception(f"File not created: {file_path}")

        document = FSInputFile(file_path)
        await callback.message.answer_document(
            document=document,
            caption=f"🔬 {topic}",
        )

        await db.update_document_order(order_id, "completed", file_path)
        price = data.get('price', 0)
        await db.update_user_balance(user.telegram_id, -price)
        _safe_remove_file(file_path)

        await callback.message.answer(get_text(user_lang, "document_ready"), reply_markup=get_main_keyboard(user_lang))
        await callback.message.answer(get_text(user_lang, "document_reminder"), parse_mode="Markdown")

        logger.info(f"Mahsus ishlanma generated: {file_path} for user {user.telegram_id}")

    except Exception as e:
        import traceback
        logger.error(f"Error generating mahsus_ishlanma: {e}\n{traceback.format_exc()}")
        await callback.message.answer(
            "❌ Xatolik yuz berdi. Iltimos, qayta urinib ko'ring.\n\nSabab: " + _friendly_error(e),
            reply_markup=get_main_keyboard(user_lang)
        )
        _safe_remove_file(locals().get("file_path"))
        if 'order_id' in locals():
            await db.update_document_order(order_id, "failed")