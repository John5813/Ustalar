import os
import io
import logging
import asyncio
import uuid

from aiogram import Router, F
from aiogram.types import Message, FSInputFile
from aiogram.fsm.context import FSMContext

from bot.states import PptxToPdfStates
from translations import get_text

router = Router()
logger = logging.getLogger(__name__)

PPTX_CONVERT_BUTTON_TEXTS = {
    "📑 PPTX → PDF",
}

DEJAVU_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
DEJAVU_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

DPI = 150
EMU_PER_INCH = 914400


@router.message(F.text.in_(PPTX_CONVERT_BUTTON_TEXTS))
async def handle_pptx_to_pdf_menu(message: Message, state: FSMContext, user_lang: str):
    await state.clear()
    await state.set_state(PptxToPdfStates.waiting_for_pptx)
    await message.answer(get_text(user_lang, "pptx_to_pdf_send_file"))


@router.message(PptxToPdfStates.waiting_for_pptx)
async def handle_pptx_file(message: Message, state: FSMContext, user_lang: str):
    if not message.document:
        await message.answer(get_text(user_lang, "pptx_to_pdf_not_pptx"))
        return

    file_name = message.document.file_name or ""
    if not file_name.lower().endswith(".pptx"):
        await message.answer(get_text(user_lang, "pptx_to_pdf_not_pptx"))
        return

    processing_msg = await message.answer(get_text(user_lang, "pptx_to_pdf_processing"))

    pptx_path = ""
    pdf_path = ""
    try:
        os.makedirs("temp", exist_ok=True)
        uid = uuid.uuid4().hex[:8]
        pptx_path = f"temp/pptx_{uid}.pptx"
        pdf_path = f"temp/pptx_{uid}.pdf"

        file = await message.bot.get_file(message.document.file_id)
        await message.bot.download_file(file.file_path, pptx_path)

        result = await asyncio.get_event_loop().run_in_executor(
            None, _convert_pptx_to_pdf, pptx_path, pdf_path
        )

        if result and os.path.exists(pdf_path):
            base_name = os.path.splitext(file_name)[0]
            doc = FSInputFile(pdf_path, filename=f"{base_name}.pdf")
            await message.answer_document(doc, caption=get_text(user_lang, "pptx_to_pdf_done"))
        else:
            await message.answer(get_text(user_lang, "pptx_to_pdf_error"))

    except Exception as e:
        logger.error(f"PPTX to PDF conversion error: {e}")
        await message.answer(get_text(user_lang, "pptx_to_pdf_error"))
    finally:
        try:
            await processing_msg.delete()
        except Exception:
            pass
        for p in [pptx_path, pdf_path]:
            try:
                if p and os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass
        await state.clear()


def _emu_to_px(emu: int) -> int:
    return int(emu * DPI / EMU_PER_INCH)


def _get_rgb(color_obj):
    try:
        c = color_obj.rgb
        return (int(c.r), int(c.g), int(c.b))
    except Exception:
        return None


def _render_slide_to_image(slide, slide_w_emu: int, slide_h_emu: int, media: dict):
    from PIL import Image, ImageDraw, ImageFont
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.dml.color import RGBColor

    img_w = _emu_to_px(slide_w_emu)
    img_h = _emu_to_px(slide_h_emu)

    img = Image.new("RGB", (img_w, img_h), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # --- Slayd fonini render qilish ---
    try:
        bg_fill = slide.background.fill
        fill_type = bg_fill.type

        if fill_type is not None:
            try:
                from pptx.enum.dml import MSO_FILL as FT
                if fill_type == FT.SOLID:
                    rgb = _get_rgb(bg_fill.fore_color)
                    if rgb:
                        _old_img = img
                        img = Image.new("RGB", (img_w, img_h), rgb)
                        _old_img.close()
                        draw = ImageDraw.Draw(img)
                elif fill_type == FT.GRADIENT:
                    try:
                        stops = bg_fill.gradient_stops
                        if len(stops) >= 2:
                            c1 = _get_rgb(stops[0].color) or (255, 255, 255)
                            c2 = _get_rgb(stops[-1].color) or (255, 255, 255)
                            _old_img = img
                            img = _make_gradient(img_w, img_h, c1, c2)
                            _old_img.close()
                            draw = ImageDraw.Draw(img)
                    except Exception:
                        pass
                elif fill_type == FT.PICTURE:
                    try:
                        blob = bg_fill._fill.blipFill.blip.embed
                        if blob in media:
                            _raw = Image.open(io.BytesIO(media[blob])).convert("RGB")
                            bg_img = _raw.resize((img_w, img_h), Image.LANCZOS)
                            _raw.close()
                            _old_img = img
                            img = bg_img
                            _old_img.close()
                            draw = ImageDraw.Draw(img)
                    except Exception:
                        pass
            except Exception:
                pass
    except Exception:
        pass

    # --- Shakllarni render qilish ---
    for shape in slide.shapes:
        left = _emu_to_px(shape.left or 0)
        top = _emu_to_px(shape.top or 0)
        width = _emu_to_px(shape.width or 0)
        height = _emu_to_px(shape.height or 0)

        # Rasm shakllari
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                image_bytes = shape.image.blob
                _raw_pic = Image.open(io.BytesIO(image_bytes)).convert("RGBA")
                pic = _raw_pic.resize((max(1, width), max(1, height)), Image.LANCZOS)
                _raw_pic.close()
                img.paste(pic, (left, top), pic if pic.mode == "RGBA" else None)
                pic.close()
            except Exception as e:
                logger.debug(f"Image shape render error: {e}")
            continue

        # Matnli shakllar
        if not shape.has_text_frame:
            continue

        # Shakl fonini render qilish
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    from pptx.enum.dml import MSO_FILL as FT
                    if fill.type == FT.SOLID:
                        rgb = _get_rgb(fill.fore_color)
                        if rgb:
                            draw.rectangle([left, top, left + width, top + height], fill=rgb)
                    elif fill.type == FT.GRADIENT:
                        try:
                            stops = fill.gradient_stops
                            if len(stops) >= 2:
                                c1 = _get_rgb(stops[0].color) or (255, 255, 255)
                                c2 = _get_rgb(stops[-1].color) or (255, 255, 255)
                                grad = _make_gradient(width, height, c1, c2)
                                img.paste(grad, (left, top))
                        except Exception:
                            pass
                except Exception:
                    pass
        except Exception:
            pass

        # Matnni render qilish
        cur_y = top + _emu_to_px(shape.text_frame.margin_top or 0)
        margin_left = _emu_to_px(shape.text_frame.margin_left or 0)

        for para in shape.text_frame.paragraphs:
            full_text = para.text
            if not full_text.strip():
                cur_y += int(DPI * 0.15)
                continue

            # Paragraf shrift parametrlarini aniqlash
            font_size_pt = 12
            is_bold = False
            txt_color = (0, 0, 0)
            try:
                if para.runs:
                    run = para.runs[0]
                    if run.font.size:
                        font_size_pt = run.font.size.pt
                    if run.font.bold is True:
                        is_bold = True
                    rgb = _get_rgb(run.font.color)
                    if rgb:
                        txt_color = rgb
            except Exception:
                pass

            # Ota-shakl shrift o'lchamini fallback sifatida qo'llash
            try:
                if font_size_pt <= 0 and shape.text_frame.paragraphs:
                    pass
            except Exception:
                pass

            font_size_px = max(8, int((font_size_pt + 6) * DPI / 72))
            font = _get_font(is_bold, font_size_px)

            # Matn joylashuvini hisoblash
            x = left + margin_left
            max_w = max(1, width - margin_left * 2)

            # Matn satrlarga bo'lib chiqish
            lines = _wrap_text(draw, full_text, font, max_w)
            for line in lines:
                if cur_y + font_size_px > top + height + font_size_px * 2:
                    break
                # Tekislik
                try:
                    align = str(para.alignment).lower() if para.alignment else "left"
                except Exception:
                    align = "left"

                line_w = draw.textlength(line, font=font)
                if "center" in align:
                    lx = x + (max_w - line_w) / 2
                elif "right" in align:
                    lx = x + max_w - line_w
                else:
                    lx = x

                draw.text((lx, cur_y), line, font=font, fill=txt_color)
                cur_y += font_size_px + int(DPI * 0.04)

    return img


def _make_gradient(w: int, h: int, c1: tuple, c2: tuple) -> "Image":
    from PIL import Image
    base = Image.new("RGB", (w, h))
    for y in range(h):
        t = y / max(h - 1, 1)
        r = int(c1[0] + (c2[0] - c1[0]) * t)
        g = int(c1[1] + (c2[1] - c1[1]) * t)
        b = int(c1[2] + (c2[2] - c1[2]) * t)
        for x in range(w):
            base.putpixel((x, y), (r, g, b))
    return base


_font_cache: dict = {}


def _get_font(bold: bool, size_px: int):
    from PIL import ImageFont
    key = (bold, size_px)
    if key in _font_cache:
        return _font_cache[key]
    path = DEJAVU_BOLD if bold and os.path.exists(DEJAVU_BOLD) else DEJAVU_REGULAR
    try:
        if os.path.exists(path):
            font = ImageFont.truetype(path, size=size_px)
        else:
            font = ImageFont.load_default()
    except Exception:
        font = ImageFont.load_default()
    _font_cache[key] = font
    return font


def _wrap_text(draw, text: str, font, max_width: int) -> list:
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = f"{current} {word}".strip() if current else word
        try:
            w = draw.textlength(test, font=font)
        except Exception:
            w = len(test) * 7
        if w <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines if lines else [text]


def _extract_media(pptx_path: str) -> dict:
    import zipfile
    media = {}
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            for name in z.namelist():
                if name.startswith("ppt/media/"):
                    media[os.path.basename(name)] = z.read(name)
                    # Relationship rId → file mapping uchun
                    rkey = name.split("/")[-1]
                    media[rkey] = media[os.path.basename(name)]
    except Exception as e:
        logger.debug(f"Media extract error: {e}")
    return media


def _convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    try:
        import img2pdf
        from pptx import Presentation
        from PIL import Image

        prs = Presentation(pptx_path)
        slide_w_emu = prs.slide_width or int(9144000)
        slide_h_emu = prs.slide_height or int(5143500)

        media = _extract_media(pptx_path)
        img_bytes_list = []

        for slide_num, slide in enumerate(prs.slides):
            try:
                img = _render_slide_to_image(slide, slide_w_emu, slide_h_emu, media)
            except Exception as e:
                logger.warning(f"Slide {slide_num + 1} render error: {e}")
                img = Image.new("RGB", (_emu_to_px(slide_w_emu), _emu_to_px(slide_h_emu)), (255, 255, 255))

            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=92, dpi=(DPI, DPI))
            img.close()
            img_bytes_list.append(buf.getvalue())
            buf.close()
            logger.debug(f"Slide {slide_num + 1} rendered OK")

        media.clear()

        if not img_bytes_list:
            return False

        with open(pdf_path, "wb") as f:
            f.write(img2pdf.convert(img_bytes_list))

        return os.path.exists(pdf_path)

    except Exception as e:
        logger.error(f"PPTX to PDF conversion failed: {e}", exc_info=True)
        return False
